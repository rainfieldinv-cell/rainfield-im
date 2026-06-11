"""
diag_tri_xml2.py — 커버/장식 슬라이드의 모든 distinct srgbClr 수집 (평면 iter).
"""
import sys, io
from collections import Counter
from lxml import etree
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

src = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_diag.pptx"
prs = Presentation(src)
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

print("=" * 90)
print("슬라이드별 solidFill srgbClr 분포 (장식 색 후보)")
print("=" * 90)
for si, sl in enumerate(prs.slides):
    cnt = Counter()
    for sf in sl._element.iter(f"{{{A}}}solidFill"):
        for c in sf:
            if etree.QName(c).localname == "srgbClr":
                cnt[c.get("val")] += 1
    if cnt:
        top = ", ".join(f"{k}×{v}" for k, v in cnt.most_common(12))
        print(f"slide{si:02d}: {top}")

# 장식 그룹(L<7)의 각 freeform 별 색 상세
print("\n" + "=" * 90)
print("좌측 장식 그룹 내부 freeform 별 fill 색")
print("=" * 90)
for si, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        l = sh.left/360000 if sh.left is not None else 99
        if tp == "GROUP" and l < 7:
            cols = []
            for sp in sh._element.iter(f"{{{A}}}sp"):
                nv = sp.find(f".//{{{A}}}cNvPr")
                nm = nv.get("name") if nv is not None else "?"
                sf = sp.find(f".//{{{A}}}solidFill/{{{A}}}srgbClr")
                if sf is not None:
                    cols.append(f"{nm}={sf.get('val')}")
            # freeform 은 sp 가 아니라 별도일 수 있어 cxnSp/sp 모두
            print(f"slide{si} '{sh.name}': {cols}")
            break
