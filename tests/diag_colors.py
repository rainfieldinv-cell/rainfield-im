"""읽기 전용 — PPTX 색상/폰트 추출. 수정/저장 없음."""
import sys, io
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_color_src.pptx"
prs = Presentation(PATH)
print(f"열기: {PATH}  슬라이드 {len(prs.slides)}장\n")


def q(t):
    return "{%s}%s" % (A, t)


color_counter = Counter()
alpha_hits = []
fonts_ko, fonts_en = set(), set()


def scan_colors(el):
    """element 하위 모든 srgbClr 수집 + alpha."""
    for c in el.iter(q("srgbClr")):
        v = c.get("val")
        if v:
            color_counter[v] += 1
            al = c.find(q("alpha"))
            if al is not None:
                alpha_hits.append((v, al.get("val")))


def cell_fill_hex(cell):
    tcPr = cell._tc.find(q("tcPr"))
    if tcPr is not None:
        sf = tcPr.find(q("solidFill"))
        if sf is not None:
            sc = sf.find(q("srgbClr"))
            if sc is not None:
                return sc.get("val")
    return "-"


def run_info(r):
    nm = r.font.name
    if nm:
        if any('가' <= ch <= '힣' for ch in nm):
            fonts_ko.add(nm)
        else:
            fonts_en.add(nm)
    sz = r.font.size.pt if r.font.size else None
    try:
        col = str(r.font.color.rgb) if r.font.color and r.font.color.type is not None else "-"
    except Exception:
        col = "-"
    return nm, sz, r.font.bold, col


# ── 1. 표 색상 ──
print("=" * 100)
print("1. 표(table) 셀 색상/폰트")
print("=" * 100)
header_navy = Counter()
gray_cols = Counter()
tbl_count = 0
for si, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        scan_colors(sh._element)
        if not sh.has_table:
            continue
        tbl_count += 1
        tbl = sh.table
        if tbl_count <= 6:   # 처음 6개 표만 상세
            print(f"\n[슬라이드{si}] 표 {len(tbl.rows)}r×{len(tbl.columns)}c")
            for ri in range(min(len(tbl.rows), 4)):
                for ci in range(len(tbl.columns)):
                    cell = tbl.cell(ri, ci)
                    fh = cell_fill_hex(cell)
                    t = cell.text.strip()[:16]
                    rinfo = ""
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            if r.text.strip():
                                rinfo = run_info(r)
                                break
                        if rinfo:
                            break
                    print(f"   cell[{ri},{ci}] bg={fh:<8} '{t}'  {rinfo}")
        # 헤더(첫 행) 남색 / 첫 열 회색 수집
        for ci in range(len(tbl.columns)):
            fh = cell_fill_hex(tbl.cell(0, ci))
            if fh != "-":
                header_navy[fh] += 1
        for ri in range(len(tbl.rows)):
            fh = cell_fill_hex(tbl.cell(ri, 0))
            if fh != "-":
                gray_cols[fh] += 1

print(f"\n총 표 개수: {tbl_count}")
print(f"첫 행(헤더) 배경색 빈도: {dict(header_navy.most_common(6))}")
print(f"첫 열 배경색 빈도: {dict(gray_cols.most_common(6))}")

# ── 2. 표지 ──
print("\n" + "=" * 100)
print("2. 표지(slides[0]) 텍스트 폰트/색 + 도형 테두리색")
print("=" * 100)
for i, sh in enumerate(prs.slides[0].shapes):
    if sh.has_text_frame and sh.text_frame.text.strip():
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    nm, sz, b, col = run_info(r)
                    print(f"  [{i}] '{r.text.strip()[:24]}' font={nm} {sz}pt bold={b} color={col}")
    # 테두리
    ln = sh._element.find('.//' + q("ln"))
    if ln is not None:
        sc = ln.find('.//' + q("srgbClr"))
        if sc is not None:
            print(f"  [{i}] {sh.name}: line색={sc.get('val')}")

# ── 3. 전체 색 top 15 ──
print("\n" + "=" * 100)
print("3. 전체 srgbClr 빈도 top 15")
print("=" * 100)
for hexv, cnt in color_counter.most_common(15):
    print(f"   #{hexv}  ×{cnt}")

# ── 4. 폰트 목록 ──
print("\n" + "=" * 100)
print("4. 사용 폰트")
print("=" * 100)
print(f"  한글: {sorted(fonts_ko)}")
print(f"  영문/기타: {sorted(fonts_en)}")

# ── 5. alpha ──
print("\n" + "=" * 100)
print("5. 투명도(alpha) 적용 fill")
print("=" * 100)
if alpha_hits:
    for v, al in Counter(alpha_hits).most_common(15):
        print(f"   #{v[0]} alpha={int(v[1])/1000:.0f}% ×{al}")
else:
    print("   없음")
