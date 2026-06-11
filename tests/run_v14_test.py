"""run_v14_test.py — 빌더 재생성(save_presentation)으로 v14 + freeform idx16 검증."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v14.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"저장 완료: {OUTPUT}")

    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v14] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    sp = p2.slides[2].shapes
    print(f"\n[slides[2] 전체 {len(sp)} shape]")
    free_idx = None
    for i, s in enumerate(sp):
        tp = str(s.shape_type).split("(")[0].split(".")[-1].strip()
        l = Emu(s.left).inches if s.left is not None else 0
        t = Emu(s.top).inches if s.top is not None else 0
        w = Emu(s.width).inches if s.width is not None else 0
        h = Emu(s.height).inches if s.height is not None else 0
        txt = s.text_frame.text.replace("\n", " ")[:22] if s.has_text_frame else ""
        mark = ""
        if tp == "FREEFORM":
            free_idx = i; mark = " ◀FREEFORM"
        print(f"  idx={i:2d} {tp:<12} L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f} {txt}{mark}")

    if free_idx is not None:
        s = sp[free_idx]
        el = s._element
        path = el.find('.//' + qn('a:custGeom') + '/' + qn('a:pathLst') + '/' + qn('a:path'))
        from lxml import etree as ET
        L0, T0 = Emu(s.left).inches, Emu(s.top).inches
        print(f"\n[freeform idx={free_idx}] 정점 (절대 인치):")
        if path is not None:
            for seg in path:
                nm = ET.QName(seg).localname
                pt = seg.find(qn('a:pt'))
                if pt is not None:
                    x = int(pt.get('x')); y = int(pt.get('y'))
                    print(f"    {nm}: ({L0 + x/914400:.2f}, {T0 + y/914400:.2f})")
        ln = el.find('.//' + qn('a:ln'))
        clr = el.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
        tail = el.find('.//' + qn('a:tailEnd'))
        miter = el.find('.//' + qn('a:ln') + '/' + qn('a:miter'))
        print(f"  bbox L/T/W/H = {L0:.2f}/{T0:.2f}/"
              f"{Emu(s.width).inches:.2f}/{Emu(s.height).inches:.2f}")
        print(f"  선 두께={ln.get('w') if ln is not None else '?'} "
              f"색={clr.get('val') if clr is not None else '?'} "
              f"tailEnd={tail.get('type') if tail is not None else '없음'} "
              f"miter={'있음 lim=' + miter.get('lim') if miter is not None else '없음'}")
        print(f"  P3→P4 수평길이={7.90 - 7.50:.2f}\"  시작띄움={6.10 - 6.00:.2f}\"  끝띄움={8.00 - 7.90:.2f}\"")

    print("\n[다른 화살표/박스 불변]")
    for i in (10, 12, 14):
        g = sp[i]._element.find('.//' + qn('a:prstGeom'))
        print(f"  idx={i} prst={g.get('prst') if g is not None else '?'}")
    for i in range(4, 8):
        print(f"  idx={i} has_table={sp[i].has_table}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
