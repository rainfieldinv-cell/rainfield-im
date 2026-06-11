"""틀-먼저 시제품 (A-1): 8p 사업개요 → 빈 틀만 배치한 슬라이드 생성.
   ① 구성요소 명세 추출(fitz) → ② PDF 상대좌표를 가로 슬라이드 본문영역에 비율 매핑
   → ③ 빈 표 격자 + 빈 글상자 placeholder 배치(내용 없음) → 저장.
   내용 채우기(B단계)는 다음. 기존 빌더 미수정."""
import sys, io, os, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)
from modules import page_builders
from modules.page_builders import create_presentation_from_template, finalize_presentation, clone_slide_layout
from modules.ppt_generator import save_presentation
from modules.ai_slide_builders import style_table, PALETTE, add_footer

PDF = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUT = os.path.join(ROOT, "test_frame_proto_p8.pptx")

ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
if ov:
    page_builders.LAYOUT_PPTX_PATH = ov


def _inside(i, o, tol=2):
    return i[0] >= o[0]-tol and i[1] >= o[1]-tol and i[2] <= o[2]+tol and i[3] <= o[3]+tol


def extract_components(page):
    W, H = page.rect.width, page.rect.height
    comps, tbl_boxes = [], []
    for t in page.find_tables().tables:
        bbox = tuple(t.bbox); tbl_boxes.append(bbox)
        ext = t.extract() or []
        comps.append({"kind": "table", "bbox": bbox,
                      "rows": len(ext), "cols": max((len(r) for r in ext), default=0)})
    for b in page.get_text("dict").get("blocks", []):
        bbox = tuple(b["bbox"])
        if any(_inside(bbox, tb) for tb in tbl_boxes):
            continue
        if b.get("type") == 1:
            comps.append({"kind": "image", "bbox": bbox, "rows": 0, "cols": 0})
        else:
            txt = "".join(s["text"] for ln in b.get("lines", []) for s in ln.get("spans", [])).strip()
            if txt:
                comps.append({"kind": "text", "bbox": bbox, "rows": 0, "cols": 0, "text": txt})
    comps.sort(key=lambda c: (round(c["bbox"][1] / 12), c["bbox"][0]))
    return comps, W, H


# ── 8p 명세 추출 ──
doc = fitz.open(PDF)
page = doc[7]   # 8p
comps, W, H = extract_components(page)
doc.close()

# 푸터(페이지번호)·문서 상단 섹션 제목은 틀 대상에서 제외 (템플릿 헤더/푸터가 담당)
body = [c for c in comps if not (c["kind"] == "text" and ("페이지" in c.get("text", "")))]
print(f"[8p] 틀 대상 구성요소 {len(body)}개:")
for c in body:
    print(f"   {c['kind']:<6} bbox={tuple(round(v) for v in c['bbox'])} "
          f"{'%dx%d' % (c['rows'], c['cols']) if c['kind']=='table' else c.get('text','')[:24]}")

# ── 본문 영역 정의(가로 슬라이드, 인치) + 콘텐츠 밴드 정규화 매핑 ──
BODY_L, BODY_T, BODY_W, BODY_H = 0.6, 1.75, 9.6, 5.2   # 헤더 아래 ~ 푸터 위
xs = [c["bbox"][0] for c in body] + [c["bbox"][2] for c in body]
ys = [c["bbox"][1] for c in body] + [c["bbox"][3] for c in body]
xmin, xmax, ymin, ymax = min(xs), max(xs), min(ys), max(ys)


def mapx(px):
    return BODY_L + (px - xmin) / (xmax - xmin) * BODY_W


def mapy(py):
    return BODY_T + (py - ymin) / (ymax - ymin) * BODY_H


# ── 슬라이드 생성 (브랜드 템플릿 content, 표/이미지 제거한 깨끗한 베이스) ──
prs = create_presentation_from_template()
n = len(prs.slides)
slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)

placed = []
for c in body:
    x0, y0, x1, y1 = c["bbox"]
    L, T = mapx(x0), mapy(y0)
    Wd, Ht = mapx(x1) - L, mapy(y1) - T
    if c["kind"] == "table" and c["rows"] >= 1 and c["cols"] >= 1:
        gf = slide.shapes.add_table(c["rows"], c["cols"],
                                    Inches(L), Inches(T), Inches(Wd), Inches(max(Ht, 0.3)))
        # 빈 틀 — 내용 없이 격자/디자인만
        style_table(gf, has_header=True, label_cols=(), header_fill=PALETTE["navy_dark"])
        placed.append(("table(빈격자)", c["rows"], c["cols"], L, T, Wd, Ht))
    elif c["kind"] == "image":
        sp = slide.shapes.add_shape(1, Inches(L), Inches(T), Inches(max(Wd, 0.5)), Inches(max(Ht, 0.5)))
        sp.fill.solid(); sp.fill.fore_color.rgb = PALETTE["gray"]
        sp.line.color.rgb = PALETTE["gray_text"]
        sp.text_frame.text = "[사진 자리]"
        placed.append(("image(placeholder)", 0, 0, L, T, Wd, Ht))
    else:  # text — 위치/크기만 잡힌 빈 글상자(외곽선만)
        tb = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(max(Wd, 0.5)), Inches(max(Ht, 0.25)))
        tb.line.color.rgb = PALETTE["steel"]; tb.line.width = Pt(0.5)
        placed.append(("text(빈글상자)", 0, 0, L, T, Wd, Ht))

add_footer(slide, 1, business_name="천안 부성2지구 도시개발사업")
finalize_presentation(prs, n)
save_presentation(prs, OUT)

print(f"\n저장: {OUT}")
dup = [nm for nm, c in Counter(zipfile.ZipFile(OUT).namelist()).items() if c > 1]
print(f"중복파트: {'있음' if dup else '없음(정상)'}")
print("\n배치된 빈 틀 (종류 / 표행열 / L,T,W,H 인치):")
for kind, r, cc, L, T, Wd, Ht in placed:
    dim = f"{r}x{cc}" if r else "-"
    print(f"   {kind:<18} {dim:<6} L={L:.2f} T={T:.2f} W={Wd:.2f} H={Ht:.2f}")
