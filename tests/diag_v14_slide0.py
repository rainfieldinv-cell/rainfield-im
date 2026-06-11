import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v14_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[0]
for i, sp in enumerate(slide.shapes):
    txt = ""
    if sp.has_text_frame:
        txt = sp.text_frame.text.replace("\n", " / ")[:60]
    elif sp.has_table:
        txt = "[TABLE]"
    print(f"idx={i} type={sp.shape_type} name={sp.name} "
          f"L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
          f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f} <{txt}>")
print("총 shape:", len(slide.shapes))
