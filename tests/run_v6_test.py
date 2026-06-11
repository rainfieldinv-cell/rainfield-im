"""
run_v6_test.py — test_ai_builders_v6.pptx 생성 + 슬라이드 7 검증 (방식 A 표 구조).
빌더 코드 수정 없음. 캐시 HIT. LAYOUT_OVERRIDE 로 템플릿 잠금 회피.
"""
import sys, os, io, traceback
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v6.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def hexfill(cell):
    try:
        f = cell.fill
        if f.type is not None and f.fore_color is not None and f.fore_color.type is not None:
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return "-"


def main():
    print("=" * 80)
    print("STEP 5-3: test_ai_builders_v6.pptx 생성")
    print("=" * 80)

    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov
        print(f"[override] LAYOUT_PPTX_PATH → {ov}")

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    print("\n--- API ---")
    for lb, r in [("S2", r2), ("S5", r5), ("S7", r7)]:
        print(f"  {lb}: {'OK' if r.get('ok') else 'FAIL'} {'캐시HIT' if r.get('cached') else 'API'} {r.get('usage', {})}")
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    s7 = build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}\n총 슬라이드: {len(prs.slides)}")

    # ── 전체 shape 표 ──
    print(f"\n{'='*96}\n슬라이드 7 shape (총 {len(s7.shapes)}개)\n{'='*96}")
    LBL = {"신탁계약", "담보신탁 우선수익권", "공사도급계약", "대출약정"}
    tables, conns, labels = [], [], []
    for i, sh in enumerate(s7.shapes):
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        li = (sh.left or 0)/914400; ti = (sh.top or 0)/914400
        wi = (sh.width or 0)/914400; hi = (sh.height or 0)/914400
        if sh.has_table:
            tables.append(sh)
            tb = sh.table
            print(f"[{i:02d}] TABLE {len(tb.rows)}r L={li:.2f}\" T={ti:.2f}\" W={wi:.2f}\" H={hi:.2f}\" "
                  f"헤더={hexfill(tb.cell(0,0))} 헤더텍스트='{tb.cell(0,0).text.strip()}'")
            for r in range(len(tb.rows)):
                print(f"        row{r}: 배경={hexfill(tb.cell(r,0)):<8} '{tb.cell(r,0).text.strip()[:18]}'")
        else:
            if "Connector" in (sh.name or ""):
                conns.append(sh)
            txt = sh.text_frame.text.replace("\n", " ")[:24] if sh.has_text_frame else ""
            if txt.strip() in LBL:
                labels.append(txt.strip())
            nm = "Conn" if "Connector" in (sh.name or "") else ""
            print(f"[{i:02d}] {tp:<11}{nm:<5} L={li:5.2f}\" T={ti:5.2f}\" W={wi:5.2f}\" H={hi:5.2f}\" {txt}")

    # ── 검증 ──
    hdr_map = {tb.table.cell(0,0).text.strip(): hexfill(tb.table.cell(0,0)) for tb in tables}
    daeju = next((tb for tb in tables if tb.table.cell(0,0).text.strip() == "대주"), None)
    box_label_text = []
    for tb in tables:
        for r in range(len(tb.table.rows)):
            box_label_text.append(tb.table.cell(r,0).text.strip())
    box_label_text += labels
    amount_hit = [t for t in box_label_text if ("억" in t or "원" in t)]
    all_txt = " ".join(sh.text_frame.text for sh in s7.shapes if sh.has_text_frame)
    asset_mgr = "자산관리자" in all_txt
    investor = any(k in all_txt for k in ["삼성증권", "KT Estate", "우미건설", "금성백조", "Equity"])
    pics = [sh for sh in s7.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]

    print(f"\n{'='*96}\nPDF p4 비교\n{'='*96}")
    exp = {"신탁사": "013A73", "시공사": "0070C0", "차주": "2E75B6", "대주": "8C4A59"}
    color_ok = all(hdr_map.get(k) == v for k, v in exp.items())
    checks = [
        ("박스 4개", len(tables) == 4, f"{len(tables)}개"),
        ("화살표 4개", len(conns) == 4, f"{len(conns)}개"),
        ("라벨 4개", len(set(labels)) == 4, f"{sorted(set(labels))}"),
        ("헤더색 구분(013A73/0070C0/2E75B6/8C4A59)", color_ok, f"{hdr_map}"),
        ("대주 3r", daeju is not None and len(daeju.table.rows) == 3, f"{len(daeju.table.rows) if daeju else '?'}r"),
        ("금액 없음(박스/라벨)", len(amount_hit) == 0, f"{amount_hit or '없음'}"),
        ("자산관리자 박스 없음", not asset_mgr, "없음" if not asset_mgr else "있음!"),
        ("투자자 박스 없음", not investor, "없음" if not investor else "있음!"),
        ("사업명 이미지 없음", len(pics) == 0, f"PICTURE {len(pics)}"),
    ]
    for nm, ok, val in checks:
        print(f"  [{'OK ' if ok else 'FAIL'}] {nm:<38} v6={val}")
    print(f"\n  결과: {sum(1 for _,ok,_ in checks if ok)}/{len(checks)} PASS")
    print(f"\n완료 — {OUTPUT}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
