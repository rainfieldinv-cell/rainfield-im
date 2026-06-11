"""읽기 전용 — 대전 PDF(세로 원본) vs 사람이 만든 대전 PPT(가로) 구성 비교.
   디자인 말고 '구성'(요소 배치/텍스트 묶음)을 정독하기 위함."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
import fitz
from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PDF = os.path.join(ROOT, "[신영증권] 대전중구 서남부터미널 토지담보대출_IM_v3.0.pdf")
PPT = os.path.join(ROOT, "[Rainfield] 대전서남터미널_토지담보대출_260511.pptx")

MODE = sys.argv[1] if len(sys.argv) > 1 else "overview"

if MODE == "overview":
    doc = fitz.open(PDF)
    print(f"[PDF] {len(doc)}페이지 (세로 {doc[0].rect.width:.0f}x{doc[0].rect.height:.0f})")
    for pi in range(len(doc)):
        p = doc[pi]
        nt = len(p.find_tables().tables)
        ni = len([1 for im in p.get_images(full=True) if p.get_image_rects(im[0])])
        first = ""
        for b in p.get_text("dict").get("blocks", []):
            if b.get("type") == 0:
                t = "".join(s["text"] for ln in b.get("lines", []) for s in ln.get("spans", [])).strip()
                if len(t) >= 3:
                    first = t[:48]; break
        print(f"  p{pi+1:>2} | 표{nt} 이미지{ni} | {first}")
    doc.close()
    prs = Presentation(PPT)
    print(f"\n[사람 PPT] {len(prs.slides)}슬라이드 (가로 {Emu(prs.slide_width).inches:.2f}x{Emu(prs.slide_height).inches:.2f}in)")
    for si, sl in enumerate(prs.slides):
        nt = sum(1 for s in sl.shapes if s.has_table)
        npic = sum(1 for s in sl.shapes if s.shape_type == 13)
        ntx = sum(1 for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip())
        titles = [s.text_frame.text.strip().replace("\n", " ")[:30] for s in sl.shapes
                  if s.has_text_frame and s.text_frame.text.strip() and Emu(s.top).inches < 1.4]
        print(f"  s{si:>2} | 표{nt} 그림{npic} 글{ntx} | {' / '.join(titles[:2])}")

elif MODE == "pdf":
    doc = fitz.open(PDF)
    for pi in [int(x) - 1 for x in sys.argv[2:]]:
        p = doc[pi]
        print(f"\n{'='*88}\n[PDF p{pi+1}] {p.rect.width:.0f}x{p.rect.height:.0f}")
        for b in sorted(p.get_text("dict").get("blocks", []), key=lambda b: b["bbox"][1]):
            if b.get("type") != 0:
                continue
            t = " ".join(s["text"] for ln in b.get("lines", []) for s in ln.get("spans", [])).strip()
            if t:
                y = b["bbox"][1]
                print(f"  T{y:6.0f} | {t[:74]}")
    doc.close()

elif MODE == "ppt":
    prs = Presentation(PPT)
    for si in [int(x) for x in sys.argv[2:]]:
        sl = prs.slides[si]
        print(f"\n{'='*88}\n[PPT s{si}] 요소 {len(sl.shapes)}개")
        for sh in sorted(sl.shapes, key=lambda s: (Emu(s.top or 0).inches, Emu(s.left or 0).inches)):
            L = Emu(sh.left or 0).inches; T = Emu(sh.top or 0).inches
            W = Emu(sh.width or 0).inches; H = Emu(sh.height or 0).inches
            if sh.has_table:
                t = sh.table
                hdr = " | ".join(t.cell(0, c).text.strip()[:10] for c in range(min(len(t.columns), 6)))
                info = f"[표 {len(t.rows)}r×{len(t.columns)}c] {hdr}"
            elif sh.shape_type == 13:
                info = "[그림]"
            elif sh.has_text_frame and sh.text_frame.text.strip():
                info = "[글] " + sh.text_frame.text.replace("\n", " ⏎ ").strip()[:60]
            else:
                continue
            print(f"  L{L:5.2f} T{T:5.2f} W{W:5.2f} H{H:5.2f} | {info}")
