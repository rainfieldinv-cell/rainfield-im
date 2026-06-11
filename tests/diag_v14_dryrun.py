"""v14 dry-run — 새 정점 + miter join freeform 구성 후 XML/정점/길이만 출력. 저장 안 함."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v13_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]

PTS_IN = [(6.10, 2.55), (7.50, 2.55), (7.50, 3.55), (7.90, 3.55)]
pts = [(int(x * 914400), int(y * 914400)) for (x, y) in PTS_IN]

fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
fb.add_line_segments(pts[1:], close=False)
shp = fb.convert_to_shape()
shp.fill.background()
shp.line.color.rgb = RGBColor(0x40, 0x40, 0x40)
shp.line.width = Pt(1.5)
shp.shadow.inherit = False
ln = shp.line._get_or_add_ln()
# miter join (round 제거 → 칼직각). 순서상 tailEnd 보다 앞에 위치
ln.append(ln.makeelement(qn('a:miter'), {'lim': '800000'}))
ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))

print("[새 freeform 정점]")
path = shp._element.find('.//' + qn('a:custGeom') + '/' + qn('a:pathLst') + '/' + qn('a:path'))
L0, T0 = Emu(shp.left).inches, Emu(shp.top).inches
for seg in path:
    nm = etree.QName(seg).localname
    pt = seg.find(qn('a:pt'))
    if pt is not None:
        x = int(pt.get('x')); y = int(pt.get('y'))
        print(f"  {nm}: ({L0 + x/914400:.2f}, {T0 + y/914400:.2f}) in")

print(f"\nbbox L/T/W/H = {L0:.2f}/{T0:.2f}/{Emu(shp.width).inches:.2f}/{Emu(shp.height).inches:.2f}")
print(f"P3→P4 수평구간 길이 = {7.90 - 7.50:.2f} in")
print(f"시작 띄움(신탁사 우변6.00→P1 6.10) = {6.10 - 6.00:.2f} in")
print(f"끝 띄움(P4 7.90→대주 좌변8.00) = {8.00 - 7.90:.2f} in")

print("\n--- ln XML (miter 적용) ---")
print(etree.tostring(ln, pretty_print=True).decode())
print("[저장 안 함 — dry-run]")
