"""v10 dry-run — idx16 bentConnector3 bbox + adj1=85000 조정. 저장 안 함."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v9_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]
sp = slide.shapes[16]

print("변경 전 L/T/W/H:", f"{Emu(sp.left).inches:.2f}/{Emu(sp.top).inches:.2f}/"
      f"{Emu(sp.width).inches:.2f}/{Emu(sp.height).inches:.2f}")

sp.left, sp.top, sp.width, sp.height = Inches(6.00), Inches(2.50), Inches(2.00), Inches(1.40)

prstGeom = sp._element.find('.//' + qn('a:prstGeom'))
for av in prstGeom.findall(qn('a:avLst')):
    prstGeom.remove(av)
avLst = etree.SubElement(prstGeom, qn('a:avLst'))
gd = etree.SubElement(avLst, qn('a:gd'))
gd.set('name', 'adj1'); gd.set('fmla', 'val 85000')

print("변경 후 L/T/W/H:", f"{Emu(sp.left).inches:.2f}/{Emu(sp.top).inches:.2f}/"
      f"{Emu(sp.width).inches:.2f}/{Emu(sp.height).inches:.2f}")
print("prst:", prstGeom.get('prst'))
print("\n--- prstGeom XML ---")
print(etree.tostring(prstGeom, pretty_print=True).decode())

# 선 스타일 불변 확인
ln = sp._element.find('.//' + qn('a:ln'))
clr = sp._element.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
tail = sp._element.find('.//' + qn('a:tailEnd'))
print("ln(w,색,tail):", ln.get('w'), clr.get('val') if clr is not None else None,
      tail.get('type') if tail is not None else None)
print("\n[저장 안 함 — dry-run]")
