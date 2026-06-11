"""읽기 전용 — 사람이 만든 3개 제안서(대전/천안/넷마블)의 슬라이드 구성 정독.
   항상 들어가는 표준 페이지(표지/하이라이트/목차/섹션/사모사채개요/투자구조도/연락처) 파악용."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTS = {
    "대전": "[Rainfield] 대전서남터미널_토지담보대출_260511.pptx",
    "천안": "[Rainfield]천안 부성2지구_토지담보출_사모사채_제안서_260310.pptx",
    "넷마블": "[Rainfield]로 넷마블 G타워 1종 수익증권 제안서_260526.pptx",
}
SHOW = sys.argv[1] if len(sys.argv) > 1 else "all"
LIMIT = int(sys.argv[2]) if len(sys.argv) > 2 else 999

for name, fn in PPTS.items():
    if SHOW != "all" and SHOW != name:
        continue
    prs = Presentation(os.path.join(ROOT, fn))
    print(f"\n{'#'*92}\n[{name}] {len(prs.slides)}슬라이드\n{'#'*92}")
    for si, sl in enumerate(prs.slides):
        if si >= LIMIT:
            break
        nt = sum(1 for s in sl.shapes if s.has_table)
        npic = sum(1 for s in sl.shapes if s.shape_type == 13)
        # 상단 텍스트(섹션라벨/소제목) + 첫 표 헤더
        tops = []
        for s in sorted(sl.shapes, key=lambda x: Emu(x.top or 0).inches):
            if s.has_text_frame and s.text_frame.text.strip():
                t = s.text_frame.text.replace("\n", " ⏎ ").strip()
                if t:
                    tops.append(t[:46])
            if len(tops) >= 3:
                break
        first_tbl = ""
        for s in sl.shapes:
            if s.has_table:
                tb = s.table
                first_tbl = " / 표:" + " | ".join(
                    tb.cell(0, c).text.strip()[:8] for c in range(min(len(tb.columns), 6)))
                break
        print(f" s{si:>2} | 표{nt} 그림{npic} | {' ‖ '.join(tops[:3])}{first_tbl}")
