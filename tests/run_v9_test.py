"""run_v9_test.py — 수정된 빌더로 v9 빌드 + idx16 bentConnector3 검증."""
import sys, os, io, traceback
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v9.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"저장 완료: {OUTPUT}")

    # 재오픈 검증
    p2 = Presentation(OUTPUT)
    sp = p2.slides[2].shapes

    def prst(s):
        g = s._element.find('.//' + qn('a:prstGeom'))
        return g.get('prst') if g is not None else "?"

    def offext(s):
        x = s._element.find('.//' + qn('a:off')); e = s._element.find('.//' + qn('a:ext'))
        return (x.get('x'), x.get('y'), e.get('cx'), e.get('cy'))

    def lninfo(s):
        ln = s._element.find('.//' + qn('a:ln'))
        clr = s._element.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
        tail = s._element.find('.//' + qn('a:tailEnd'))
        head = s._element.find('.//' + qn('a:headEnd'))
        return (ln.get('w') if ln is not None else None,
                clr.get('val') if clr is not None else None,
                tail.get('type') if tail is not None else None,
                head.get('type') if head is not None else None)

    s16 = sp[16]
    print(f"\n{'='*70}\nv9 slides[2] idx16 검증\n{'='*70}")
    print(f"  prst        = {prst(s16)}  (기대: bentConnector3)")
    print(f"  off/ext     = {offext(s16)}  (기대: 5486400/2468880/1828800/640080)")
    print(f"  L/T/W/H(in) = {Emu(s16.left).inches:.2f}/{Emu(s16.top).inches:.2f}/"
          f"{Emu(s16.width).inches:.2f}/{Emu(s16.height).inches:.2f}")
    print(f"  ln(w,색,tail,head) = {lninfo(s16)}  (기대: 19050/404040/triangle/None)")

    print(f"\n  [다른 커넥터 prst — line 유지 확인]")
    for i in (10, 12, 14):
        print(f"    idx={i} prst={prst(sp[i])}")
    print(f"  [박스 표 idx4~7 불변]")
    for i in range(4, 8):
        print(f"    idx={i} has_table={sp[i].has_table} L={Emu(sp[i].left).inches:.2f} T={Emu(sp[i].top).inches:.2f}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
