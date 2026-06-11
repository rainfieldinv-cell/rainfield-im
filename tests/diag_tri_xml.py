"""
diag_tri_xml.py — 좌측 장식 GROUP 내부 FREEFORM 들의 fill 색(srgbClr/schemeClr)을 XML에서 추출.
"""
import sys, io
from lxml import etree
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

src = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_diag.pptx"
prs = Presentation(src)
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def walk(el, depth, slabel):
    """spPr 안의 solidFill 색을 추출."""
    tag = etree.QName(el).localname
    if tag in ("sp", "pic", "cxnSp"):
        nv = el.find(f".//{{{A}}}cNvPr")
        name = nv.get("name") if nv is not None else "?"
        # solidFill 찾기
        colors = []
        for sf in el.iter(f"{{{A}}}solidFill"):
            for c in sf:
                ln = etree.QName(c).localname
                val = c.get("val")
                # 자식 lumMod/lumOff/shade/tint 등 변조 표시
                mods = [f"{etree.QName(m).localname}={m.get('val')}" for m in c]
                colors.append(f"{ln}:{val}" + (f"[{','.join(mods)}]" if mods else ""))
            break  # 첫 solidFill만
        if colors:
            print(f"  {'  '*depth}[{slabel}] '{name[:24]}' fill={colors}")
    if tag in ("grpSp", "sp") or tag.endswith("Sp"):
        for ch in el:
            walk(ch, depth + 1, slabel)


print("=" * 90)
print("좌측 장식 그룹 내부 색 추출 (cover/슬라이드0~3)")
print("=" * 90)
for si in range(min(4, len(prs.slides))):
    sl = prs.slides[si]
    for sh in sl.shapes:
        nm = sh.name or ""
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        l = sh.left/360000 if sh.left is not None else 99
        t = sh.top/360000 if sh.top is not None else 99
        if tp == "GROUP" and l < 7:
            print(f"\nslide{si} GROUP '{nm}' (L={l:.2f} T={t:.2f}) ─ 내부 색:")
            # 그룹 element 순회
            for ch in sh._element:
                walk(ch, 0, f"s{si}")

# theme1.xml 색 매핑도 출력 (schemeClr 해석용)
print("\n" + "=" * 90)
print("테마 색상 팔레트 (schemeClr → srgbClr 매핑)")
print("=" * 90)
try:
    part = prs.part.package.part_related_by(  # noqa
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
except Exception:
    part = None
# 더 단순하게: 슬라이드 마스터의 theme part
for p in prs.part.package.iter_parts():
    if "theme" in p.partname:
        root = etree.fromstring(p.blob)
        clrs = root.find(f".//{{{A}}}clrScheme")
        if clrs is not None:
            for c in clrs:
                nm = etree.QName(c).localname
                child = list(c)[0]
                val = child.get("val") or child.get("lastClr")
                kind = etree.QName(child).localname
                print(f"  {nm:<8} = {kind}:{val}")
        break
