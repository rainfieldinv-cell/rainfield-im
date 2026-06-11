"""run_v19_test.py — 섹션2 본문 글머리 정리(clean_bullet) 반영. page_num=1/2/3."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v19.pptx")
BIZ = "천안 부성2지구 도시개발사업"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def q(t):
    return "{%s}%s" % (A, t)


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ, page_num=1)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ, page_num=2)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ, page_num=3)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}")
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v19] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    sl = p2.slides[0]
    s2 = next((sh for sh in sl.shapes if sh.name == "TextBox 26"), None)
    print(f"\n[섹션2 본문 TextBox 26 단락 상태]")
    for pi, para in enumerate(s2.text_frame.paragraphs):
        txt = "".join(r.text for r in para.runs)
        pPr = para._pPr
        bu = "기본"
        marL = indent = lvl = None
        if pPr is not None:
            if pPr.find(q('buNone')) is not None:
                bu = "buNone"
            elif pPr.find(q('buChar')) is not None:
                bu = "buChar='%s'" % pPr.find(q('buChar')).get('char')
            elif pPr.find(q('buAutoNum')) is not None:
                bu = "buAutoNum"
            marL = pPr.get('marL'); indent = pPr.get('indent'); lvl = pPr.get('lvl')
        print(f"  p{pi} level={para.level} lvl속성={lvl} bullet={bu} "
              f"marL={marL} indent={indent} text='{txt[:30]}'")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
