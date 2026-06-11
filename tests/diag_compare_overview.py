"""읽기 전용 — 원본 PDF(세로) vs 사람이 만든 제안서(가로) 개요 비교."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
import fitz
from pptx import Presentation
from pptx.util import Emu

ROOT = r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im"
PDF = os.path.join(ROOT, "신영증권_구로_넷마블_지타워_담보대출_및_1종_수익증권_IM_v3_2.pdf")
PPT = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_netmarble_prop.pptx"

# ── 원본 PDF: 페이지별 구성요소 요약 ──
doc = fitz.open(PDF)
print(f"[원본 PDF] {len(doc)}페이지 (세로 {doc[0].rect.width:.0f}x{doc[0].rect.height:.0f})")
print(f"{'pg':>3} | 표 | 이미지 | 첫 텍스트")
for pi in range(len(doc)):
    p = doc[pi]
    nt = len(p.find_tables().tables)
    ni = len([1 for im in p.get_images(full=True) if p.get_image_rects(im[0])])
    txt = ""
    for b in p.get_text("dict").get("blocks", []):
        if b.get("type") == 0:
            t = "".join(s["text"] for ln in b.get("lines", []) for s in ln.get("spans", [])).strip()
            if len(t) >= 3:
                txt = t[:36]; break
    print(f"{pi+1:>3} | {nt:>2} | {ni:>4}  | {txt}")
doc.close()

# ── 제안서 PPT: 슬라이드별 구성요소 요약 ──
prs = Presentation(PPT)
print(f"\n[제안서 PPT] {len(prs.slides)}슬라이드 (가로 {Emu(prs.slide_width).inches:.2f}x{Emu(prs.slide_height).inches:.2f}in)")
print(f"{'sl':>3} | 표 | 그림 | 글상자 | 제목/소제목")
for si, sl in enumerate(prs.slides):
    nt = sum(1 for s in sl.shapes if s.has_table)
    npic = sum(1 for s in sl.shapes if s.shape_type == 13)
    ntx = sum(1 for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip())
    # 소제목 추정: 짧은 텍스트 상단
    titles = [s.text_frame.text.strip().replace("\n", " ")[:28] for s in sl.shapes
              if s.has_text_frame and s.text_frame.text.strip() and Emu(s.top).inches < 1.5]
    print(f"{si:>3} | {nt:>2} | {npic:>3}  | {ntx:>4}  | {' / '.join(titles[:2])}")
