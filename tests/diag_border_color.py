"""읽기 전용 — 레이아웃.pptx 테마 accent3(회색 강조3) HEX + 표 셀 테두리 실제 색."""
import sys, io
from lxml import etree
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
SRC = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_diag.pptx"
prs = Presentation(SRC)


def q(t):
    return "{%s}%s" % (A, t)


# 1) 테마 색상표
print("=== 테마 clrScheme ===")
for p in prs.part.package.iter_parts():
    if "theme" in p.partname:
        root = etree.fromstring(p.blob)
        cs = root.find(".//" + q("clrScheme"))
        if cs is not None:
            for c in cs:
                nm = etree.QName(c).localname
                child = list(c)[0]
                val = child.get("val") or child.get("lastClr")
                print(f"  {nm:<8} = {val}")
        break

# 2) 표 셀 테두리 색 (레이아웃 모든 슬라이드의 첫 표)
print("\n=== 표 셀 테두리(tcBorders) 색·두께 샘플 ===")
found = 0
for si, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        if not sh.has_table:
            continue
        tbl = sh.table
        cell = tbl.cell(0, 0)
        tcPr = cell._tc.find(q("tcPr"))
        print(f"\n슬{si} 표 cell[0,0] '{cell.text.strip()[:12]}'")
        if tcPr is None:
            print("   tcPr 없음")
        else:
            for edge in ("lnL", "lnR", "lnT", "lnB"):
                ln = tcPr.find(q(edge))
                if ln is None:
                    print(f"   {edge}: 없음(테마 기본)")
                    continue
                w = ln.get("w")
                clr = ln.find(".//" + q("srgbClr"))
                sch = ln.find(".//" + q("schemeClr"))
                cinfo = (f"srgb {clr.get('val')}" if clr is not None
                         else (f"scheme {sch.get('val')}" if sch is not None else "색없음/none"))
                print(f"   {edge}: w={w} ({int(w)/12700:.2f}pt) " if w else f"   {edge}: w=기본 ",
                      cinfo)
        found += 1
        break
    if found >= 3:
        break
