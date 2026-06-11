"""
diag_triangles.py — 레이아웃 전 슬라이드에서 삼각형 장식 도형의 fill RGB 추출.
코드 수정 없음. 로컬 복사본을 직접 연다.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

src = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_diag.pptx"
prs = Presentation(src)


def fill_hex(sh):
    try:
        f = sh.fill
        if f.type is not None and f.fore_color is not None and f.fore_color.type is not None:
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return "-"


def autoshape_name(sh):
    try:
        return str(sh.adjustments) if False else ""
    except Exception:
        return ""


print("=" * 100)
print("삼각형/장식 도형 fill 색 스캔 (AUTO_SHAPE / FREEFORM)")
print("=" * 100)
for si, sl in enumerate(prs.slides):
    for i, sh in enumerate(sl.shapes):
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        nm = (sh.name or "")
        is_tri = ("삼각형" in nm) or ("Triangle" in nm) or ("triangle" in nm)
        if tp in ("AUTO_SHAPE", "FREEFORM") or is_tri:
            l = sh.left/360000 if sh.left is not None else 0
            t = sh.top/360000 if sh.top is not None else 0
            w = sh.width/360000 if sh.width is not None else 0
            h = sh.height/360000 if sh.height is not None else 0
            mark = " ◀삼각형" if is_tri else ""
            print(f"slide{si:02d}[{i:02d}] {tp:<11} name='{nm[:24]}' fill={fill_hex(sh):<8} "
                  f"L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f}{mark}")

# 좌측 상단(특히 cover/슬라이드0~2) 삼각형 집중 출력
print("\n" + "=" * 100)
print("좌측 상단 영역(L<6, T<8) 도형 색 — 삼각형 장식 후보")
print("=" * 100)
seen = []
for si in range(min(3, len(prs.slides))):
    sl = prs.slides[si]
    for i, sh in enumerate(sl.shapes):
        l = sh.left/360000 if sh.left is not None else 99
        t = sh.top/360000 if sh.top is not None else 99
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        if l < 6 and t < 8 and tp in ("AUTO_SHAPE", "FREEFORM", "GROUP"):
            fh = fill_hex(sh)
            print(f"slide{si}[{i:02d}] {tp:<10} name='{(sh.name or '')[:22]}' fill={fh} L={l:.2f} T={t:.2f}")
            # 그룹이면 내부 자식도
            if tp == "GROUP":
                for j, ch in enumerate(sh.shapes):
                    chtp = str(ch.shape_type).split("(")[0].split(".")[-1].strip()
                    print(f"      └ child[{j}] {chtp:<10} name='{(ch.name or '')[:20]}' fill={fill_hex(ch)}")
