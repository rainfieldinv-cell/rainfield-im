"""run_v20_test.py — 슬7 박스 글씨 13/12pt + 다이어그램 세로중앙. page_num=1/2/3."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v20.pptx")
BIZ = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.oxml.ns import qn

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ, page_num=1)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ, page_num=2)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ, page_num=3)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}")
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v20] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    sl = p2.slides[2]
    LBLS = {"신탁계약", "공사도급계약", "대출약정", "담보신탁 우선수익권"}
    diag_tops = []
    print(f"\n[slides[2] 전체 {len(sl.shapes)} shape]")
    for i, sp in enumerate(sl.shapes):
        tp = str(sp.shape_type).split("(")[0].split(".")[-1].strip()
        l = Emu(sp.left).inches if sp.left is not None else 0
        t = Emu(sp.top).inches if sp.top is not None else 0
        w = Emu(sp.width).inches if sp.width is not None else 0
        h = Emu(sp.height).inches if sp.height is not None else 0
        fonts = []
        marg = ""
        if sp.has_table:
            tb = sp.table
            for r in range(len(tb.rows)):
                c0 = tb.cell(r, 0)
                for p in c0.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            fonts.append(run.font.size.pt)
                marg = f"margL={Emu(c0.margin_left).inches:.2f}"
            txt = "/".join(tb.cell(r, 0).text.strip() for r in range(len(tb.rows)))
        elif sp.has_text_frame:
            txt = sp.text_frame.text.replace("\n", " ")[:24]
            for p in sp.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.size:
                        fonts.append(run.font.size.pt)
        else:
            txt = ""
        is_diag = (sp.has_table or "Connector" in (sp.name or "") or tp == "FREEFORM"
                   or (sp.has_text_frame and sp.text_frame.text.strip() in LBLS))
        if is_diag:
            diag_tops.append((t, t + h))
        print(f"  idx={i:2d} {tp:<11} {sp.name:<16} L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f} "
              f"font={sorted(set(fonts))} {marg} <{txt}>")

    dmin = min(t for t, _ in diag_tops)
    dmax = max(b for _, b in diag_tops)
    print(f"\n다이어그램 T 범위: {dmin:.2f} ~ {dmax:.2f} (중심 {(dmin+dmax)/2:.2f})")
    print("헤더/푸터 구분선: 상단 1.53, 하단 7.15 → 그 사이 중심 4.34")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
