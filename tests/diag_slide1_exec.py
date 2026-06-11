"""
diag_slide1_exec.py — executive_summary(템플릿 index 1) 구조 정밀 덤프.
<> 컨테이너/섹션/그룹/텍스트박스 위치·크기 파악. 코드 수정 없음.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

src = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_diag.pptx"
prs = Presentation(src)
slide = prs.slides[1]


def cm(v):
    return v/360000 if v is not None else 0


def dump(sh, depth=0, idx=""):
    tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
    nm = (sh.name or "")[:26]
    l, t, w, h = cm(sh.left), cm(sh.top), cm(sh.width), cm(sh.height)
    txt = ""
    if sh.has_text_frame:
        txt = sh.text_frame.text.replace("\n", " | ")[:40]
    pad = "  " * depth
    print(f"{pad}[{idx:>4}] {tp:<12} '{nm:<26}' L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f} | {txt}")
    if tp == "GROUP":
        for j, ch in enumerate(sh.shapes):
            dump(ch, depth+1, f"{idx}.{j}")


print("=" * 120)
print(f"executive_summary (slide index 1) — 총 {len(slide.shapes)} shape")
print("=" * 120)
for i, sh in enumerate(slide.shapes):
    dump(sh, 0, str(i))

# 섹션 구분: top 순으로 정렬해 세로 간격 파악
print("\n" + "=" * 120)
print("top 좌표 순 정렬 (세로 배치 / 간격 파악)")
print("=" * 120)
rows = []
for i, sh in enumerate(slide.shapes):
    rows.append((cm(sh.top), i,
                 str(sh.shape_type).split("(")[0].split(".")[-1].strip(),
                 (sh.text_frame.text.replace("\n", " ")[:30] if sh.has_text_frame else ""),
                 cm(sh.left), cm(sh.width), cm(sh.height)))
for t, i, tp, txt, l, w, h in sorted(rows):
    print(f"  T={t:5.2f} [{i:02d}] {tp:<11} L={l:5.2f} W={w:5.2f} H={h:5.2f} | {txt}")
