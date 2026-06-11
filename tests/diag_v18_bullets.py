import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from lxml import etree
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v18_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[0]


def qn(t):
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    p, e = t.split(':')
    return '{%s}%s' % (ns[p], e)


for sp in slide.shapes:
    if sp.has_text_frame and "낮은 인허가 리스크" in sp.text_frame.text:
        print("대상 shape:", sp.name)
        for pi, p in enumerate(sp.text_frame.paragraphs):
            txt = "".join(r.text for r in p.runs)
            pPr = p._pPr
            buinfo = "기본/없음"
            if pPr is not None:
                if pPr.find(qn('a:buNone')) is not None:
                    buinfo = "buNone(없음)"
                elif pPr.find(qn('a:buChar')) is not None:
                    buinfo = "buChar='%s'" % pPr.find(qn('a:buChar')).get('char')
                elif pPr.find(qn('a:buAutoNum')) is not None:
                    buinfo = "buAutoNum"
            print(f"--- p{pi} level={p.level} bullet={buinfo}")
            print(f"    text='{txt[:50]}'")
            marL = pPr.get('marL') if pPr is not None else None
            indent = pPr.get('indent') if pPr is not None else None
            print(f"    marL={marL} indent={indent}")
            if pPr is not None:
                print("    pPr 전체:", etree.tostring(pPr).decode())
        break

