# -*- coding: utf-8 -*-
"""원본 PDF 금융조건 전문 + 사람 PPT 기초자산개요 셀 + 내 v13 셀 전체 비교."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
ROOT = r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im"

def ct(c):
    return " / ".join(p.text for p in c.text_frame.paragraphs).strip()

def dump_ppt_slides(path, slide_nums, label):
    prs = Presentation(path)
    print(f"\n########## {label} : {os.path.basename(path)} ##########")
    for sn in slide_nums:
        if sn < 1 or sn > len(prs.slides):
            continue
        sl = prs.slides[sn-1]
        # 상단 텍스트
        from pptx.util import Emu
        txt = sorted((Emu(s.top or 0).inches, s.text_frame.text.strip().replace("\n"," "))
                     for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip())
        print(f"\n=== [{label}] 슬라이드 {sn} ===")
        print("  상단:", " | ".join(t for _,t in txt[:2])[:110])
        ti = 0
        for s in sl.shapes:
            if not s.has_table:
                continue
            ti += 1
            t = s.table; nr=len(t.rows); nc=len(t.columns)
            print(f"  ── 표{ti} ({nr}행x{nc}열) ──")
            for r in range(nr):
                cells = [ct(t.cell(r,c)) for c in range(nc)]
                print(f"     r{r}: {cells}")

# 사람 대전 PPT : 기초자산개요 = 슬8,9,10
dump_ppt_slides(os.path.join(ROOT,"[Rainfield] 대전서남터미널_토지담보대출_260511.pptx"),
                [8,9,10], "사람PPT")
# 내 v13 : 기초자산개요 = 슬7,8,9
dump_ppt_slides(os.path.join(ROOT,"변환결과_대전_v14.pptx"),
                [7,8,9,10], "내v14")
