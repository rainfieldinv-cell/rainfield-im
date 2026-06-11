"""idx16 직선→bentConnector3 dry-run. 저장 안 함. 변경 전/후 prst + 좌표 불변 확인."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v8_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]
sp = slide.shapes[16]


def offext(s):
    x = s._element.find('.//' + qn('a:off'))
    e = s._element.find('.//' + qn('a:ext'))
    return (x.get('x'), x.get('y'), e.get('cx'), e.get('cy'))


def ln_info(s):
    ln = s._element.find('.//' + qn('a:ln'))
    w = ln.get('w') if ln is not None else None
    clr = s._element.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
    tail = s._element.find('.//' + qn('a:tailEnd'))
    head = s._element.find('.//' + qn('a:headEnd'))
    return (w, clr.get('val') if clr is not None else None,
            tail.get('type') if tail is not None else None,
            head.get('type') if head is not None else None)


prstGeom = sp._element.find('.//' + qn('a:prstGeom'))
print("변경 전 prst:", prstGeom.get('prst'))
off_before = offext(sp); ln_before = ln_info(sp)

prstGeom.set('prst', 'bentConnector3')
print("변경 후 prst:", prstGeom.get('prst'))
off_after = offext(sp); ln_after = ln_info(sp)

print("\noff/ext 변경 전:", off_before)
print("off/ext 변경 후:", off_after)
print("→ 좌표 동일:", off_before == off_after)
print(f"  (L/T/W/H = {Emu(sp.left).inches:.2f}/{Emu(sp.top).inches:.2f}/"
      f"{Emu(sp.width).inches:.2f}/{Emu(sp.height).inches:.2f})")

print("\nln(두께,색,tailEnd,headEnd) 전:", ln_before)
print("ln(두께,색,tailEnd,headEnd) 후:", ln_after)
print("→ 선 스타일 동일:", ln_before == ln_after)

# 다른 도형 불변 확인
print("\n[다른 커넥터 prst 확인 — 미변경이어야]")
for i in (10, 12, 14):
    g = slide.shapes[i]._element.find('.//' + qn('a:prstGeom'))
    print(f"  idx={i} prst={g.get('prst')}")
print("[박스 idx4~7]")
for i in range(4, 8):
    s = slide.shapes[i]
    print(f"  idx={i} has_table={s.has_table} L={Emu(s.left).inches:.2f} T={Emu(s.top).inches:.2f}")

print("\n[저장 안 함 — 메모리 변경만]")
