"""v8 좌표 수정 dry-run — 화살표/라벨 5개 좌표만 변경, 저장 안 함. 변경 전→후 출력."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Emu

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v7_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]
sp = slide.shapes


def ltwh(s):
    return (Emu(s.left).inches, Emu(s.top).inches, Emu(s.width).inches, Emu(s.height).inches)


targets = {
    10: (5.25, 3.00, 0.00, 0.50),
    12: (2.50, 4.00, 2.00, 0.00),
    14: (6.00, 4.00, 2.00, 0.00),
    16: (6.00, 2.70, 2.00, 0.70),
    17: (6.30, 2.45, 1.95, 0.30),
}

before = {i: ltwh(sp[i]) for i in targets}

for i, (l, t, w, h) in targets.items():
    sp[i].left, sp[i].top, sp[i].width, sp[i].height = Inches(l), Inches(t), Inches(w), Inches(h)

after = {i: ltwh(sp[i]) for i in targets}

label = {10: "세로 신탁사↔차주", 12: "가로 시공사↔차주", 14: "가로 차주↔대주",
         16: "사선 신탁사→대주", 17: "라벨 담보신탁 우선수익권"}
print(f"{'idx':>3} {'설명':<18} {'전 L/T/W/H':<26} → {'후 L/T/W/H':<26} 변경")
print("-" * 90)
for i in targets:
    b = "/".join(f"{v:.2f}" for v in before[i])
    a = "/".join(f"{v:.2f}" for v in after[i])
    changed = "동일" if before[i] == after[i] else "변경됨"
    print(f"{i:>3} {label[i]:<18} {b:<26} → {a:<26} {changed}")

# 박스(표) idx4~7 불변 확인
print("\n[박스 표 idx4~7 좌표 — 불변 확인]")
for i in range(4, 8):
    l, t, w, h = ltwh(sp[i])
    head = sp[i].table.cell(0, 0).text.strip() if sp[i].has_table else ""
    print(f"  idx={i} {head:<6} L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

print("\n[참고] 저장하지 않음 (메모리 상에서만 변경).")
