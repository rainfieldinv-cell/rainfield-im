"""틀→채움 시제품 (A개선+B): 8p 사업개요 → 2열(구분|내용) 정규화 표를 내용까지 채우고
   제목/소제목/인트로의 템플릿 샘플 텍스트를 이 페이지 내용으로 교체. 표 서식은 style_table(테두리/세로중앙 포함)."""
import sys, io, os, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
import fitz
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)
from modules import page_builders
from modules.page_builders import create_presentation_from_template, finalize_presentation, clone_slide_layout
from modules.ppt_generator import save_presentation
from modules.ai_slide_builders import style_table, PALETTE, add_footer, _replace_text_keep_runs, _find_shape_by_pos

PDF = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUT = os.path.join(ROOT, "test_frame_p8_filled5.pptx")
ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
if ov:
    page_builders.LAYOUT_PPTX_PATH = ov


# ── 8p 표 추출(fitz) ──
doc = fitz.open(PDF)
page = doc[7]
tbl = None
for t in page.find_tables().tables:
    ext = t.extract()
    if ext and len(ext) >= 5:
        tbl = ext
        break
doc.close()


def to_label_value(rows):
    """과분할된 표를 2열(구분|내용)로 정규화. 각 행: 첫 비어있지않은셀=구분, 나머지 합침=내용('•' 제거)."""
    out = []
    for row in rows:
        cells = [str(c or "").strip() for c in row]
        ne = [c for c in cells if c and c != "•"]
        if not ne:
            continue
        label = ne[0]
        value = " ".join(ne[1:]).replace("•", "").strip()
        out.append([label, value])
    return out


norm = to_label_value(tbl)
print(f"[8p 표 정규화] {len(tbl)}행(원본) → {len(norm)}행 × 2열")
for r in norm:
    print(f"   {r[0]:<12} | {r[1][:40]}")

# ── 슬라이드 생성 ──
prs = create_presentation_from_template()
n = len(prs.slides)
slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)

# 제목/소제목/인트로 — 템플릿 샘플 텍스트를 8p 내용으로 교체
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if "사모사채개요" in t and t[:2].strip().isdigit() is False and t.startswith("0"):
        _replace_text_keep_runs(sh.text_frame, "03  본건 사업 개요")     # 섹션 라벨
    elif "본 건 사모사채 개요" in t or "사모사채 개요" in t:
        _replace_text_keep_runs(sh.text_frame, "3.1  사업 개요")          # 소제목
    elif t.startswith("본건 사모사채는"):
        _replace_text_keep_runs(sh.text_frame, "본건 사업의 개요는 아래와 같다.")  # 인트로

# ── 표 채움 (2열, 헤더=첫 행) ──
rows = len(norm)
gf = slide.shapes.add_table(rows, 2, Inches(0.6), Inches(2.0), Inches(9.6), Inches(max(rows * 0.42, 1.0)))
t2 = gf.table
t2.columns[0].width = Inches(2.6)
t2.columns[1].width = Inches(7.0)
for ri, (lab, val) in enumerate(norm):
    _replace_text_keep_runs(t2.cell(ri, 0).text_frame, lab)
    _replace_text_keep_runs(t2.cell(ri, 1).text_frame, val)
# col1(값)은 데이터 행만 왼쪽정렬(헤더 행은 style_table이 가운데로 통일)
from pptx.enum.text import PP_ALIGN
for ri in range(1, rows):
    t2.cell(ri, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

style_table(gf, has_header=True, label_cols=(0,),
            header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])

add_footer(slide, 1, business_name="천안 부성2지구 도시개발사업")
finalize_presentation(prs, n)
save_presentation(prs, OUT)

print(f"\n저장: {OUT}")
dup = [nm for nm, c in Counter(zipfile.ZipFile(OUT).namelist()).items() if c > 1]
print(f"중복파트: {'있음' if dup else '없음(정상)'}")

# 검증 재오픈
from pptx import Presentation
from pptx.oxml.ns import qn
p2 = Presentation(OUT)
sl = p2.slides[0]
print("\n[제목/소제목/인트로]")
for sh in sl.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() and not sh.has_table:
        print("   ", sh.text_frame.text.strip()[:40])
for sh in sl.shapes:
    if sh.has_table:
        tb = sh.table
        c0 = tb.cell(0, 0)
        tcPr = c0._tc.find(qn('a:tcPr'))
        ln = tcPr.find(qn('a:lnL')) if tcPr is not None else None
        bcol = ln.find('.//' + qn('a:srgbClr')).get('val') if ln is not None and ln.find('.//' + qn('a:srgbClr')) is not None else "?"
        print(f"\n[표] {len(tb.rows)}행×{len(tb.columns)}열, 헤더fill={c0.fill.fore_color.rgb if c0.fill.type else '-'}, 테두리={bcol}/{int(ln.get('w'))/12700:.2f}pt")
        for ri in range(min(4, len(tb.rows))):
            print("   ", [tb.cell(ri, ci).text.strip()[:28] for ci in range(len(tb.columns))])
