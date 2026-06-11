import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v14_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[0]
for sh in slide.shapes:
    if sh.name in ("TextBox 42", "TextBox 26", "TextBox 50", "육각형 5", "육각형 17",
                   "육각형 24", "육각형 25", "육각형 48", "육각형 49"):
        info = ""
        if sh.has_text_frame:
            tf = sh.text_frame
            try:
                anc = tf.vertical_anchor
            except Exception:
                anc = "?"
            sizes = []
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        sizes.append(r.font.size.pt if r.font.size else None)
            info = f"anchor={anc} word_wrap={tf.word_wrap} sizes={sizes[:3]}"
        print(f"{sh.name:<12} L={Emu(sh.left).inches:.2f} T={Emu(sh.top).inches:.2f} "
              f"W={Emu(sh.width).inches:.2f} H={Emu(sh.height).inches:.2f}  {info}")
