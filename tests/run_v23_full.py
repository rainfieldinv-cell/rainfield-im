"""run_v23_full.py — 천안 PDF 전체를 build_full_presentation 으로 변환 + 표 채움 검증."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT = os.path.join(ROOT, "test_ai_builders_v23.pptx")
BIZ = "천안 부성2지구 도시개발사업"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def fill_hex(cell):
    tcPr = cell._tc.find("{%s}tcPr" % A)
    if tcPr is not None:
        sf = tcPr.find("{%s}solidFill" % A)
        if sf is not None:
            sc = sf.find("{%s}srgbClr" % A)
            if sc is not None:
                return sc.get("val")
    return "-"


def main():
    from modules import page_builders
    from modules.content_parser import parse_document
    from modules.page_builders import build_full_presentation

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    print("PDF 파싱...")
    pages = parse_document(PDF)
    tbl_total = sum(len(p.get("tables") or []) for p in pages)
    print(f"파싱 완료: {len(pages)} 페이지, 추출 표 {tbl_total}개")

    ppt_bytes = build_full_presentation(
        business_name=BIZ, year="2026", month_en="May", pages=pages,
    )
    with open(OUTPUT, "wb") as f:
        f.write(ppt_bytes)
    print(f"저장: {OUTPUT} ({len(ppt_bytes):,} bytes)")

    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v23] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    # 재오픈(읽기) 검증 — 슬라이드별 표/채움 셀
    from pptx import Presentation
    prs = Presentation(OUTPUT)
    print(f"\n총 슬라이드 {len(prs.slides)}장. 표 포함 슬라이드:")
    print(f"{'sl':>3} | 표 | (행x열, 채움셀/전체)")
    keyword_hits = {}
    for si, sl in enumerate(prs.slides):
        tbls = [sh for sh in sl.shapes if sh.has_table]
        if not tbls:
            continue
        parts = []
        for sh in tbls:
            t = sh.table
            r, c = len(t.rows), len(t.columns)
            filled = sum(1 for ri in range(r) for ci in range(c) if t.cell(ri, ci).text.strip())
            parts.append(f"{r}x{c} {filled}/{r*c}")
            # 키워드 표 식별
            joined = " ".join(t.cell(ri, ci).text for ri in range(r) for ci in range(c))
            for kw, label in [("계좌", "계좌표"), ("Cash", "Cash흐름"), ("주주", "주주표"),
                              ("총 자산", "재무제표"), ("사모사채명", "사모사채표"),
                              ("Cash-in", "Cash-in/out")]:
                if kw in joined and label not in keyword_hits:
                    keyword_hits[label] = (si, sh)
        print(f"{si:>3} | {len(tbls):>2} | " + "  ".join(parts))

    # 주요 표 셀 내용 일부
    print("\n[주요 표 셀 내용 확인]")
    for label, (si, sh) in keyword_hits.items():
        t = sh.table
        print(f"\n--- {label} (슬라이드 {si}, {len(t.rows)}x{len(t.columns)}) 헤더fill={fill_hex(t.cell(0,0))} ---")
        for ri in range(min(5, len(t.rows))):
            print("   ", [t.cell(ri, ci).text.strip()[:14] for ci in range(len(t.columns))])


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
