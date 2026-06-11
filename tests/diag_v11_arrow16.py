"""읽기 전용 — v11 idx16 ㄱ자 커넥터 시작/끝점·flip + idx7/14 좌표 진단."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v11_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]

for i in [7, 14, 16]:
    sp = slide.shapes[i]
    L, T, W, H = (Emu(sp.left).inches, Emu(sp.top).inches,
                  Emu(sp.width).inches, Emu(sp.height).inches)
    print(f"idx={i} {sp.name}")
    print(f"  bbox: L={L:.2f} T={T:.2f} W={W:.2f} H={H:.2f}  "
          f"(오른쪽끝 X={L+W:.2f}, 아래끝 Y={T+H:.2f})")
    xfrm = sp._element.find('.//' + qn('a:xfrm'))
    if xfrm is not None:
        print(f"  flipH={xfrm.get('flipH')} flipV={xfrm.get('flipV')} rot={xfrm.get('rot')}")
    g = sp._element.find('.//' + qn('a:prstGeom'))
    if g is not None:
        adj = g.find(qn('a:avLst') + '/' + qn('a:gd'))
        print(f"  prst={g.get('prst')} adj1={adj.get('fmla') if adj is not None else '없음'}")

print("\n[대주 박스 idx7] 왼쪽변 X=8.00, Y범위 3.20~4.60, 중앙Y=3.90")
print("[대출약정 idx14] 가로, Y=4.00 → ㄱ자 끝점이 여기랑 겹치면 안 됨")
