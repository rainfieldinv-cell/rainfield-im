import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v20_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[1]
for i, sp in enumerate(slide.shapes):
    txt = ""
    if sp.has_text_frame:
        txt = sp.text_frame.text.replace("\n", " / ")[:40]
    elif sp.has_table:
        txt = "[TABLE]"
    print(f"idx={i} type={sp.shape_type} name={sp.name} "
          f"L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
          f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f} <{txt}>")
    if sp.has_table:
        tbl = sp.table
        print(f"   표크기 L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
              f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f}, 행수={len(tbl.rows)}, 열수={len(tbl.columns)}")
        for ri, row in enumerate(tbl.rows):
            c0 = row.cells[0].text[:18]
            c1 = row.cells[1].text[:22] if len(row.cells) > 1 else ""
            print(f"      row{ri} H={Emu(row.height).inches:.2f} col0='{c0}' col1='{c1}'")

print("\n슬라이드 세로:", Emu(prs.slide_height).inches)
# 인트로 본문 끝 Y, 푸터 T
for sp in slide.shapes:
    if sp.has_text_frame and sp.text_frame.text.strip().startswith("본건"):
        print(f"인트로 본문: T={Emu(sp.top).inches:.2f} H={Emu(sp.height).inches:.2f} → 끝 Y={Emu(sp.top).inches + Emu(sp.height).inches:.2f}")
    if sp.has_text_frame and "|" in sp.text_frame.text:
        print(f"푸터: T={Emu(sp.top).inches:.2f} '{sp.text_frame.text.strip()}'")
