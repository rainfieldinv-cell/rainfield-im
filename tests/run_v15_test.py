"""run_v15_test.py — 빌더 재생성(save_presentation)으로 v15 + slides[0] 잡요소 삭제 검증."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v15.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장 완료: {OUTPUT}")

    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v15] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    sl = p2.slides[0]
    print(f"\n[slides[0] 전체 {len(sl.shapes)} shape]")
    names = []
    for i, sp in enumerate(sl.shapes):
        tp = str(sp.shape_type).split("(")[0].split(".")[-1].strip()
        l = Emu(sp.left).inches if sp.left is not None else 0
        t = Emu(sp.top).inches if sp.top is not None else 0
        w = Emu(sp.width).inches if sp.width is not None else 0
        h = Emu(sp.height).inches if sp.height is not None else 0
        txt = sp.text_frame.text.replace("\n", " / ")[:60] if sp.has_text_frame else ""
        names.append(sp.name)
        print(f"  idx={i:2d} {tp:<11} {sp.name:<18} L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f} <{txt}>")

    print("\n[삭제 대상 4개 부재 확인]")
    for nm in ["그룹 2", "그룹 51", "TextBox 56", "TextBox 54"]:
        print(f"  '{nm}': {'남아있음 FAIL' if nm in names else '삭제됨 OK'}")
    print("\n[유지 대상 존재 확인]")
    for nm in ["그룹 1", "TextBox 55", "TextBox 26", "TextBox 42", "TextBox 50"]:
        print(f"  '{nm}': {'존재 OK' if nm in names else '없음 FAIL'}")
    hexes = [n for n in names if "육각형" in n]
    print(f"  육각형(<>) 개수: {len(hexes)}개 {hexes}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
