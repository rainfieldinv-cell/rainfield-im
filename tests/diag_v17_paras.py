import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v17_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[0]
for nm in ("TextBox 42", "TextBox 26", "TextBox 50"):
    sh = next((s for s in slide.shapes if s.name == nm), None)
    if not sh:
        continue
    tf = sh.text_frame
    paras = tf.paragraphs
    print(f"=== {nm}: 단락 {len(paras)}개 ===")
    for i, p in enumerate(paras):
        t = "".join(r.text for r in p.runs)
        print(f"  p{i}: runs={len(p.runs)} len={len(t)} '{t[:40]}'")
