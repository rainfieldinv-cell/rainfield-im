"""run_v11_test.py — 수정된 빌더로 v11 빌드(save_presentation) + idx16/중복파트 검사."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v11.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)   # raw prs.save() 아님 — 빌더 저장 경로
    print(f"저장 완료: {OUTPUT}")

    # 중복파트 자체검사
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v11] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    # idx16 검증
    p2 = Presentation(OUTPUT)
    sp = p2.slides[2].shapes
    s = sp[16]
    g = s._element.find('.//' + qn('a:prstGeom'))
    adj = g.find(qn('a:avLst') + '/' + qn('a:gd'))
    clr = s._element.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
    ln = s._element.find('.//' + qn('a:ln'))
    tail = s._element.find('.//' + qn('a:tailEnd'))
    print(f"\n[v11 idx16] prst={g.get('prst')} adj1={adj.get('fmla') if adj is not None else '없음'}")
    print(f"  L/T/W/H = {Emu(s.left).inches:.2f}/{Emu(s.top).inches:.2f}/"
          f"{Emu(s.width).inches:.2f}/{Emu(s.height).inches:.2f}")
    print(f"  선 두께={ln.get('w')} 색={clr.get('val') if clr is not None else None} "
          f"tailEnd={tail.get('type') if tail is not None else None}")

    print("\n[다른 커넥터 prst]")
    for i in (10, 12, 14):
        gg = sp[i]._element.find('.//' + qn('a:prstGeom'))
        print(f"  idx={i} prst={gg.get('prst')}")
    print("[박스 idx4~7]")
    for i in range(4, 8):
        print(f"  idx={i} has_table={sp[i].has_table} L={Emu(sp[i].left).inches:.2f} T={Emu(sp[i].top).inches:.2f}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
