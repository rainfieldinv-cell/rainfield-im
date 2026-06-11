"""
verify_v5.py — 저장된 test_ai_builders_v5.pptx 를 직접 열어 슬라이드 7 정밀 검증.
(run_v5_test 검증 로직의 오탐 수정: shape_type 정확 비교 + 금액 체크를 박스/라벨로 한정)
"""
import sys, os, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(os.path.join(ROOT, "test_ai_builders_v5.pptx"))
s7 = prs.slides[2]   # S2, S5, S7 순서 → 3번째

BOX_ROLES = {"신탁사", "시공사", "차주", "대주"}
COMPANIES = {"신한자산신탁", "포스코이앤씨", "더함도시개발", "Tr.A", "Tr.B"}
LABELS = {"신탁계약", "담보신탁 우선수익권", "공사도급계약", "대출약정"}

autoshapes, connectors = [], []
box_texts, label_texts = [], []
diagram_fonts = set()

for sh in s7.shapes:
    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        autoshapes.append(sh)
        t = sh.text_frame.text.strip()
        box_texts.append(t)
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() and r.font.name:
                    diagram_fonts.add(r.font.name)
    if "Connector" in (sh.name or ""):
        connectors.append(sh)
    if sh.has_text_frame and sh.text_frame.text.strip() in LABELS:
        label_texts.append(sh.text_frame.text.strip())
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() and r.font.name:
                    diagram_fonts.add(r.font.name)

roles_found = sorted(set(box_texts) & BOX_ROLES)
comp_found = sorted(set(box_texts) & COMPANIES)
labels_found = sorted(set(label_texts))

# 금액 표기: 박스 셀 + 라벨 텍스트에만 한정 (인트로 문장 제외)
diagram_only = box_texts + label_texts
amount_in_diagram = [t for t in diagram_only if ("억" in t or "원" in t)]

# 금지 요소 (박스/라벨 한정)
asset_mgr = [t for t in diagram_only if "자산관리" in t]
investor = [t for t in diagram_only if any(k in t for k in
            ["투자자", "Equity", "삼성증권", "KT Estate", "우미건설", "금성백조"])]
pics = [sh for sh in s7.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]

print("=" * 90)
print("PDF 천안 부성2지구 p4 구조도 시각 일치 — 정밀 재검증")
print("=" * 90)
checks = [
    ("박스 4개 (신탁사/시공사/차주/대주)", "있음", f"{len(roles_found)}개 {roles_found}", len(roles_found) == 4),
    ("  └ 회사명 5개", "있음", f"{len(comp_found)}개 {comp_found}", len(comp_found) == 5),
    ("화살표(연결선) 3개", "있음", f"{len(connectors)}개", len(connectors) == 3),
    ("라벨 4개", "있음", f"{labels_found}", len(labels_found) == 4),
    ("폰트: 피플폰트 Bold", "적용", f"{sorted(diagram_fonts)}", diagram_fonts == {"피플폰트 Bold"}),
    ("금액 표기 없음 (박스/라벨)", "없음", f"{amount_in_diagram or '없음'}", len(amount_in_diagram) == 0),
    ("자산관리자 박스 없음", "없음", f"{asset_mgr or '없음'}", len(asset_mgr) == 0),
    ("사업명 이미지(PICTURE) 없음", "없음", f"{len(pics)}개", len(pics) == 0),
    ("우측 투자자 박스 없음", "없음", f"{investor or '없음'}", len(investor) == 0),
]
for name, pdf, actual, ok in checks:
    print(f"  [{'OK ' if ok else 'FAIL'}] {name:<34} | PDF={pdf:<4} | v5={actual}")
print(f"\n  결과: {sum(1 for *_, ok in checks if ok)}/{len(checks)} PASS")

# 참고: 인트로 문장의 금액(정상)
intro = [sh.text_frame.text.strip() for sh in s7.shapes
         if sh.has_text_frame and sh.text_frame.text.strip().startswith("본건은")]
if intro:
    print(f"\n[참고] 인트로 문장(베이스 유지, PDF 동일)에는 금액이 정상 포함됨:")
    print(f"   '{intro[0][:80]}...'")
