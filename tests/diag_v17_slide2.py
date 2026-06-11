import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v17_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]
tops = []
for i, sp in enumerate(slide.shapes):
    txt = ""
    sizes = []
    if sp.has_text_frame:
        txt = sp.text_frame.text.replace("\n", " / ")[:60]
        for p in sp.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sizes.append(r.font.size.pt)
    elif sp.has_table:
        txt = "[TABLE]"
        for row in sp.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            sizes.append(r.font.size.pt)
    tops.append(Emu(sp.top).inches)
    print(f"idx={i} type={sp.shape_type} name={sp.name} "
          f"L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
          f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f} "
          f"font={sorted(set(sizes))} <{txt}>")
print("총 shape:", len(slide.shapes))
print("T 최소:", round(min(tops), 2), "T 최대:", round(max(tops), 2))
print("슬라이드 세로:", Emu(prs.slide_height).inches)

# freeform(ㄱ자) 정점 확인
print("\n[freeform/커넥터 정점·prst 확인]")
for i, sp in enumerate(slide.shapes):
    tp = str(sp.shape_type).split("(")[0].split(".")[-1].strip()
    if tp == "FREEFORM":
        path = sp._element.find('.//' + qn('a:custGeom') + '/' + qn('a:pathLst') + '/' + qn('a:path'))
        L0, T0 = Emu(sp.left).inches, Emu(sp.top).inches
        pts = []
        if path is not None:
            for seg in path:
                pt = seg.find(qn('a:pt'))
                if pt is not None:
                    pts.append((round(L0 + int(pt.get('x')) / 914400, 2),
                                round(T0 + int(pt.get('y')) / 914400, 2)))
        print(f"  idx={i} FREEFORM 정점={pts}")
    elif tp == "LINE" and "Connector" in (sp.name or ""):
        g = sp._element.find('.//' + qn('a:prstGeom'))
        print(f"  idx={i} {sp.name} prst={g.get('prst') if g is not None else '?'} "
              f"L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
              f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f}")
