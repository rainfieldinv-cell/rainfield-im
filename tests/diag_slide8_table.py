"""
diag_slide8_table.py — 레이아웃 슬라이드 8(투자구조도) 표 구조 정밀 진단.
코드 수정 없음. 잠금 회피를 위해 인자로 받은 pptx 경로(로컬 복사본)를 직접 연다.
"""
import sys, os, io
from lxml import etree
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

src = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_layout_v5.pptx"
prs = Presentation(src)
print(f"열기: {src}  (총 {len(prs.slides)} 슬라이드)")

# investment_structure = SLIDE_INDEX_MAP 기준 index 7 (slide #8)
idx = 7
slide = prs.slides[idx]
print(f"investment_structure = 슬라이드 index {idx} (slide #{idx+1})\n")


def hexfill(cell):
    """셀 배경색 추출."""
    try:
        f = cell.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return "-"


print("=" * 110)
print(f"슬라이드 8 전체 shape — has_table / GRAPHIC_FRAME 여부")
print("=" * 110)
tables = []
for i, sh in enumerate(slide.shapes):
    tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
    ht = sh.has_table
    l = sh.left/360000 if sh.left is not None else 0
    t = sh.top/360000 if sh.top is not None else 0
    w = sh.width/360000 if sh.width is not None else 0
    h = sh.height/360000 if sh.height is not None else 0
    extra = ""
    if ht:
        tb = sh.table
        extra = f"  TABLE {len(tb.rows)}r×{len(tb.columns)}c"
        tables.append((i, sh))
    elif sh.has_text_frame:
        extra = f"  '{sh.text_frame.text.strip()[:20]}'"
    print(f"[{i:02d}] {tp:<14} has_table={ht!s:<5} L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f}{extra}")

print(f"\n총 GRAPHIC_FRAME(table): {len(tables)}개  → idx {[i for i,_ in tables]}")

# "구조도" 라벨 탐색
print(f"\n{'='*110}\n'구조도' 텍스트 탐색 (좌측 라벨 칸 여부)\n{'='*110}")
hits = []
for i, sh in enumerate(slide.shapes):
    if sh.has_text_frame and "구조도" in sh.text_frame.text:
        hits.append(f"[{i}] TEXT_BOX '{sh.text_frame.text.strip()}'")
    if sh.has_table:
        tb = sh.table
        for r in range(len(tb.rows)):
            for c in range(len(tb.columns)):
                if "구조도" in tb.cell(r, c).text:
                    hits.append(f"[{i}] TABLE cell[{r},{c}] '{tb.cell(r,c).text.strip()}'")
print("\n".join(hits) if hits else "  '구조도' 라벨 못 찾음 (제목 '2.1 투자구조도' 제외)")

# 각 표 상세
print(f"\n{'='*110}\n표(table) 상세 — 셀 내용 / 크기 / 배경색\n{'='*110}")
for i, sh in tables:
    tb = sh.table
    l = sh.left/360000; t = sh.top/360000; w = sh.width/360000; h = sh.height/360000
    print(f"\n■ shape[{i}] TABLE {len(tb.rows)}r×{len(tb.columns)}c  L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")
    for r in range(len(tb.rows)):
        rh = tb.rows[r].height/360000
        for c in range(len(tb.columns)):
            cell = tb.cell(r, c)
            print(f"   cell[{r},{c}] rowH={rh:.2f}  배경={hexfill(cell):<8} 텍스트='{cell.text.strip()[:30]}'")

# 첫 표 raw XML 인용
if tables:
    i0, sh0 = tables[0]
    print(f"\n{'='*110}\nshape[{i0}] raw XML (graphicFrame, 일부)\n{'='*110}")
    xml = etree.tostring(sh0._element, pretty_print=True).decode()
    print(xml[:2500])
