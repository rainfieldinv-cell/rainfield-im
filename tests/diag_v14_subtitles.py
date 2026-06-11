import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v14_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[0]

targets = [4, 15, 16, 17, 3, 14]
for i in targets:
    sp = slide.shapes[i]
    print(f"--- idx={i} type={sp.shape_type} name={sp.name} ---")
    if sp.has_text_frame:
        for p in sp.text_frame.paragraphs:
            for r in p.runs:
                c = r.font.color
                rgb = c.rgb if c and c.type is not None else "None"
                print(f"   text='{r.text}' size={r.font.size} color={rgb} font={r.font.name}")
    # GROUP 내부도 확인
    if str(sp.shape_type).startswith("GROUP") or "GROUP" in str(sp.shape_type):
        for j, ch in enumerate(sp.shapes):
            t = ch.text_frame.text if ch.has_text_frame else ""
            print(f"   child[{j}] {ch.shape_type} name={ch.name} <{t[:30]}>")
    print()

print("=== 작은 도형(체크박스 후보) ===")
for i, sp in enumerate(slide.shapes):
    if sp.shape_type and "AUTO_SHAPE" in str(sp.shape_type):
        w = Emu(sp.width).inches; h = Emu(sp.height).inches
        if w < 0.5 and h < 0.5:
            print(f"idx={i} name={sp.name} L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} W={w:.2f} H={h:.2f}")

# GROUP 내부 PICTURE(체크마크 이미지) 후보
print("\n=== GROUP 내부 구성(체크마크 후보) ===")
for i in [3, 4, 14]:
    sp = slide.shapes[i]
    print(f"GROUP idx={i} name={sp.name}")
    for j, ch in enumerate(sp.shapes):
        t = ch.text_frame.text[:25] if ch.has_text_frame else ""
        print(f"   child[{j}] {ch.shape_type} name={ch.name} "
              f"L={Emu(ch.left).inches:.2f} T={Emu(ch.top).inches:.2f} <{t}>")
