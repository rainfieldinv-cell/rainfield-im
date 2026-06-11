"""읽기 전용 — v13 idx16 freeform 전체 XML + 박스 좌표. 저장 금지."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v13_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]

sp = slide.shapes[16]
print("=== idx16", sp.shape_type, sp.name, "===")
print(f"bbox L/T/W/H: {Emu(sp.left).inches:.3f} {Emu(sp.top).inches:.3f} "
      f"{Emu(sp.width).inches:.3f} {Emu(sp.height).inches:.3f}")

# path 정점 + 인치 환산 (path-local EMU → 인치)
path = sp._element.find('.//' + qn('a:custGeom') + '/' + qn('a:pathLst') + '/' + qn('a:path'))
print("\n[path 정점]")
if path is not None:
    print(f"  path w={path.get('w')} h={path.get('h')}")
    for seg in path:
        nm = etree.QName(seg).localname
        pt = seg.find(qn('a:pt'))
        if pt is not None:
            x = int(pt.get('x')); y = int(pt.get('y'))
            print(f"  {nm}: x={x} y={y}  → +bbox = ({Emu(sp.left).inches + x/914400:.3f}, "
                  f"{Emu(sp.top).inches + y/914400:.3f}) in")
        else:
            print(f"  {nm}")

# 선 모서리/조인 스타일
ln = sp._element.find('.//' + qn('a:ln'))
if ln is not None:
    print("\n[ln 속성]", "cap=", ln.get('cap'))
    for c in ln:
        print("   ", etree.QName(c).localname,
              {k.split('}')[-1]: v for k, v in c.attrib.items()})

print("\n--- idx16 전체 XML ---")
print(etree.tostring(sp._element, pretty_print=True).decode())

print("\n=== 박스 좌표 ===")
for i in [4, 5, 6, 7]:
    b = slide.shapes[i]
    L, T, W, H = (Emu(b.left).inches, Emu(b.top).inches,
                  Emu(b.width).inches, Emu(b.height).inches)
    head = b.table.cell(0, 0).text.strip() if b.has_table else ""
    print(f"idx={i} {head:<6} L={L:.2f} T={T:.2f} W={W:.2f} H={H:.2f}  "
          f"우변X={L+W:.2f} 하변Y={T+H:.2f} 중앙Y={T+H/2:.2f}")
