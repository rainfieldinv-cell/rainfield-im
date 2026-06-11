"""읽기 전용 진단 — v7 슬라이드 좌표 출력. 파일 수정/저장 없음."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

import os
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im\test_ai_builders_v7.pptx"
prs = Presentation(PATH)
print(f"읽은 파일: {PATH}")

print(f"slide_width  = {Emu(prs.slide_width).inches:.2f} in")
print(f"slide_height = {Emu(prs.slide_height).inches:.2f} in")
print(f"총 슬라이드 수 = {len(prs.slides)}")

# 투자구조도 슬라이드 자동 탐색 ('2.1 투자구조도' 포함)
target = None
for idx, sl in enumerate(prs.slides):
    joined = " ".join(s.text_frame.text for s in sl.shapes if s.has_text_frame)
    if "투자구조도" in joined:
        target = idx
        break
print(f"'2.1 투자구조도' 슬라이드 인덱스 = {target}")


def dump(idx, label):
    print(f"\n===== {label}: slides[{idx}] =====")
    slide = prs.slides[idx]
    for i, sp in enumerate(slide.shapes):
        txt = ""
        if sp.has_text_frame:
            txt = sp.text_frame.text.replace("\n", " / ")[:30]
        elif sp.has_table:
            tbl = sp.table
            txt = " | ".join(c.text for r in tbl.rows for c in r.cells)[:40]
        print(f"idx={i} type={sp.shape_type} "
              f"L={Emu(sp.left).inches:.2f} T={Emu(sp.top).inches:.2f} "
              f"W={Emu(sp.width).inches:.2f} H={Emu(sp.height).inches:.2f}  <{txt}>")
    print("총 shape 개수:", len(slide.shapes))


if target is not None:
    dump(target, "투자구조도(실제)")
# 사용자가 지정한 인덱스 6 도 가능하면 함께 출력
if len(prs.slides) > 6:
    dump(6, "요청 인덱스 6")
else:
    print(f"\n[참고] 요청한 slides[6] 은 존재하지 않음 (이 PPT는 {len(prs.slides)}장).")
