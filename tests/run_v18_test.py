"""run_v18_test.py — 공통 푸터(add_footer) 전 슬라이드 적용. page_num=슬라이드순서(1/2/3)."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v18.pptx")
BIZ = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.util import Emu, Pt

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    # page_num = 최종 슬라이드 순서 (S2→1, S5→2, S7→3)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ, page_num=1)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ, page_num=2)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ, page_num=3)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}")
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v18] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    print(f"\n총 슬라이드: {len(p2.slides)}")
    for si, sl in enumerate(p2.slides):
        # 푸터(피플폰트 Light 회색) 텍스트박스 찾기
        footers = []
        leftover = []
        for sp in sl.shapes:
            if "슬라이드 번호" in (sp.name or ""):
                leftover.append(sp.name)
            if sp.has_text_frame:
                t = sp.text_frame.text.strip()
                if "|" in t and BIZ in t:
                    # 폰트/색/여백
                    r = sp.text_frame.paragraphs[0].runs[0]
                    col = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
                    ml = sp.text_frame.margin_left
                    footers.append((t, r.font.name, r.font.size.pt if r.font.size else None,
                                    str(col), Emu(sp.left).inches, Emu(sp.top).inches,
                                    Emu(sp.width).inches, ml))
        print(f"\n[슬라이드 {si} (= page {si+1})] shape {len(sl.shapes)}개")
        for f in footers:
            print(f"   오른쪽 푸터: '{f[0]}' font={f[1]} {f[2]}pt color={f[3]} "
                  f"L={f[4]:.2f} T={f[5]:.2f} W={f[6]:.2f} margin_left={f[7]}")
        print(f"   잔존 '슬라이드 번호' 박스: {leftover if leftover else '없음(정리됨)'}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
