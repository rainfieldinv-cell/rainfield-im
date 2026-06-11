"""유형 E 시제품: 8p → [데이터 표 | 이미지 표] 2단.
   - 왼쪽: 구분|내용 데이터 표 (위치도/조감도 행 제외)
   - 오른쪽: 구분|내용 이미지 표 (내용칸에 추출한 조감도·위치도 JPEG 삽입)"""
import sys, io, os, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
import fitz
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)
from modules import page_builders
from modules.page_builders import create_presentation_from_template, finalize_presentation, clone_slide_layout
from modules.ppt_generator import save_presentation
from modules.ai_slide_builders import style_table, PALETTE, add_footer, _replace_text_keep_runs

PDF = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUT = os.path.join(ROOT, "test_frame_p8_E2.pptx")
ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
if ov:
    page_builders.LAYOUT_PPTX_PATH = ov

IMG_LABELS = {"위치도", "조감도", "조 감 도", "위 치 도"}

# ── 8p 표 + 이미지 추출 ──
doc = fitz.open(PDF)
page = doc[7]
raw_tbl = None
for t in page.find_tables().tables:
    ext = t.extract()
    if ext and len(ext) >= 5:
        raw_tbl = ext
        break
# 이미지: y(top) 기준 위→아래 = 조감도, 위치도
imgs = []
for im in page.get_images(full=True):
    xref = im[0]
    rects = page.get_image_rects(xref)
    top = rects[0][1] if rects else 0
    base = doc.extract_image(xref)
    imgs.append({"y": top, "bytes": base["image"], "ext": base["ext"],
                 "w": base["width"], "h": base["height"]})
imgs.sort(key=lambda d: d["y"])     # 위→아래
doc.close()


def to_label_value(rows):
    out = []
    for row in rows:
        cells = [str(c or "").strip() for c in row]
        ne = [c for c in cells if c and c != "•"]
        if not ne:
            continue
        out.append([ne[0], " ".join(ne[1:]).replace("•", "").strip()])
    return out


norm = to_label_value(raw_tbl)
data_rows = [r for r in norm if r[0].replace(" ", "") not in {x.replace(" ", "") for x in IMG_LABELS}]
print(f"[8p] 데이터 행 {len(data_rows)}개(이미지행 제외), 이미지 {len(imgs)}장")
for r in data_rows:
    print(f"   {r[0]:<10} | {r[1][:34]}")

# ── 슬라이드 ──
prs = create_presentation_from_template()
n = len(prs.slides)
slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)

# 제목/소제목/인트로 교체
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("01") or "사모사채개요" in t:
        _replace_text_keep_runs(sh.text_frame, "03  본건 사업 개요")
    elif "사모사채 개요" in t:
        _replace_text_keep_runs(sh.text_frame, "3.1  사업 개요")
    elif t.startswith("본건 사모사채는"):
        _replace_text_keep_runs(sh.text_frame, "본건 사업의 개요는 아래와 같다.")

# ── 소제목바 "자산 개요 (1/3)" (사람 슬9: L0.43 T1.52) ──
sub = slide.shapes.add_textbox(Inches(0.43), Inches(1.5), Inches(2.5), Inches(0.3))
_sp = sub.text_frame.paragraphs[0].add_run()
_sp.text = "자산 개요 (1/3)"
_sp.font.name = "피플폰트 Bold"
_sp.font.size = Pt(12)
_sp.font.color.rgb = PALETTE["navy_dark"]

# ── 왼쪽: 데이터 표 (사람 슬9: L0.43 W4.99) ──
L_L, L_T, L_W = 0.43, 1.83, 4.99
rows = len(data_rows)
gL = slide.shapes.add_table(rows, 2, Inches(L_L), Inches(L_T), Inches(L_W), Inches(rows * 0.45))
tL = gL.table
tL.columns[0].width = Inches(1.40)
tL.columns[1].width = Inches(3.59)
for ri, (lab, val) in enumerate(data_rows):
    _replace_text_keep_runs(tL.cell(ri, 0).text_frame, lab)
    _replace_text_keep_runs(tL.cell(ri, 1).text_frame, val)
    if ri >= 1:
        tL.cell(ri, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
style_table(gL, has_header=True, label_cols=(0,),
            header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])

# ── 오른쪽: 이미지 표 (사람 슬9: L5.56 W4.85) ──
R_L, R_T, R_W = 5.56, 1.83, 4.85
labels = ["조감도", "위치도"]
img_rows = len(imgs)
ROW_H = 2.30
gR = slide.shapes.add_table(img_rows + 1, 2, Inches(R_L), Inches(R_T), Inches(R_W), Inches(0.4 + img_rows * ROW_H))
tR = gR.table
tR.columns[0].width = Inches(0.95)
tR.columns[1].width = Inches(3.90)
_replace_text_keep_runs(tR.cell(0, 0).text_frame, "구 분")
_replace_text_keep_runs(tR.cell(0, 1).text_frame, "내 용")
tR.rows[0].height = Inches(0.4)
for i in range(img_rows):
    _replace_text_keep_runs(tR.cell(i + 1, 0).text_frame, labels[i] if i < len(labels) else "사진")
    tR.rows[i + 1].height = Inches(ROW_H)
style_table(gR, has_header=True, label_cols=(0,),
            header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])

# 이미지 삽입 (내용칸 위에 비율 유지·중앙)
col0_w = 0.95
content_L = R_L + col0_w
content_W = 3.90
for i, im in enumerate(imgs):
    cell_T = R_T + 0.4 + i * ROW_H
    cell_H = ROW_H
    pad = 0.1
    boxL, boxT = content_L + pad, cell_T + pad
    boxW, boxH = content_W - 2 * pad, cell_H - 2 * pad
    scale = min(boxW / im["w"], boxH / im["h"])
    iw, ih = im["w"] * scale, im["h"] * scale
    px = boxL + (boxW - iw) / 2
    py = boxT + (boxH - ih) / 2
    slide.shapes.add_picture(io.BytesIO(im["bytes"]), Inches(px), Inches(py),
                             Inches(iw), Inches(ih))

add_footer(slide, 1, business_name="천안 부성2지구 도시개발사업")
finalize_presentation(prs, n)
save_presentation(prs, OUT)
print(f"\n저장: {OUT}")
dup = [nm for nm, c in Counter(zipfile.ZipFile(OUT).namelist()).items() if c > 1]
print(f"중복파트: {'있음' if dup else '없음(정상)'}")

from pptx import Presentation
p2 = Presentation(OUT)
sl = p2.slides[0]
ntab = sum(1 for s in sl.shapes if s.has_table)
npic = sum(1 for s in sl.shapes if s.shape_type == 13)
print(f"검증: 표 {ntab}개, 사진 {npic}장")
