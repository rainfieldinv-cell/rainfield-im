"""run_v8_test.py — 수정된 빌더로 v8 새로 빌드 + slides[2] 화살표/라벨 좌표 검증."""
import sys, os, io, traceback
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v8.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"저장 완료: {OUTPUT}")

    # 저장된 v8 재오픈 검증
    from pptx import Presentation
    p2 = Presentation(OUTPUT)
    sl = p2.slides[2]
    sp = sl.shapes

    def rot(s):
        try:
            return s.rotation
        except Exception:
            return "n/a"

    print(f"\n{'='*80}\nv8 slides[2] 화살표/라벨 최종 좌표 검증\n{'='*80}")
    label = {10: "세로 신탁사↔차주", 12: "가로 시공사↔차주", 14: "가로 차주↔대주",
             16: "사선 신탁사→대주", 17: "라벨 담보신탁 우선수익권"}
    exp = {10: (5.25, 3.00, 0.00, 0.50), 12: (2.50, 4.00, 2.00, 0.00),
           14: (6.00, 4.00, 2.00, 0.00), 16: (6.00, 2.70, 2.00, 0.70),
           17: (6.30, 2.45, 1.95, 0.30)}
    print(f"{'idx':>3} {'설명':<16} {'L':>5} {'T':>5} {'W':>5} {'H':>5}  기대일치  rot")
    for i in [10, 12, 14, 16, 17]:
        l = Emu(sp[i].left).inches; t = Emu(sp[i].top).inches
        w = Emu(sp[i].width).inches; h = Emu(sp[i].height).inches
        e = exp[i]
        ok = all(abs(a-b) < 0.01 for a, b in zip((l, t, w, h), e))
        print(f"{i:>3} {label[i]:<16} {l:5.2f} {t:5.2f} {w:5.2f} {h:5.2f}  {'OK' if ok else 'FAIL':<6}  {rot(sp[i])}")

    print("\n[박스 표 idx4~7 불변 확인]")
    for i in range(4, 8):
        l = Emu(sp[i].left).inches; t = Emu(sp[i].top).inches
        head = sp[i].table.cell(0, 0).text.strip() if sp[i].has_table else ""
        print(f"  idx={i} {head:<6} L={l:.2f} T={t:.2f}")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
