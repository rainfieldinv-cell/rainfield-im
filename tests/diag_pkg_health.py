"""읽기 전용 — pptx 패키지 중복파트 진단 + 폴더 산출물 목록. 삭제/수정 없음.
잠금 회피: 각 pptx 를 임시 복사본으로 검사."""
import sys, io, os, glob, zipfile, shutil, datetime, tempfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

folder = r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im"
TMP = tempfile.gettempdir()


def safe_copy(f):
    dst = os.path.join(TMP, "_chk_" + os.path.basename(f))
    shutil.copyfile(f, dst)
    return dst


def dup_parts(zpath):
    zz = zipfile.ZipFile(zpath)
    return [n for n, c in Counter(zz.namelist()).items() if c > 1]


# ===== (A) v10 진단 =====
v10 = os.path.join(folder, "test_ai_builders_v10.pptx")
c10 = safe_copy(v10)
prs = Presentation(c10)
print("[v10] python-pptx 인식 슬라이드 수:", len(prs.slides))
for i, s in enumerate(prs.slides):
    first = ""
    for sp in s.shapes:
        if sp.has_text_frame and sp.text_frame.text.strip():
            first = sp.text_frame.text.strip()[:30]; break
    print(f"  slide[{i}] <{first}>")
print("[v10] 중복 등록 파트:", dup_parts(c10) or "없음")

# ===== 버전별 건강검진 =====
print("\n--- 버전별 중복 파트 검사 ---")
results = {}
for f in sorted(glob.glob(os.path.join(folder, "test_ai_builders_v*.pptx"))):
    try:
        c = safe_copy(f)
        d = dup_parts(c)
        pp = Presentation(c)
        results[os.path.basename(f)] = (len(pp.slides), d)
        print(f"  {os.path.basename(f)}: 슬라이드 {len(pp.slides)}장, "
              f"중복파트={'있음 ' + str(d) if d else '없음'}")
    except Exception as e:
        print(f"  {os.path.basename(f)}: 열기 실패 - {e}")

# ===== (B) 폴더 내 pptx 전체 목록 =====
print("\n--- 폴더 내 .pptx 전체 (크기/수정시각) ---")
for f in sorted(glob.glob(os.path.join(folder, "*.pptx"))):
    sz = os.path.getsize(f)
    mt = datetime.datetime.fromtimestamp(os.path.getmtime(f)).strftime("%Y-%m-%d %H:%M")
    print(f"  {os.path.basename(f):<30} {sz:>12,}B  {mt}")
print("\n--- templates 폴더 ---")
for f in sorted(glob.glob(os.path.join(folder, "templates", "*.pptx"))):
    print(f"  templates/{os.path.basename(f)}  (디자인 원본 — 보존)")
