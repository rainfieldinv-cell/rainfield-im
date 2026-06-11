import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v20_readonly.pptx"
prs = Presentation(PATH)
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def q(t):
    return "{%s}%s" % (A, t)


def fill_hex(cell):
    tcPr = cell._tc.find(q("tcPr"))
    if tcPr is not None:
        sf = tcPr.find(q("solidFill"))
        if sf is not None:
            sc = sf.find(q("srgbClr"))
            if sc is not None:
                return sc.get("val")
    return "-"


sl = prs.slides[1]
for sh in sl.shapes:
    if not sh.has_table:
        continue
    tbl = sh.table
    print(f"표 {len(tbl.rows)}r×{len(tbl.columns)}c")
    for ri in range(len(tbl.rows)):
        for ci in range(len(tbl.columns)):
            cell = tbl.cell(ri, ci)
            p0 = cell.text_frame.paragraphs[0]
            algn = p0.alignment
            run = next((r for r in p0.runs if r.text.strip()), None)
            fn = run.font.name if run else None
            sz = run.font.size.pt if run and run.font.size else None
            bd = run.font.bold if run else None
            ml = Emu(cell.margin_left).inches
            mr = Emu(cell.margin_right).inches
            print(f"  [{ri},{ci}] align={algn} fill={fill_hex(cell)} font={fn} {sz}pt bold={bd} "
                  f"marL={ml:.2f} marR={mr:.2f} '{cell.text.strip()[:16]}'")
