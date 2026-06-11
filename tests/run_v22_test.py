"""run_v22_test.py — 표 서식 표준(style_table) 적용. 슬1 표 + 슬7 박스 불변 확인."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v22.pptx")
BIZ = "천안 부성2지구 도시개발사업"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def fill_hex(cell):
    tcPr = cell._tc.find("{%s}tcPr" % A)
    if tcPr is not None:
        sf = tcPr.find("{%s}solidFill" % A)
        if sf is not None:
            sc = sf.find("{%s}srgbClr" % A)
            if sc is not None:
                return sc.get("val")
    return "-"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ, page_num=1)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ, page_num=2)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ, page_num=3)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}")
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v22] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)

    # 슬1 표
    print("\n[슬라이드1 사모사채 표 — 표준 적용 확인]")
    sl1 = p2.slides[1]
    tbl = next((sh.table for sh in sl1.shapes if sh.has_table), None)
    for ri in range(len(tbl.rows)):
        for ci in range(len(tbl.columns)):
            cell = tbl.cell(ri, ci)
            p0 = cell.text_frame.paragraphs[0]
            run = next((r for r in p0.runs if r.text.strip()), None)
            fn = run.font.name if run else None
            sz = run.font.size.pt if run and run.font.size else None
            print(f"  [{ri},{ci}] align={str(p0.alignment)[:6]} font={fn} {sz}pt "
                  f"marL={Emu(cell.margin_left).inches:.2f} marR={Emu(cell.margin_right).inches:.2f} "
                  f"fill={fill_hex(cell)} '{cell.text.strip()[:14]}'")
        if ri >= 2:
            print("   ...")
            break

    # 슬7 박스 불변
    print("\n[슬라이드7 구조도 박스 — 12~13pt·4색 헤더 불변 확인]")
    sl7 = p2.slides[2]
    for sh in sl7.shapes:
        if sh.has_table:
            t = sh.table
            head = t.cell(0, 0)
            sizes = sorted({r.font.size.pt for ri in range(len(t.rows)) for r in t.cell(ri, 0).text_frame.paragraphs[0].runs if r.font.size})
            print(f"  박스 '{head.text.strip()}' 헤더색={fill_hex(head)} font={sizes}pt")


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
