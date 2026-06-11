"""v10 저장 — 입력(v9 복사본)에 idx16 조정 적용 후 test_ai_builders_v10.pptx 저장."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

ROOT = r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im"
SRC = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v9_readonly.pptx"
OUT = os.path.join(ROOT, "test_ai_builders_v10.pptx")

prs = Presentation(SRC)
slide = prs.slides[2]
sp = slide.shapes[16]

sp.left, sp.top, sp.width, sp.height = Inches(6.00), Inches(2.50), Inches(2.00), Inches(1.40)
prstGeom = sp._element.find('.//' + qn('a:prstGeom'))
for av in prstGeom.findall(qn('a:avLst')):
    prstGeom.remove(av)
avLst = etree.SubElement(prstGeom, qn('a:avLst'))
gd = etree.SubElement(avLst, qn('a:gd'))
gd.set('name', 'adj1'); gd.set('fmla', 'val 85000')

prs.save(OUT)
print(f"저장: {OUT}")

# 재오픈 검증
p2 = Presentation(OUT)
s = p2.slides[2].shapes[16]
g = s._element.find('.//' + qn('a:prstGeom'))
adj = g.find(qn('a:avLst') + '/' + qn('a:gd'))
print(f"idx16: prst={g.get('prst')} adj1={adj.get('fmla') if adj is not None else '없음'} "
      f"L/T/W/H={Emu(s.left).inches:.2f}/{Emu(s.top).inches:.2f}/"
      f"{Emu(s.width).inches:.2f}/{Emu(s.height).inches:.2f}")
