"""run_v16_test.py — 빌더 재생성으로 v16 + slides[0] 섹션 정리/꺾쇠 재배치 검증."""
import sys, os, io, traceback, zipfile
from collections import Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
os.chdir(ROOT)

PDF_PATH = os.path.join(ROOT, "★ TS_천안 부성2지구 도시개발사업 BL_v5.2.pdf")
OUTPUT   = os.path.join(ROOT, "test_ai_builders_v16.pptx")
BIZ_NAME = "천안 부성2지구 도시개발사업"


def main():
    import fitz
    doc = fitz.open(PDF_PATH)
    pdf_text = "\n".join(p.get_text() for p in doc)
    doc.close()

    from modules import page_builders
    from modules.ai_slide_builders import (
        generate_executive_summary, generate_sasae_overview, generate_investment_structure,
        build_slide_2_executive_summary, build_slide_5_sasae_overview,
        build_slide_7_investment_structure,
    )
    from modules.page_builders import create_presentation_from_template, finalize_presentation
    from modules.ppt_generator import save_presentation
    from pptx import Presentation
    from pptx.util import Emu

    ov = os.environ.get("LAYOUT_OVERRIDE", "").strip()
    if ov:
        page_builders.LAYOUT_PPTX_PATH = ov

    r2 = generate_executive_summary(pdf_text)
    r5 = generate_sasae_overview(pdf_text)
    r7 = generate_investment_structure(pdf_text)
    if not all(r.get("ok") for r in [r2, r5, r7]):
        print("[FAIL] API"); sys.exit(1)

    prs = create_presentation_from_template()
    n = len(prs.slides)
    build_slide_2_executive_summary(prs, r2["data"], business_name=BIZ_NAME, page_num=2)
    build_slide_5_sasae_overview(prs, r5["data"], business_name=BIZ_NAME, page_num=6)
    build_slide_7_investment_structure(prs, r7["data"], business_name=BIZ_NAME, page_num=8)
    finalize_presentation(prs, n)
    save_presentation(prs, OUTPUT)
    print(f"\n저장: {OUTPUT}")
    dup = [nm for nm, c in Counter(zipfile.ZipFile(OUTPUT).namelist()).items() if c > 1]
    print(f"[v16] 중복파트: {'있음 ' + str(dup) if dup else '없음 (정상)'}")

    p2 = Presentation(OUTPUT)
    sl = p2.slides[0]
    names = [sp.name for sp in sl.shapes]

    print(f"\n[slides[0] 전체 {len(sl.shapes)} shape]")
    for i, sp in enumerate(sl.shapes):
        tp = str(sp.shape_type).split("(")[0].split(".")[-1].strip()
        l = Emu(sp.left).inches if sp.left is not None else 0
        t = Emu(sp.top).inches if sp.top is not None else 0
        w = Emu(sp.width).inches if sp.width is not None else 0
        h = Emu(sp.height).inches if sp.height is not None else 0
        txt = sp.text_frame.text.replace("\n", " / ")[:60] if sp.has_text_frame else ""
        print(f"  idx={i:2d} {tp:<11} {sp.name:<14} L={l:5.2f} T={t:5.2f} W={w:5.2f} H={h:5.2f} <{txt}>")

    print("\n[부제 삭제 확인]")
    for nm in ["TextBox 56", "TextBox 54"]:
        print(f"  '{nm}': {'남음 FAIL' if nm in names else '삭제됨 OK'}")
    print("[제목 GROUP 유지]")
    for nm in ["그룹 1", "그룹 2", "그룹 51"]:
        sh = next((s for s in sl.shapes if s.name == nm), None)
        t = ""
        if sh is not None:
            t = " ".join(c.text_frame.text for c in sh.shapes if c.has_text_frame).strip()
        print(f"  '{nm}': {'존재' if sh else '없음 FAIL'}  제목='{t[:20]}'")

    # 꺾쇠 가로폭 통일 + 세로 높이/간격
    def g(nm):
        s = next((x for x in sl.shapes if x.name == nm), None)
        return (Emu(s.left).inches, Emu(s.top).inches, Emu(s.width).inches, Emu(s.height).inches) if s else None
    print("\n[꺾쇠 바깥 육각형 L/T/W/H]")
    outers = [("S1", "육각형 5"), ("S2", "육각형 24"), ("S3", "육각형 48")]
    prev_bot = None
    for tag, nm in outers:
        v = g(nm)
        if v:
            print(f"  {tag} {nm}: L={v[0]:.2f} T={v[1]:.2f} W={v[2]:.2f} H={v[3]:.2f} (하단 {v[1]+v[3]:.2f})")
    print("\n[섹션 세로 간격 (앞 꺾쇠 하단 → 다음 꺾쇠 상단)]")
    o1, o2, o3 = g("육각형 5"), g("육각형 24"), g("육각형 48")
    if o1 and o2 and o3:
        print(f"  S1하단 {o1[1]+o1[3]:.2f} → S2상단 {o2[1]:.2f} : gap {o2[1]-(o1[1]+o1[3]):.2f}")
        print(f"  S2하단 {o2[1]+o2[3]:.2f} → S3상단 {o3[1]:.2f} : gap {o3[1]-(o2[1]+o2[3]):.2f}")
    print("\n[바깥 육각형 가로폭 통일?]", {round(g(n)[2],2) for _, n in outers})


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc(); sys.exit(1)
