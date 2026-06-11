"""v12 dry-run — idx16 bbox(H1.05) + adj1=50000. 저장 안 함."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v11_readonly.pptx"
prs = Presentation(PATH)
sp = prs.slides[2].shapes[16]

print("변경 전 bbox:", f"{Emu(sp.left).inches:.2f}/{Emu(sp.top).inches:.2f}/"
      f"{Emu(sp.width).inches:.2f}/{Emu(sp.height).inches:.2f}")

sp.left, sp.top, sp.width, sp.height = Inches(6.00), Inches(2.50), Inches(2.00), Inches(1.05)
prstGeom = sp._element.find('.//' + qn('a:prstGeom'))
for av in prstGeom.findall(qn('a:avLst')):
    prstGeom.remove(av)
avLst = etree.SubElement(prstGeom, qn('a:avLst'))
gd = etree.SubElement(avLst, qn('a:gd'))
gd.set('name', 'adj1'); gd.set('fmla', 'val 50000')

L, T, W, H = (Emu(sp.left).inches, Emu(sp.top).inches, Emu(sp.width).inches, Emu(sp.height).inches)
print("변경 후 bbox:", f"L={L:.2f} T={T:.2f} W={W:.2f} H={H:.2f}")
print(f"끝점 예상: X={L+W:.2f}(대주 왼쪽변 8.00), Y={T+H:.2f}(대출약정 Y4.00과 {4.00-(T+H):.2f} 분리)")
print("\n--- prstGeom XML ---")
print(etree.tostring(prstGeom, pretty_print=True).decode())

# 선 스타일 불변 확인
ln = sp._element.find('.//' + qn('a:ln'))
clr = sp._element.find('.//' + qn('a:solidFill') + '/' + qn('a:srgbClr'))
tail = sp._element.find('.//' + qn('a:tailEnd'))
print("ln(w,색,tail):", ln.get('w'), clr.get('val') if clr is not None else None,
      tail.get('type') if tail is not None else None, " prst=", prstGeom.get('prst'))
print("[저장 안 함 — dry-run]")
