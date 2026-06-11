"""읽기 전용 — 제안서 대표 슬라이드의 요소 배치(좌표) 분석. 세로→가로 구성 학습용."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

PPT = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_netmarble_prop.pptx"
prs = Presentation(PPT)
TARGETS = [int(x) for x in (sys.argv[2].split(",") if len(sys.argv) > 2 else ["9", "14", "17", "26"])]

for si in TARGETS:
    sl = prs.slides[si]
    print(f"\n{'='*92}\n[슬라이드 {si}] 요소 {len(sl.shapes)}개")
    print(f"{'종류':<10} {'L':>5} {'T':>5} {'W':>5} {'H':>5}  내용/크기")
    for sh in sorted(sl.shapes, key=lambda s: (Emu(s.top or 0).inches, Emu(s.left or 0).inches)):
        tp = str(sh.shape_type).split("(")[0].split(".")[-1].strip()
        L = Emu(sh.left or 0).inches; T = Emu(sh.top or 0).inches
        W = Emu(sh.width or 0).inches; H = Emu(sh.height or 0).inches
        info = ""
        if sh.has_table:
            t = sh.table
            info = f"표 {len(t.rows)}r×{len(t.columns)}c : " + " | ".join(
                t.cell(0, c).text.strip()[:8] for c in range(min(len(t.columns), 6)))
        elif sh.shape_type == 13:
            info = "[그림]"
        elif sh.has_text_frame:
            info = sh.text_frame.text.replace("\n", " ").strip()[:42]
        print(f"{tp:<10} {L:5.2f} {T:5.2f} {W:5.2f} {H:5.2f}  {info}")
