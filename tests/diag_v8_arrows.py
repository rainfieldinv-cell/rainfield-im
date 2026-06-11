"""읽기 전용 — v8 slides[2] 화살표/커넥터 타입·XML 진단. 수정/저장 없음."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Temp\_v8_readonly.pptx"
prs = Presentation(PATH)
slide = prs.slides[2]

targets = [10, 12, 14, 16]
for i in targets:
    sp = slide.shapes[i]
    print(f"\n===== idx={i} =====")
    print("shape_type:", sp.shape_type)
    print("name:", sp.name)
    print("L/T/W/H:", f"{Emu(sp.left).inches:.2f}/{Emu(sp.top).inches:.2f}/"
                      f"{Emu(sp.width).inches:.2f}/{Emu(sp.height).inches:.2f}")
    # 커넥터 연결점(stCxn/endCxn) 탐색
    el = sp._element
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    st = el.find(f".//{{{P}}}stCxn") or el.find(f".//{{{A}}}stCxn")
    en = el.find(f".//{{{P}}}endCxn") or el.find(f".//{{{A}}}endCxn")
    print("stCxn:", etree.tostring(st).decode().strip() if st is not None else "없음(좌표 고정)")
    print("endCxn:", etree.tostring(en).decode().strip() if en is not None else "없음(좌표 고정)")
    print("tag:", etree.QName(el).localname)
    limit = 2600 if i == 16 else 900
    print("--- XML ---")
    print(etree.tostring(el, pretty_print=True).decode()[:limit])
