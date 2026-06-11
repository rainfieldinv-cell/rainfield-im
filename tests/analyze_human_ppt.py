# -*- coding: utf-8 -*-
"""사람이 만든 제안서 PPT 3개의 표/제목 구조를 덤프해 원본 변환 방식 분석."""
import sys, io, os, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

ROOT = r"C:\Users\jbzle\OneDrive\Desktop\종합\자동화\rainfield-im"
TARGETS = {
    "천안": "[Rainfield]천안 부성2지구_토지담보출_사모사채_제안서_260310.pptx",
    "대전": "[Rainfield] 대전서남터미널_토지담보대출_260511.pptx",
    "넷마블": "[Rainfield]로 넷마블 G타워 1종 수익증권 제안서_260526.pptx",
}

def cell_txt(c):
    return " ".join(p.text for p in c.text_frame.paragraphs).strip()

def shape_txt(sh):
    if not sh.has_text_frame:
        return ""
    return " / ".join(p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip())

key = sys.argv[1] if len(sys.argv) > 1 else "대전"
path = os.path.join(ROOT, TARGETS[key])
prs = Presentation(path)
print(f"=== {key} : {os.path.basename(path)} ({len(prs.slides)} slides) ===\n")

for i, slide in enumerate(prs.slides, 1):
    # 제목 후보: 가장 위쪽 텍스트박스들
    texts = []
    tables = []
    for sh in slide.shapes:
        if sh.has_table:
            tables.append(sh)
        elif sh.has_text_frame and shape_txt(sh):
            top = sh.top if sh.top is not None else 0
            texts.append((Emu(top).inches, shape_txt(sh)))
    texts.sort()
    # 상단 텍스트 2개만(섹션라벨/제목)
    head = " | ".join(t for _, t in texts[:3])
    print(f"[슬라이드 {i}]  표{len(tables)}개")
    print(f"   상단텍스트: {head[:160]}")
    for ti, tb in enumerate(tables, 1):
        t = tb.table
        nr = len(t.rows); nc = len(t.columns)
        # 위치/크기
        L = Emu(tb.left).inches; T = Emu(tb.top).inches
        W = Emu(tb.width).inches; H = Emu(tb.height).inches
        print(f"   └ 표{ti}: {nr}행 x {nc}열  pos=({L:.2f},{T:.2f}) size=({W:.2f}x{H:.2f})")
        # 첫 행(헤더) + 처음 3개 데이터행
        for r in range(min(nr, 4)):
            row = [cell_txt(t.cell(r, c))[:18] for c in range(nc)]
            print(f"        r{r}: {row}")
    print()
