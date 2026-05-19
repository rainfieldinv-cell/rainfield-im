"""
page_builders.py
─────────────────────────────────────────────────────────
실제 PPT 슬라이드를 만드는 빌더 함수 모음.
ppt_generator.py의 기본 함수를 조합해서 슬라이드를 구성합니다.

현재 구현:
  - clone_slide_layout()         : 레이아웃.pptx 슬라이드 복제 (공통 헬퍼)
  - replace_placeholder_text()   : 복제된 슬라이드의 텍스트 교체
  - replace_placeholder_image()  : 복제된 슬라이드의 이미지 교체
  - auto_resize_text_to_fit()    : 텍스트 길이에 따라 폰트 자동 축소
  - build_cover_slide()          : 표지 슬라이드 생성
  - build_test_cover()           : 빠른 테스트용 표지 생성
─────────────────────────────────────────────────────────
"""

import io
import os
import re
import copy
from datetime import datetime

from lxml import etree
from pptx import Presentation as _LayoutPrs
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn as _qn
from pptx.util import Cm as _Cm, Pt as _Pt
from PIL import Image

from modules.constants import (
    FONT_BOLD,
    COLOR_DARK, COLOR_GRAY, COLOR_BG_LIGHT,
    COVER_IM_LABEL_COLOR, COVER_DATE_COLOR,
    OUTPUT_DIR, OUTPUT_FILENAME_FORMAT, TEMPLATES_DIR,
    LAYOUT_PPTX_PATH, SLIDE_INDEX_MAP,
    DEFAULT_TOC_MAP,
)

# ─────────────────────────────────────────────
# [고정 이미지 경로 — templates/ 폴더]
# 여기를 수정하면 고정 이미지 파일명이 바뀝니다
# ─────────────────────────────────────────────
IMG_DIVIDER        = os.path.join(TEMPLATES_DIR, "표지.png")               # 연두색 세로 구분선
IMG_DISCLAIMER     = os.path.join(TEMPLATES_DIR, "[Disclaimer].png")       # 면책 문구 박스
IMG_CONFIDENTIAL   = os.path.join(TEMPLATES_DIR, "Strictly Confidential.png")  # 상단 좌측 라벨

from modules.ppt_generator import (
    create_blank_presentation,
    add_text_box,
    add_image,
    add_oval_with_image,
    make_circular_image_png,
    save_presentation,
)

# ─────────────────────────────────────────────
# [rId 관계 속성 네임스페이스 상수]
# ─────────────────────────────────────────────
_OOXML_R      = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
_R_EMBED      = f'{{{_OOXML_R}}}embed'
_R_LINK       = f'{{{_OOXML_R}}}link'
_R_ID         = f'{{{_OOXML_R}}}id'
_HEADER_TAGS  = {_qn('p:nvGrpSpPr'), _qn('p:grpSpPr')}   # spTree 필수 헤더
_COPY_REL_TYPES = {
    f'{_OOXML_R}/image',
    f'{_OOXML_R}/audio',
    f'{_OOXML_R}/video',
    f'{_OOXML_R}/oleObject',
}


# ─────────────────────────────────────────────
# [템플릿 기반 프레젠테이션 생성]
# 레이아웃.pptx를 기반으로 prs를 만듭니다.
# 기존 create_blank_presentation() 대신 이 함수를 사용합니다.
# 레이아웃·마스터·폰트 등이 레이아웃.pptx와 완전히 동일하게 유지됩니다.
# ─────────────────────────────────────────────
def create_presentation_from_template():
    """
    레이아웃.pptx를 그대로 복사해서 새 프레젠테이션 객체를 반환합니다.
    템플릿 슬라이드(0~N)는 clone_slide_layout() 복제 원본으로 사용됩니다.
    모든 빌더 호출이 끝나면 finalize_presentation()으로 원본 슬라이드를 삭제합니다.
    """
    with open(LAYOUT_PPTX_PATH, 'rb') as f:
        prs = _LayoutPrs(io.BytesIO(f.read()))
    return prs


# ─────────────────────────────────────────────
# [템플릿 슬라이드 삭제 — 빌드 완료 후 호출]
# ─────────────────────────────────────────────
def finalize_presentation(prs, template_count: int):
    """
    빌드가 끝난 후 앞쪽의 원본 템플릿 슬라이드를 모두 삭제합니다.
    template_count : create_presentation_from_template() 직후 len(prs.slides) 값

    사용법:
        prs = create_presentation_from_template()
        n   = len(prs.slides)   # 템플릿 슬라이드 수 기록
        build_cover_slide(prs, ...)
        build_toc_slide(prs)
        ...
        finalize_presentation(prs, n)
        save_presentation(prs, output_path)
    """
    # 뒤에서부터 제거해야 인덱스가 밀리지 않음
    for i in range(template_count - 1, -1, -1):
        _delete_slide(prs, i)


def _delete_slide(prs, idx: int):
    """프레젠테이션에서 idx 번째 슬라이드를 삭제합니다 (OOXML 수준 제거)."""
    sldIdLst = prs.slides._sldIdLst
    if idx >= len(sldIdLst):
        return
    rId_attr = f'{{{_OOXML_R}}}id'
    sld_elem = sldIdLst[idx]
    rId = sld_elem.get(rId_attr)
    sldIdLst.remove(sld_elem)
    # Part 관계 제거
    if rId:
        try:
            del prs.part._rels[rId]
        except Exception:
            pass


# ─────────────────────────────────────────────
# (공통 헬퍼 a) 슬라이드 복제
# ─────────────────────────────────────────────
_GRAPHIC_FRAME_TAG = _qn('p:graphicFrame')   # 테이블·차트 컨테이너 태그


def clone_slide_layout(prs, slide_type: str, skip_graphic_frames: bool = False):
    """
    prs(레이아웃.pptx 기반) 안의 템플릿 슬라이드를 복제합니다.
    같은 prs 안에서 복제하므로 슬라이드 레이아웃 참조가 그대로 유지됩니다.
    → 로고·삼각형·폰트·위치가 원본과 완전히 동일하게 나옵니다.

    Parameters
    ----------
    prs                  : create_presentation_from_template()으로 만든 Presentation 객체
    slide_type           : SLIDE_INDEX_MAP의 키 (예: "content", "toc_4" 등)
    skip_graphic_frames  : True이면 <p:graphicFrame> (테이블·차트)을 복사하지 않음
                           content/content_with_image 슬라이드에서 템플릿 테이블이
                           그대로 복사되는 Bug #2/#7 방지용

    Returns
    -------
    추가된 슬라이드 객체
    """
    idx = SLIDE_INDEX_MAP.get(slide_type)
    if idx is None:
        raise ValueError(f"알 수 없는 슬라이드 타입: {slide_type!r} — constants.py의 SLIDE_INDEX_MAP 확인")

    source = prs.slides[idx]   # 같은 prs 안의 템플릿 슬라이드

    # 소스와 동일한 레이아웃으로 새 슬라이드 추가
    # → 레이아웃이 같으므로 로고·삼각형 등이 자동으로 포함됨
    new_slide = prs.slides.add_slide(source.slide_layout)

    # 새 슬라이드의 기본 도형 제거 (헤더는 유지)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        if child.tag not in _HEADER_TAGS:
            sp_tree.remove(child)

    # 소스 슬라이드 이미지 관계를 새 슬라이드에 등록
    rId_map = {}
    for rel in source.part.rels.values():
        try:
            if not rel.is_external and rel.reltype in _COPY_REL_TYPES:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rel.rId] = new_rId
        except Exception:
            pass

    # 소스 슬라이드 도형을 새 슬라이드에 복사
    for child in source.shapes._spTree:
        if child.tag in _HEADER_TAGS:
            continue
        # 테이블·차트(graphicFrame)는 skip_graphic_frames=True 일 때 건너뜀
        if skip_graphic_frames and child.tag == _GRAPHIC_FRAME_TAG:
            continue
        node = copy.deepcopy(child)
        # 이미지 rId 재매핑
        for elem in node.iter():
            for attr_key in (_R_EMBED, _R_LINK, _R_ID):
                if attr_key in elem.attrib and elem.attrib[attr_key] in rId_map:
                    elem.attrib[attr_key] = rId_map[elem.attrib[attr_key]]
        sp_tree.append(node)

    return new_slide


# ─────────────────────────────────────────────
# (공통 헬퍼 b-내부) 텍스트 프레임 내용 교체 (스타일 보존)
# ─────────────────────────────────────────────
def _replace_text_frame_content(tf, new_text: str):
    """
    텍스트 프레임의 내용을 new_text로 교체합니다.
    첫 번째 런(run)의 폰트·크기·색상 스타일을 모든 줄에 적용합니다.
    \\n으로 줄바꿈 처리됩니다.
    """
    # 첫 번째 런의 rPr(폰트 스타일) 및 단락의 pPr(정렬 등) 저장
    first_rPr = None
    first_pPr = None
    if tf.paragraphs:
        first_p = tf.paragraphs[0]
        pPr_elem = first_p._p.find(_qn('a:pPr'))
        if pPr_elem is not None:
            first_pPr = copy.deepcopy(pPr_elem)
        for run in first_p.runs:
            rPr_elem = run._r.find(_qn('a:rPr'))
            if rPr_elem is not None:
                first_rPr = copy.deepcopy(rPr_elem)
            break

    # 기존 단락 모두 제거
    txBody = tf._txBody
    for old_p in list(txBody.findall(_qn('a:p'))):
        txBody.remove(old_p)

    # 새 텍스트로 단락 재생성 (\n 기준 줄 분리)
    for line in (new_text or '').split('\n'):
        new_p = etree.SubElement(txBody, _qn('a:p'))
        if first_pPr is not None:
            new_p.insert(0, copy.deepcopy(first_pPr))
        new_r = etree.SubElement(new_p, _qn('a:r'))
        if first_rPr is not None:
            new_r.insert(0, copy.deepcopy(first_rPr))
        new_t = etree.SubElement(new_r, _qn('a:t'))
        new_t.text = line


# ─────────────────────────────────────────────
# (공통 헬퍼 b-2) 하단 사업명 교체
# ─────────────────────────────────────────────
def _replace_footer_business_name(slide, business_name: str):
    """
    슬라이드 하단 푸터의 '사업명' 자리표시자를 실제 사업명으로 교체합니다.
    left < 1.0cm, top > 17.5cm 조건으로 푸터 도형을 찾습니다.
    """
    if not business_name:
        return
    FOOTER_TOP_MIN = 17.5
    for shape in slide.shapes:
        if (shape.has_text_frame
                and shape.left / 360000 < 1.0
                and shape.top  / 360000 > FOOTER_TOP_MIN):
            _replace_text_frame_content(shape.text_frame, business_name)
            return


# ─────────────────────────────────────────────
# (공통 헬퍼 b) 텍스트 교체
# ─────────────────────────────────────────────
def replace_placeholder_text(slide, placeholder_hint: str, new_text: str,
                              layout_shapes: list = None):
    """
    슬라이드에서 placeholder_hint에 해당하는 텍스트 박스를 찾아 텍스트를 교체합니다.
    폰트·크기·색상은 원본 유지, 텍스트 내용만 바꿉니다.
    못 찾으면 경고만 출력하고 계속 실행합니다 (오류로 멈추지 않음).

    Parameters
    ----------
    slide            : 수정할 슬라이드 객체
    placeholder_hint : 찾을 힌트 문자열 (layout_data.json의 placeholder_hint 값)
    new_text         : 새로 넣을 텍스트 (\\n으로 줄바꿈 가능)
    layout_shapes    : layout_data.json 해당 슬라이드의 shapes 리스트 (위치 매칭용)

    Returns
    -------
    True (성공) / False (실패)
    """
    TOLERANCE = _Cm(0.5)   # 0.5cm 오차 허용 — 여기를 수정하면 허용 오차가 바뀝니다

    if not layout_shapes:
        print(f"[경고] replace_placeholder_text: layout_shapes 없음 — '{placeholder_hint}' 매칭 불가")
        return False

    # layout_shapes에서 힌트에 해당하는 좌표(left_cm, top_cm) 찾기
    target_left = target_top = None
    for sh in layout_shapes:
        if sh.get('placeholder_hint') == placeholder_hint:
            target_left = sh['left_cm']
            target_top  = sh['top_cm']
            break

    if target_left is None:
        print(f"[경고] layout_shapes에서 hint='{placeholder_hint}' 를 찾지 못했습니다.")
        return False

    # 슬라이드에서 해당 좌표의 텍스트 도형 탐색
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if (abs(shape.left - _Cm(target_left)) < TOLERANCE and
                abs(shape.top  - _Cm(target_top))  < TOLERANCE):
            _replace_text_frame_content(shape.text_frame, new_text)
            return True

    print(f"[경고] hint='{placeholder_hint}' 위치({target_left:.2f}, {target_top:.2f})cm 의 도형을 슬라이드에서 찾지 못했습니다.")
    return False


# ─────────────────────────────────────────────
# (공통 헬퍼 c) 이미지 교체
# ─────────────────────────────────────────────
def replace_placeholder_image(slide, left_cm: float, top_cm: float,
                               image_bytes: bytes,
                               is_circular: bool = False,
                               border_color: RGBColor = None,
                               border_width_pt: float = None):
    """
    슬라이드의 지정 위치(left_cm, top_cm)에 있는 이미지 도형을 새 이미지로 교체합니다.
    기존 이미지의 크기·위치는 그대로 유지됩니다.

    Parameters
    ----------
    slide            : 수정할 슬라이드
    left_cm, top_cm  : 교체할 이미지의 위치 (layout_data.json 좌표 사용)
    image_bytes      : 새 이미지 bytes
    is_circular      : True이면 원형 마스킹 적용 (add_oval_with_image 사용)
    border_color     : 테두리 색 (RGBColor, None이면 테두리 없음)
    border_width_pt  : 테두리 두께 (pt)

    Returns
    -------
    추가된 Picture 도형 또는 None (실패 시)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    TOLERANCE = _Cm(0.5)

    # 기존 이미지 도형 찾기 → 크기 파악 후 제거
    old_w_cm = old_h_cm = 3.0   # 도형 못 찾을 때 기본값
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if (abs(shape.left - _Cm(left_cm)) < TOLERANCE and
                abs(shape.top  - _Cm(top_cm))  < TOLERANCE):
            old_w_cm = shape.width  / 360000   # EMU → cm
            old_h_cm = shape.height / 360000
            shape.element.getparent().remove(shape.element)
            break
    else:
        print(f"[경고] replace_placeholder_image: ({left_cm:.2f}, {top_cm:.2f})cm 에 이미지 도형이 없습니다. 기본 크기로 추가합니다.")

    # 새 이미지 삽입
    if is_circular:
        return add_oval_with_image(
            slide, image_bytes,
            left_cm, top_cm,
            diameter_cm=old_w_cm,
            border_color=border_color,
            border_width_pt=border_width_pt or 2.0,
        )
    else:
        return add_image(
            slide, image_bytes,
            left_cm, top_cm,
            width_cm=old_w_cm, height_cm=old_h_cm,
            border_color=border_color,
            border_width_pt=border_width_pt,
            fit_mode="cover",
        )


# ─────────────────────────────────────────────
# (공통 헬퍼 d) 텍스트 자동 축소
# ─────────────────────────────────────────────
def auto_resize_text_to_fit(text_frame, max_size: float = 10.5, min_size: float = 8.0):
    """
    텍스트 길이에 따라 글자 크기를 max_size~min_size 범위에서 자동 조절합니다.
    텍스트가 길수록 작아집니다.

    Parameters
    ----------
    text_frame : 크기를 조절할 텍스트 프레임 객체
    max_size   : 최대 글자 크기 pt (짧은 텍스트) — 여기를 수정하면 최대 크기가 바뀝니다
    min_size   : 최소 글자 크기 pt (긴 텍스트)  — 여기를 수정하면 최소 크기가 바뀝니다
    """
    text = text_frame.text
    char_count = len(text)

    # 문자 수 기반 선형 보간
    # 여기를 수정하면 자동 축소 기준 문자 수가 바뀝니다
    THRESHOLD_MAX = 100   # 이 수 이하면 max_size 유지
    THRESHOLD_MIN = 500   # 이 수 이상이면 min_size로 축소

    if char_count <= THRESHOLD_MAX:
        target = max_size
    elif char_count >= THRESHOLD_MIN:
        target = min_size
    else:
        ratio  = (char_count - THRESHOLD_MAX) / (THRESHOLD_MIN - THRESHOLD_MAX)
        target = max_size - ratio * (max_size - min_size)

    target = max(min_size, min(max_size, target))

    # 모든 런에 새 크기 적용
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.size = _Pt(target)
            except Exception:
                pass


# ─────────────────────────────────────────────
# [내부 헬퍼: 그룹 도형 제목 텍스트 교체]
# ─────────────────────────────────────────────
def _replace_group_title(group_shape, new_title: str):
    """
    그룹 도형(✓ 아이콘 + 제목 텍스트) 안의 텍스트를 new_title로 교체합니다.
    텍스트가 가장 긴 단락 = 제목 단락으로 판단합니다.
    기존 런의 a:t 텍스트를 직접 교체하는 방식을 사용합니다.
    """
    from pptx.oxml.ns import qn as _q

    # 그룹 내 모든 단락 중 합산 텍스트 길이가 가장 긴 것 = 제목 단락
    best_para = None
    best_len  = 0
    for para in group_shape.element.iter(_q('a:p')):
        t_elems = list(para.iter(_q('a:t')))
        total = ''.join((t.text or '') for t in t_elems).strip()
        if len(total) > best_len:
            best_len  = len(total)
            best_para = para

    if best_para is None:
        return

    # 첫 번째 a:t에 새 제목 텍스트 덮어쓰기, 나머지 a:t는 비움
    # (런 구조를 건드리지 않고 텍스트만 바꿔 폰트·크기·색상 보존)
    t_elems = list(best_para.iter(_q('a:t')))
    for idx, t_elem in enumerate(t_elems):
        t_elem.text = new_title if idx == 0 else ''


# ─────────────────────────────────────────────
# (b) Executive Summary 슬라이드 생성 함수
# ─────────────────────────────────────────────
def build_executive_summary_slide(prs, sections: list, business_name: str = ""):
    """
    Executive Summary 슬라이드를 만들어 prs에 추가합니다.
    표지 바로 다음 페이지로, 제안서 전체 내용을 3개 섹션으로 요약합니다.

    레이아웃 구조 (레이아웃.pptx 슬라이드[1] 기준):
      - 그룹 도형 3개 (✓ 아이콘 + 섹션 제목, 중앙 정렬)
      - 본문 텍스트박스 3개 (좌측 정렬, 너비 17cm 이상)
      - 서브타이틀 텍스트박스 3개 ("내용" 자리, 너비 ~13cm, 중앙 정렬)
      - 좌하단 사업명 텍스트박스 (페이지 하단 left=0)

    Parameters
    ----------
    prs           : create_blank_presentation()으로 만든 Presentation 객체
    sections      : 최대 3개 섹션 리스트. 각 항목:
                    {"title": "섹션 제목", "subtitle": "선택 소제목", "content": "본문 텍스트"}
                    - title    : 그룹 아이콘 옆 제목
                    - subtitle : 아이콘 아래 작은 설명 (없으면 빈 문자열)
                    - content  : 좌측 정렬 본문 (\\n 줄바꿈)
    business_name : 좌하단에 표시할 사업명

    Returns
    -------
    추가된 슬라이드 객체
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # ── 0. 레이아웃 슬라이드 복제 ─────────────────────
    slide = clone_slide_layout(prs, "executive_summary")

    # ── 1. 그룹 도형 수집 (top 오름차순 → 섹션 1/2/3 순)
    # 그룹 = ✓ 아이콘 + 섹션 제목 텍스트 포함 복합 도형
    groups = sorted(
        [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.GROUP],
        key=lambda s: s.top
    )

    # ── 2. 본문 콘텐츠 박스 수집 (너비 16cm 초과 + left 4cm 이상 → 실제 콘텐츠 박스)
    # 레이아웃 기준: 17.61/18.51/19.35cm 너비, left 4.49/4.91/4.50cm
    # 배경·테두리 역할의 겹친 박스들은 left 3.3~3.8cm 에 있어서 필터로 제외됨
    # 여기를 수정하면 본문 박스 선별 기준이 바뀝니다
    CONTENT_WIDTH_MIN = 16.0   # 너비 기준 (cm)
    CONTENT_LEFT_MIN  = 4.0    # left 기준 (cm) — 배경 박스와 구분
    content_boxes = sorted(
        [sh for sh in slide.shapes
         if sh.has_text_frame
         and sh.shape_type != MSO_SHAPE_TYPE.GROUP
         and sh.width / 360000 > CONTENT_WIDTH_MIN
         and sh.left  / 360000 > CONTENT_LEFT_MIN
         and sh.top   / 360000 > 1.0
         and sh.top   / 360000 < 18.0],
        key=lambda s: s.top
    )

    # ── 3. 서브타이틀 박스 수집 (너비 12~14cm → "내용" 자리)
    # 레이아웃 기준: 13.00cm 너비 / top: 3.14 / 7.33 / 15.15 cm
    # 여기를 수정하면 서브타이틀 박스 선별 너비 기준이 바뀝니다
    subtitle_boxes = sorted(
        [sh for sh in slide.shapes
         if sh.has_text_frame
         and sh.shape_type != MSO_SHAPE_TYPE.GROUP
         and 12.0 < sh.width / 360000 < 14.0
         and sh.top / 360000 > 1.0
         and sh.top / 360000 < 18.0],
        key=lambda s: s.top
    )

    # ── 4. 각 섹션 텍스트 채우기 ────────────────────
    for i, sec in enumerate(sections[:3]):
        title    = sec.get("title", "")
        subtitle = sec.get("subtitle", "")
        content  = sec.get("content", "")

        if i < len(groups):
            _replace_group_title(groups[i], title)
        if i < len(subtitle_boxes):
            _replace_text_frame_content(subtitle_boxes[i].text_frame, subtitle)
        if i < len(content_boxes):
            _replace_text_frame_content(content_boxes[i].text_frame, content)
            auto_resize_text_to_fit(content_boxes[i].text_frame, max_size=9.0, min_size=8.0)

    if len(sections) > 3:
        print(f"[경고] Executive Summary 섹션이 3개 초과 ({len(sections)}개). 3개까지만 표시됩니다.")

    # ── 5. 좌하단 사업명 텍스트박스 교체 ──────────────
    # 레이아웃 기준: left=0, top=18.35cm (슬라이드 하단)
    # 여기를 수정하면 좌하단 사업명 위치 기준이 바뀝니다
    if business_name:
        FOOTER_TOP_MIN = 17.5
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and shape.left / 360000 < 1.0          # 왼쪽 끝
                    and shape.top  / 360000 > FOOTER_TOP_MIN):  # 하단
                _replace_text_frame_content(shape.text_frame, business_name)
                break

    return slide


# ─────────────────────────────────────────────
# [Disclaimer 텍스트]
# 표지 하단 면책 문구 — 여기를 수정하면 면책 문구가 바뀝니다
# ─────────────────────────────────────────────
DISCLAIMER_TEXT = (
    "[Disclaimer]\n"
    "이 문서는 법적 구속력이 없으며 금융기관의 여신심사 과정, 금융약정 체결을 위한 검토과정 등 "
    "향후 사정에 따라 사전 통지 없이 변경될 수 있고 실제로 이루어질 거래의 내용이나 구조와 "
    "상이할 수 있습니다. 당사(들)는 본 금융조건에 대하여 최대한 정확하고 완전한 정보를 담고자 "
    "노력하였으나 이 문서에는 오류 또는 중요정보의 누락이 있을 수 있으며, 당사(들)는 정보의 "
    "정확성, 완전성 및 적정성에 대하여 어떠한 보증도 제공하지 않고 그에 관한 어떠한 의무도 "
    "부담하지 않습니다. 이 문서는 당사(들)의 주관을 전제로 한 자금조달을 위하여 작성되었으며, "
    "어떠한 경우에도 금융투자상품의 매수의 청약 권유 또는 매도의 청약으로 해석될 수 없고, "
    "귀사 이외의 제3자에게 제공, 공개 또는 유출되어서는 안됩니다."
)


# ─────────────────────────────────────────────
# [내부 헬퍼: 색상 채운 사각형 추가]
# ─────────────────────────────────────────────
def _add_filled_rect(slide, left_cm, top_cm, width_cm, height_cm, fill_color: RGBColor):
    """단색으로 채운 사각형 도형을 슬라이드에 추가합니다."""
    from pptx.util import Cm as PCm
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        PCm(left_cm), PCm(top_cm),
        PCm(width_cm), PCm(height_cm)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 테두리 제거
    return shape


# ─────────────────────────────────────────────
# (a) 표지 슬라이드 생성 함수
# ─────────────────────────────────────────────
def build_cover_slide(prs, business_name: str, year: str, month_en: str,
                      cover_image_bytes: bytes = None):
    """
    표지(Cover) 슬라이드를 만들어 prs에 추가합니다.

    layout_data.json의 cover 슬라이드 좌표를 기준으로 요소를 배치합니다.
    (templates/레이아웃.pptx 분석 결과 기반)

    Parameters
    ----------
    prs               : create_blank_presentation()으로 만든 Presentation 객체
    business_name     : 사업명 (예: "천안 부성2지구 도시개발사업")
    year              : 연도 문자열 (예: "2026")
    month_en          : 영문 월 (예: "March")
    cover_image_bytes : 표지 메인 이미지 bytes (None이면 회색 박스로 대체)

    Returns
    -------
    추가된 슬라이드 객체
    """
    # 표지 레이아웃으로 슬라이드 추가 (레이아웃.pptx의 cover 슬라이드와 동일한 layout 사용)
    cover_layout = prs.slides[SLIDE_INDEX_MAP["cover"]].slide_layout
    slide = prs.slides.add_slide(cover_layout)

    # 레이아웃에서 자동 복사된 placeholder 도형 모두 제거
    # (PowerPoint 기본 문구 "제목을 추가하려면 / 텍스트를 입력하십시오" 가 나오지 않도록)
    for ph in list(slide.placeholders):
        ph.element.getparent().remove(ph.element)

    # ── 0. 슬라이드 배경 흰색으로 설정 ──────────────────────────
    # 여기를 수정하면 슬라이드 배경색이 바뀝니다
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ── 1. 오른쪽 커버 이미지 ────────────────────────────────────
    # 좌표: (14.99, 0.0)cm / 크기: 12.53 × 19.05cm
    # 여기를 수정하면 표지 사진 위치·크기가 바뀝니다
    if cover_image_bytes:
        try:
            add_image(
                slide, cover_image_bytes,
                left_cm=14.99, top_cm=0.0,
                width_cm=12.53, height_cm=19.05,
                fit_mode="cover",
            )
        except Exception:
            _add_filled_rect(slide, 14.99, 0.0, 12.53, 19.05, COLOR_GRAY)
    else:
        _add_filled_rect(slide, 14.99, 0.0, 12.53, 19.05, COLOR_BG_LIGHT)

    # ── 2. 표지.png — 연두색 세로 구분선 (고정 이미지) ────────────
    # 좌표: (13.93, 0.0)cm / 크기: 2.86 × 18.91cm
    add_image(
        slide, IMG_DIVIDER,
        left_cm=13.93, top_cm=0.0,
        width_cm=2.86, height_cm=18.91,
    )

    # ── 3. Strictly Confidential.png (고정 이미지 — 상단 좌측) ───
    # 좌표: (0.15, 0.15)cm / 크기: 1.61 × 0.34cm
    # 여기를 수정하면 라벨 위치가 바뀝니다
    add_image(
        slide, IMG_CONFIDENTIAL,
        left_cm=0.15, top_cm=0.15,
        width_cm=1.61, height_cm=0.34,
    )

    # ── 4. "Information Memorandum" 텍스트 ───────────────────────
    # 좌표: (1.71, 7.31)cm / 크기: 10.43 × 0.73cm
    # 여기를 수정하면 IM 라벨 색상이 바뀝니다 (constants.py의 COVER_IM_LABEL_COLOR)
    add_text_box(
        slide,
        text="Information Memorandum",
        left_cm=1.71, top_cm=7.31,
        width_cm=10.43, height_cm=0.73,
        font_name=FONT_BOLD,
        font_size_pt=11.0,
        font_color=COVER_IM_LABEL_COLOR,
        align="left",
    )

    # ── 5. 사업명 제목 ────────────────────────────────────────────
    # 좌표: (1.71, 7.89)cm / 크기: 12.91 × 1.11cm
    # {business_name}을 수정하면 표지 제목이 바뀝니다
    _cover_title = business_name if business_name.strip().endswith("제안서") else f"{business_name} 제안서"
    add_text_box(
        slide,
        text=_cover_title,
        left_cm=1.71, top_cm=7.89,
        width_cm=12.91, height_cm=1.11,
        font_name=FONT_BOLD,
        font_size_pt=20.0,
        font_color=COLOR_DARK,
        align="left",
    )

    # ── 6. 날짜 텍스트 ───────────────────────────────────────────
    # 좌표: (1.71, 9.76)cm / 크기: 12.12 × 0.77cm
    # 여기를 수정하면 날짜 색상이 바뀝니다 (constants.py의 COVER_DATE_COLOR)
    add_text_box(
        slide,
        text=f"{month_en} {year}",
        left_cm=1.71, top_cm=9.76,
        width_cm=12.12, height_cm=0.77,
        font_name=FONT_BOLD,
        font_size_pt=12.0,
        font_color=COVER_DATE_COLOR,
        align="left",
    )

    # ── 7. [Disclaimer].png (고정 이미지 — 하단) ─────────────────
    # 좌표: (1.71, 15.04)cm / 크기: 12.95 × 3.30cm
    # 면책 문구 내용을 바꾸려면 templates/[Disclaimer].png 파일을 교체하세요
    add_image(
        slide, IMG_DISCLAIMER,
        left_cm=1.71, top_cm=15.04,
        width_cm=12.95, height_cm=3.30,
    )

    return slide


# ─────────────────────────────────────────────
# (b) 빠른 테스트 표지 생성 함수
# ─────────────────────────────────────────────
def build_test_cover():
    """
    더미 데이터로 표지 1장짜리 PPT를 생성합니다.
    결과 파일: output/test_cover.pptx

    터미널에서 실행:
      python -c "from modules.page_builders import build_test_cover; build_test_cover()"
    """
    print("[테스트] 표지 생성 시작...")

    # 빈 A4 프레젠테이션 생성
    prs = create_blank_presentation()

    # 테스트용 더미 이미지 (단색 배경 — 실제 조감도 대신)
    dummy_image = _make_dummy_image(width=1200, height=1600, color=(200, 210, 195))

    # 표지 슬라이드 추가
    build_cover_slide(
        prs,
        business_name="테스트 사업",    # ← 여기를 수정하면 테스트 사업명이 바뀝니다
        year="2026",
        month_en="May",
        cover_image_bytes=dummy_image,
    )

    # output/ 폴더에 저장
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, "test_cover.pptx")
    saved = save_presentation(prs, output_path)

    print(f"[완료] 파일 저장: {saved}")
    print("  → PowerPoint에서 열어 A4 크기 및 피플폰트 적용 여부를 확인하세요.")
    return saved


def _make_dummy_image(width: int, height: int, color: tuple) -> bytes:
    """테스트용 단색 이미지를 PNG bytes로 반환합니다."""
    img = Image.new("RGB", (width, height), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# (c) 목차(TOC) 슬라이드 생성 함수
# TOC oval 위치 — toc_4 / toc_5 레이아웃 공통
# ─────────────────────────────────────────────
_TOC_OVAL_LEFT = 0.8409
_TOC_OVAL_TOP  = 6.9520
_TOC_OVAL_W    = 5.2600
_TOC_OVAL_H    = 5.3000
_TOC_OVAL_TOL  = _Cm(0.4)


def build_toc_slide(prs, num_sections: int = 4, toc_map: dict = None,
                    toc_image_bytes_list: list = None):
    """
    목차(CONTENTS) 슬라이드를 만들어 prs에 추가합니다.
    '부제목' 자리표시자가 포함된 도형을 실제 소항목으로 교체합니다.
    레이아웃의 원형 oval은 항상 제거하고, 이미지가 제공된 경우 PIL 원형 사진을 삽입합니다.

    Parameters
    ----------
    prs                   : create_presentation_from_template()으로 만든 Presentation 객체
    num_sections          : 목차 항목 수 — 5개면 toc_5, 그 외 toc_4 레이아웃 사용
    toc_map               : {"01": ["1.1 ...", ...], "02": [...], ...}  None이면 DEFAULT_TOC_MAP 사용
    toc_image_bytes_list  : [bytes|None] — 목차 원형 슬롯(1개)에 넣을 이미지 bytes

    Returns
    -------
    추가된 슬라이드 객체
    """
    layout_type = "toc_5" if num_sections == 5 else "toc_4"
    slide = clone_slide_layout(prs, layout_type)

    # "부제목" 자리표시자가 포함된 도형 → 실제 소항목으로 교체
    _SEC_FROM_TXT = re.compile(r'^(\d+)\.')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "부제목" not in text:
            continue
        m = _SEC_FROM_TXT.match(text.strip())
        if not m:
            continue
        sec_num = f"0{m.group(1)}"  # "3" → "03", "4" → "04"
        subtitles = (toc_map or {}).get(sec_num) or DEFAULT_TOC_MAP.get(sec_num) or []
        _replace_text_frame_content(shape.text_frame, "\n".join(subtitles))

    # 섹션 타이틀 박스 교체 — toc_map에 "_labels" 키가 있을 때만 동작
    # 박스 식별: "0N  " 접두사(섹션 번호)로 식별, label 텍스트에 의존하지 않음
    section_labels = (toc_map or {}).get("_labels", {})
    if section_labels:
        _SEC_TITLE_PREFIX = re.compile(r'^(\d{2})\s{2}')
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            m = _SEC_TITLE_PREFIX.match(shape.text_frame.text.strip())
            if not m:
                continue
            sec_num = m.group(1)  # "01" ~ "05"
            new_label = section_labels.get(sec_num)
            if new_label:
                _replace_text_frame_content(shape.text_frame, f"{sec_num}  {new_label}")

    # 도형 삭제 없음 — 이미지 제공 시 oval 위에 z-order 상위로 덮어씌움
    img_bytes = (toc_image_bytes_list[0] if toc_image_bytes_list else None)
    if img_bytes:
        try:
            circ_png = make_circular_image_png(img_bytes, output_size=512,
                                               border_color_rgb=(255, 255, 255),
                                               border_width_px=8)
            add_image(slide, circ_png,
                      left_cm=_TOC_OVAL_LEFT, top_cm=_TOC_OVAL_TOP,
                      width_cm=_TOC_OVAL_W, height_cm=_TOC_OVAL_H)
        except Exception as exc:
            print(f"[경고] 목차 원형 이미지 삽입 실패: {exc}")

    return slide


# ─────────────────────────────────────────────
# (d-1) PPT 진짜 표 삽입 헬퍼
# ─────────────────────────────────────────────
def _add_table_to_slide(slide, table_data: list,
                         left_cm: float, top_cm: float,
                         width_cm: float, height_cm: float):
    """
    2D 문자열 배열을 PPT 편집 가능한 표 객체로 삽입합니다.
    헤더 행: Rainfield 초록 배경 + 흰 굵은 글씨
    데이터 행: 흰 배경 + 짙은 회색 글씨
    """
    if not table_data:
        return None

    rows = len(table_data)
    cols = max((len(row) for row in table_data), default=0)
    if rows == 0 or cols == 0:
        return None

    # Rainfield 초록 (#92D050)
    COLOR_HEADER_BG   = RGBColor(146, 208, 80)
    COLOR_HEADER_TEXT = 'FFFFFF'
    COLOR_BODY_TEXT   = '333333'

    try:
        gf  = slide.shapes.add_table(
            rows, cols,
            _Cm(left_cm), _Cm(top_cm),
            _Cm(width_cm), _Cm(height_cm),
        )
        tbl = gf.table

        # 열 너비 균등 분배
        col_w = int(_Cm(width_cm) / cols)
        for c_idx in range(cols):
            tbl.columns[c_idx].width = col_w

        for r, row_data in enumerate(table_data):
            is_header = (r == 0)
            for c in range(cols):
                cell_text = (row_data[c] if c < len(row_data) else "") or ""
                cell_text = str(cell_text).strip()

                cell = tbl.cell(r, c)
                tf   = cell.text_frame
                tf.word_wrap = True

                # 헤더 행 배경색
                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_HEADER_BG

                # txBody 초기화 후 텍스트 삽입
                txBody = tf._txBody
                for p_elem in list(txBody.findall(_qn('a:p'))):
                    txBody.remove(p_elem)

                new_p = etree.SubElement(txBody, _qn('a:p'))
                if cell_text:
                    new_r = etree.SubElement(new_p, _qn('a:r'))
                    rPr   = etree.SubElement(new_r, _qn('a:rPr'))
                    rPr.set('lang', 'ko-KR')
                    rPr.set('sz', '800')
                    if is_header:
                        rPr.set('b', '1')
                    # 글자 색상
                    sf  = etree.SubElement(rPr, _qn('a:solidFill'))
                    clr = etree.SubElement(sf,  _qn('a:srgbClr'))
                    clr.set('val', COLOR_HEADER_TEXT if is_header else COLOR_BODY_TEXT)
                    new_t = etree.SubElement(new_r, _qn('a:t'))
                    new_t.text = cell_text

        return gf

    except Exception as exc:
        print(f"[경고] _add_table_to_slide 실패: {exc}")
        return None


# ─────────────────────────────────────────────
# (d) 내용(Content) 슬라이드 생성 함수
# ─────────────────────────────────────────────
def build_content_slide(prs, title: str, subtitle: str, body_text: str = "",
                         has_image: bool = False, business_name: str = "",
                         tables: list = None):
    """
    내용(Content) 슬라이드를 만들어 prs에 추가합니다.
    PDF/Word에서 추출한 제목·부제목·본문을 채워 넣습니다.

    레이아웃 구조 (레이아웃.pptx 슬라이드[5] / [10] 기준):
      - 제목 플레이스홀더   (L=1.09, T=1.11cm)  : "01  섹션제목" 자리
      - 부제목 플레이스홀더 (L=1.09, T=1.64cm)  : "1.1 소항목" 자리
      - 본문 텍스트박스     (L=1.09, T=2.29cm)  : PDF에서 추출한 내용 텍스트
      - 테이블·이미지 등   : 레이아웃 고정 (건들지 않음)

    Parameters
    ----------
    prs        : create_blank_presentation()으로 만든 Presentation 객체
    title      : 제목 텍스트 (예: "01  자금조달 구조")
    subtitle   : 부제목 텍스트 (예: "1.1 현재 자금조달 구조")
    body_text  : 본문 내용 (\\n으로 줄바꿈, PDF/Word에서 추출)
    has_image  : True이면 content_with_image 레이아웃 사용, False면 content 레이아웃

    Returns
    -------
    추가된 슬라이드 객체
    """
    layout_type = "content_with_image" if has_image else "content"
    slide = clone_slide_layout(prs, layout_type, skip_graphic_frames=True)

    # 왼쪽 여백(L=1.09cm)에 위치하고 T < 4cm 이내인 텍스트 도형을 top 순으로 정렬
    # → 정렬 결과: [0]=제목, [1]=부제목, [2]=본문
    # 테이블·이미지·하단 사업명·페이지번호 등은 이 조건에 해당 안 됨 (건들지 않음)
    # 여기를 수정하면 제목/부제목/본문 박스 감지 기준이 바뀝니다
    LEFT_MARGIN_CM = 1.09   # 레이아웃 기준 왼쪽 여백
    TOP_LIMIT_CM   = 4.0    # 이 높이 아래 도형은 헤더 영역이 아님
    TOLERANCE      = _Cm(0.5)

    header_shapes = sorted(
        [sh for sh in slide.shapes
         if sh.has_text_frame
         and abs(sh.left - _Cm(LEFT_MARGIN_CM)) < TOLERANCE
         and sh.top / 360000 < TOP_LIMIT_CM],
        key=lambda s: s.top
    )

    # 제목 교체
    if len(header_shapes) >= 1:
        _replace_text_frame_content(header_shapes[0].text_frame, title)

    # 부제목 교체
    if len(header_shapes) >= 2:
        _replace_text_frame_content(header_shapes[1].text_frame, subtitle)

    # 본문 텍스트 표시 여부 — 표가 있어도 짧은 본문(▶ 불릿 등)은 표 위에 표시
    _BODY_WITH_TABLES_LIMIT = 300   # 이 글자 수 미만이면 표와 함께 본문 표시
    show_body = bool(body_text and body_text.strip())
    body_above_table = show_body and bool(tables) and len(body_text.strip()) < _BODY_WITH_TABLES_LIMIT

    if len(header_shapes) >= 3:
        if show_body and (not tables or body_above_table):
            _replace_text_frame_content(header_shapes[2].text_frame, body_text)
            auto_resize_text_to_fit(header_shapes[2].text_frame, max_size=10.5, min_size=8.0)
        else:
            _replace_text_frame_content(header_shapes[2].text_frame, "")

    # ── 표 삽입 (편집 가능한 PPT 표 객체) ────────────────────
    if tables:
        from modules.constants import SLIDE_WIDTH, SLIDE_HEIGHT
        SLIDE_W_CM      = SLIDE_WIDTH  / 360000   # 27.52
        SLIDE_H_CM      = SLIDE_HEIGHT / 360000   # 19.05
        TABLE_LEFT      = 1.09
        TABLE_WIDTH     = SLIDE_W_CM - TABLE_LEFT - 0.5   # ≈ 25.93
        # 본문이 표 위에 표시될 때는 표 시작 위치를 낮춤
        TABLE_TOP_START = 4.5 if body_above_table else 3.5
        FOOTER_RESERVE  = 1.8    # 푸터·하단 여백
        SPACING         = 0.3    # 표 간 간격

        n_tables    = len(tables)
        available_h = SLIDE_H_CM - TABLE_TOP_START - FOOTER_RESERVE  # ≈ 13.75
        per_h       = (available_h - SPACING * (n_tables - 1)) / n_tables

        current_top = TABLE_TOP_START
        for table_data in tables:
            if not table_data:
                continue
            row_count = len(table_data)
            # 최소 높이: 행당 0.55cm 보장
            table_h = max(per_h, row_count * 0.55)
            table_h = min(table_h, SLIDE_H_CM - current_top - FOOTER_RESERVE)
            if table_h < 0.5:
                break
            _add_table_to_slide(slide, table_data,
                                 TABLE_LEFT, current_top,
                                 TABLE_WIDTH, table_h)
            current_top += table_h + SPACING

    _replace_footer_business_name(slide, business_name)
    return slide


# ─────────────────────────────────────────────
# (f) 섹션 구분 슬라이드 — 안쪽 연두 원 좌표 (실측)
# outer 파란 링(#2E75B6)은 건드리지 않음.
# inner 연두 원(#92D050) 위치·크기에 사진을 덮어씌움.
# 각 튜플: (inner_left_cm, inner_top_cm, inner_w_cm, inner_h_cm)
# ─────────────────────────────────────────────
_DIV_INNER_OVALS = [
    # Slot 1 — 우측 상단  (타원 14, fill=#92D050)
    (17.0734, 1.3751, 8.2217, 8.4607),
    # Slot 2 — 우측 하단  (타원 2,  fill=#92D050)
    (17.9979, 8.7912, 7.4462, 7.5643),
    # Slot 3 — 중앙 좌측  (타원 15, fill=#92D050)
    (14.1443, 6.6548, 5.7571, 5.7933),
]

# 연두색 기준값 및 허용 오차
_GREEN_RGB   = (0x92, 0xD0, 0x50)   # #92D050
_GREEN_TOL   = 30                    # 각 채널 ±30
_OVAL_TOL    = _Cm(0.3)


def _detect_inner_green_ovals(slide) -> list:
    """
    슬라이드에서 연두색(#92D050 ±30) 채움 oval을 모두 찾아
    [(left_cm, top_cm, width_cm, height_cm), ...] 를 반환합니다.
    콘솔에 식별 결과를 출력합니다.
    """
    from pptx.dml.color import RGBColor

    result = []
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill.type != 1:          # solid fill 만 대상
                continue
            rgb = fill.fore_color.rgb   # RGBColor
            r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
            gr, gg, gb = _GREEN_RGB
            if abs(r - gr) <= _GREEN_TOL and abs(g - gg) <= _GREEN_TOL and abs(b - gb) <= _GREEN_TOL:
                l_cm = shape.left   / 914400 * 2.54
                t_cm = shape.top    / 914400 * 2.54
                w_cm = shape.width  / 914400 * 2.54
                h_cm = shape.height / 914400 * 2.54
                result.append((l_cm, t_cm, w_cm, h_cm))
        except Exception:
            pass

    # _DIV_INNER_OVALS 순서(Slot1→2→3)에 맞게 재정렬 — XML 내 shape 순서와 무관하게 고정
    if len(result) == len(_DIV_INNER_OVALS):
        ordered = []
        for hl, ht, _hw, _hh in _DIV_INNER_OVALS:
            best = min(result, key=lambda r: (r[0] - hl) ** 2 + (r[1] - ht) ** 2)
            ordered.append(best)
        result = ordered

    print(f"[_detect_inner_green_ovals] 식별된 연두 oval: {len(result)}개")
    for i, (l, t, w, h) in enumerate(result):
        print(f"  Slot {i+1}: L={l:.3f} T={t:.3f} W={w:.3f} H={h:.3f}")
    return result

# 템플릿의 소제목 그룹 위치 — 부정확한 템플릿 텍스트를 교체하기 위해 항상 제거
_DIV_SUBTITLE_GROUP_LEFT = 3.5328   # cm
_DIV_SUBTITLE_GROUP_TOP  = 9.9047   # cm
_DIV_SUBTITLE_GROUP_TOL  = _Cm(0.5)


# ─────────────────────────────────────────────
# (f) 섹션 구분 슬라이드 생성 함수
# ─────────────────────────────────────────────
def build_section_divider_slide(prs, section_number: str = "", section_title: str = "",
                                 business_name: str = "",
                                 section_image_bytes_list: list = None,
                                 subtitles: list = None):
    """
    섹션 구분 슬라이드를 만들어 prs에 추가합니다.

    Parameters
    ----------
    prs                       : create_presentation_from_template()으로 만든 Presentation 객체
    section_number            : 섹션 번호 문자열 (예: "01", "02", "03", "04")
    section_title             : 섹션 제목 문자열 (예: "사모사채 개요", "금융개요")
    section_image_bytes_list  : 원형 슬롯 3개에 넣을 이미지 bytes 리스트 (None 항목은 건너뜀)
                                _DIV_OVAL_PAIRS 순서와 1:1 대응
    subtitles                 : 이 섹션의 소제목 목록 (예: ["1.1 개요", "1.2 투자구조"])
                                None이면 소제목 표시 안 함

    Returns
    -------
    추가된 슬라이드 객체
    """
    slide = clone_slide_layout(prs, "section_divider")

    if not section_number and not section_title:
        return slide

    # 섹션 번호+제목 텍스트박스 교체
    # layout_data.json 기준: left=3.2474cm, top=8.2049cm
    TARGET_LEFT = 3.2474
    TARGET_TOP  = 8.2049
    TOLERANCE   = _Cm(0.8)

    new_text = f"{section_number}  {section_title}" if section_number else section_title

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if (abs(shape.left - _Cm(TARGET_LEFT)) < TOLERANCE and
                abs(shape.top  - _Cm(TARGET_TOP))  < TOLERANCE):
            _replace_text_frame_content(shape.text_frame, new_text)
            break

    # 템플릿 소제목 그룹 제거 후 올바른 소제목 텍스트박스 삽입
    for shape in list(slide.shapes):
        if (shape.shape_type == 6 and   # GROUP
                abs(shape.left - _Cm(_DIV_SUBTITLE_GROUP_LEFT)) < _DIV_SUBTITLE_GROUP_TOL and
                abs(shape.top  - _Cm(_DIV_SUBTITLE_GROUP_TOP))  < _DIV_SUBTITLE_GROUP_TOL):
            shape.element.getparent().remove(shape.element)
            break

    if subtitles:
        sub_text   = "\n".join(subtitles)
        sub_height = max(0.8, 0.65 * len(subtitles))
        add_text_box(slide, sub_text,
                     left_cm=_DIV_SUBTITLE_GROUP_LEFT, top_cm=_DIV_SUBTITLE_GROUP_TOP,
                     width_cm=12.0, height_cm=sub_height,
                     font_name=FONT_BOLD, font_size_pt=16.0)

    # 원형 이미지 슬롯 처리
    # 도형 삭제 없음 — 안쪽 연두 oval 위에 사진을 z-order 상위로 덮어씌움
    if section_image_bytes_list:
        inner_ovals = _detect_inner_green_ovals(slide)

        # 색상 감지 실패 시 하드코딩 fallback
        if not inner_ovals:
            print("[경고] _detect_inner_green_ovals 실패 → fallback 좌표 사용")
            inner_ovals = _DIV_INNER_OVALS

        for slot_idx, img_bytes in enumerate(section_image_bytes_list):
            if not img_bytes or slot_idx >= len(inner_ovals):
                continue

            il, it, iw, ih = inner_ovals[slot_idx]

            try:
                circ_png = make_circular_image_png(img_bytes, output_size=512,
                                                    border_color_rgb=(255, 255, 255),
                                                    border_width_px=8)
                add_image(slide, circ_png,
                          left_cm=il, top_cm=it,
                          width_cm=iw, height_cm=ih)
            except Exception as exc:
                print(f"[경고] 섹션 원형 이미지 삽입 실패 slot={slot_idx}: {exc}")

    _replace_footer_business_name(slide, business_name)
    return slide


# ─────────────────────────────────────────────
# (g) 연락처 슬라이드 생성 함수
# ─────────────────────────────────────────────
def build_contact_slide(prs):
    """
    연락처(Contact) 슬라이드를 만들어 prs에 추가합니다.
    레이아웃.pptx의 마지막 슬라이드(index 12)를 그대로 복제합니다.
    배경 이미지·테이블·텍스트 모두 고정 디자인이므로 교체하지 않습니다.

    Parameters
    ----------
    prs : create_presentation_from_template()으로 만든 Presentation 객체

    Returns
    -------
    추가된 슬라이드 객체
    """
    return clone_slide_layout(prs, "contact")


# ─────────────────────────────────────────────
# [출력 파일명 생성 헬퍼]
# ─────────────────────────────────────────────
def make_output_filename(business_name: str) -> str:
    """
    OUTPUT_FILENAME_FORMAT에 따라 실제 파일명을 생성합니다.
    예) "[Rainfield] 천안 부성2지구 도시개발사업_260512.pptx"

    여기를 수정하면 파일명 날짜 형식이 바뀝니다
    """
    date_str = datetime.now().strftime("%y%m%d")  # YYMMDD 형식
    filename = OUTPUT_FILENAME_FORMAT.format(
        business_name=business_name,
        date=date_str,
    )
    return filename


# ─────────────────────────────────────────────
# (h) 전체 PPT 조립 — 핵심 오케스트레이터
# ─────────────────────────────────────────────
def build_full_presentation(
    business_name: str,
    year: str,
    month_en: str,
    pages: list,
    cover_image_bytes: bytes = None,
    executive_summary_sections: list = None,
    section_image_bytes_list: list = None,
    toc_count: int = None,
    toc_image_bytes_list: list = None,
    toc_map: dict = None,
) -> bytes:
    """
    content_parser 결과를 받아 완성된 PPT를 bytes로 반환합니다.

    Parameters
    ----------
    business_name              : 사업명 (표지·하단 푸터 텍스트)
    year                       : 연도 문자열 (예: "2026")
    month_en                   : 영문 월 (예: "May")
    pages                      : content_parser.parse_document() 반환값.
                                 5형식의 경우 remap_pages_for_5sections() 결과를 전달.
    cover_image_bytes          : 표지 메인 이미지 bytes (None이면 회색 박스)
    executive_summary_sections : Executive Summary 섹션 목록 (최대 3개)
                                 None이면 ES 슬라이드 생략
    section_image_bytes_list   : [bytes|None, bytes|None, bytes|None]
                                 섹션 divider 원형 슬롯 3개 이미지 (모든 섹션 공통)
    toc_count                  : 목차 항목 수 강제 지정 (4 또는 5).
                                 None이면 pages에서 자동 감지.
    toc_map                    : {"01": ["1.1 ...", ...], ...} — 섹션별 소제목 목록.
                                 None이면 pages에서 자동 추출 (_build_toc_map).
                                 split_into_5_sections() 결과를 그대로 전달하면 됩니다.

    Returns
    -------
    완성된 PPT 파일 bytes (BytesIO → Streamlit download_button에 바로 전달 가능)
    """
    # ── 0. 템플릿 기반 프레젠테이션 생성 ─────────────────────
    prs = create_presentation_from_template()
    template_count = len(prs.slides)   # 제거할 원본 슬라이드 수 기록

    # ── 1. 표지 ────────────────────────────────────────────
    build_cover_slide(prs, business_name, year, month_en, cover_image_bytes)

    # ── 2. Executive Summary (선택) ───────────────────────
    if executive_summary_sections:
        build_executive_summary_slide(prs, executive_summary_sections, business_name)

    # ── 3. 목차 ────────────────────────────────────────────
    unique_sections = _count_unique_sections(pages)
    # toc_map이 외부에서 주어지면 그대로 사용; 없으면 pages에서 자동 추출
    computed_toc_map = toc_map if toc_map is not None else _build_toc_map(pages)
    num_toc = toc_count if toc_count in (4, 5) else len(unique_sections)
    if unique_sections:
        num_toc = min(num_toc, len(unique_sections))  # 더미 섹션 방지
    build_toc_slide(prs, num_sections=num_toc, toc_map=computed_toc_map,
                    toc_image_bytes_list=toc_image_bytes_list)

    # ── 4. 섹션 구분 + 내용 슬라이드 ─────────────────────
    _build_content_block(prs, pages, business_name=business_name,
                          section_image_bytes_list=section_image_bytes_list,
                          toc_map=computed_toc_map)

    # ── 5. 연락처 슬라이드 ────────────────────────────────
    build_contact_slide(prs)

    # ── 6. 원본 템플릿 슬라이드 제거 ─────────────────────
    finalize_presentation(prs, template_count)

    # ── 7. bytes로 반환 ───────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _count_unique_sections(pages: list) -> list:
    """pages에서 중복 없는 section_title 목록을 순서 유지하며 반환합니다."""
    seen = []
    for p in pages:
        t = p.get("section_title", "").strip()
        if t and t not in seen:
            seen.append(t)
    return seen


_SEC_NUM_EXTRACT = re.compile(r'^(0[1-9])\b')


def _build_toc_map(pages: list) -> dict:
    """pages에서 섹션 번호별 부제목 목록을 추출합니다."""
    toc: dict = {}
    seen: dict = {}
    for page in pages:
        section_title = page.get("section_title", "").strip()
        subtitle      = page.get("subtitle", "").strip()
        m = _SEC_NUM_EXTRACT.match(section_title)
        if not m or not subtitle:
            continue
        sec_num = m.group(1)
        if sec_num not in seen:
            seen[sec_num] = set()
        if subtitle not in seen[sec_num]:
            seen[sec_num].add(subtitle)
            toc.setdefault(sec_num, []).append(subtitle)
    return toc


def _build_content_block(prs, pages: list, business_name: str = "",
                          section_image_bytes_list: list = None,
                          toc_map: dict = None):
    """
    pages 목록을 순회하며 섹션 구분 슬라이드 + 내용 슬라이드를 삽입합니다.

    섹션 구분 슬라이드 생성 기준:
      - section_title의 앞 부분에 "01", "02" … "09" 형태의 번호가 있고
      - 그 번호가 이전 섹션 번호와 달라졌을 때만 1회 생성
      - 전체 텍스트가 달라도 번호가 같으면 같은 섹션으로 간주

    content 슬라이드:
      - 섹션 구분자가 방금 생성됐고 본문·부제목이 모두 비어있으면 스킵
        (PDF 섹션 헤더 페이지 → divider가 대신함)
      - 그 외에는 PDF 1페이지 = PPT content 1장

    Parameters
    ----------
    section_image_bytes_list : [bytes|None, bytes|None, bytes|None]
                               섹션 divider 원형 슬롯 3개에 넣을 이미지 (4개 섹션 공통)
    toc_map                  : {"01": ["1.1 ...", ...], ...} — 섹션별 소제목 목록
    """
    _toc = toc_map or DEFAULT_TOC_MAP
    current_sec_num = ""   # 현재 섹션 번호 ("01", "02" …)
    section_num = 0

    for page in pages:
        section_title = page.get("section_title", "").strip()
        subtitle      = page.get("subtitle", "").strip()
        body_text     = page.get("body_text", "").strip()

        # section_num: content_parser가 자동 분류한 값 우선
        # 없으면 section_title 앞 부분 "0N" 패턴으로 fallback
        this_sec_num = page.get("section_num", "")
        if not this_sec_num:
            m = _SEC_NUM_EXTRACT.match(section_title)
            this_sec_num = m.group(1) if m else ""

        # 새 섹션 판단
        is_new_section = bool(this_sec_num and this_sec_num != current_sec_num)

        if is_new_section:
            section_num += 1
            num_str = f"{section_num:02d}"
            # divider 제목: section_label 우선, 없으면 section_title
            divider_title = page.get("section_label") or section_title or ""
            build_section_divider_slide(prs, num_str, divider_title,
                                        business_name=business_name,
                                        section_image_bytes_list=section_image_bytes_list,
                                        subtitles=_toc.get(this_sec_num, []))
            current_sec_num = this_sec_num
            # 섹션 첫 페이지에 본문·부제목 없으면 content slide 스킵
            if not body_text and not subtitle:
                continue
        else:
            num_str = f"{section_num:02d}" if section_num > 0 else ""

        # 표시할 제목: 현재 섹션 번호 + section_title
        display_title = f"{num_str}  {section_title}" if num_str and section_title else section_title

        # 표가 있는 페이지: text_without_tables 사용
        # 표가 없는 페이지: 기존 body_text 그대로
        tables_data = page.get("tables") or []
        body_to_use = (page.get("text_without_tables") or body_text) if tables_data else body_text

        build_content_slide(
            prs,
            title=display_title,
            subtitle=subtitle,
            body_text=body_to_use,
            tables=tables_data,
            business_name=business_name,
        )


# ─────────────────────────────────────────────
# (i) 미리보기 PPT — 표지 + 섹션 divider 4장
# ─────────────────────────────────────────────
_PREVIEW_SECTIONS = [
    ("01", "사모사채 개요"),
    ("02", "금융개요"),
    ("03", "본건 사업 개요"),
    ("04", "Appendix"),
]


def build_preview_presentation(
    business_name: str,
    year: str,
    month_en: str,
    cover_image_bytes: bytes = None,
    section_image_bytes_list: list = None,
    toc_image_bytes_list: list = None,
) -> bytes:
    """
    표지 1장 + 목차 1장 + 섹션 구분 슬라이드 4장 = 총 6장짜리 미리보기 PPT를 반환합니다.
    실제 콘텐츠 슬라이드 없이 디자인 확인용으로 사용합니다.

    Parameters
    ----------
    business_name            : 사업명
    year / month_en          : 표지 날짜
    cover_image_bytes        : 표지 이미지 bytes
    section_image_bytes_list : 섹션 divider 원형 슬롯 3개 이미지 리스트
    toc_image_bytes_list     : 목차 원형 슬롯 1개 이미지 리스트

    Returns
    -------
    완성된 PPT 파일 bytes
    """
    prs = create_presentation_from_template()
    template_count = len(prs.slides)

    build_cover_slide(prs, business_name, year, month_en, cover_image_bytes)

    build_toc_slide(prs, num_sections=4, toc_map=DEFAULT_TOC_MAP,
                    toc_image_bytes_list=toc_image_bytes_list)

    for sec_num, sec_label in _PREVIEW_SECTIONS:
        build_section_divider_slide(
            prs, sec_num, sec_label,
            business_name=business_name,
            section_image_bytes_list=section_image_bytes_list,
            subtitles=DEFAULT_TOC_MAP.get(sec_num, []),
        )

    finalize_presentation(prs, template_count)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
