"""
ppt_generator.py
─────────────────────────────────────────────────────────
PPT 슬라이드에 텍스트·이미지·도형을 추가하는 기본 함수 모음.
page_builders.py에서 이 함수들을 호출해서 실제 슬라이드를 만듭니다.
─────────────────────────────────────────────────────────
"""

import io
import os
import re

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw

from modules.constants import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    FONT_BOLD, FONT_LIGHT,
)


# ─────────────────────────────────────────────
# [파일명에 사용할 수 없는 문자 치환]
# Windows 파일명 금지 문자: \ / : * ? " < > |
# ─────────────────────────────────────────────
_INVALID_FILENAME_CHARS = re.compile(r'[\\/:*?"<>|]')

def _sanitize_filename(name: str) -> str:
    """파일명에 사용 불가한 특수문자를 '_'로 치환합니다."""
    return _INVALID_FILENAME_CHARS.sub("_", name)


# ─────────────────────────────────────────────
# [단위 변환 헬퍼]
# ─────────────────────────────────────────────
def _cm(value: float) -> Emu:
    """cm 값을 python-pptx EMU 단위로 변환"""
    return Cm(value)

def _pt(value: float) -> Emu:
    """pt 값을 python-pptx EMU 단위로 변환"""
    return Pt(value)

def _align_enum(align: str):
    """문자열 정렬값을 PP_ALIGN enum으로 변환"""
    mapping = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(align, PP_ALIGN.LEFT)


# ─────────────────────────────────────────────
# (a) 빈 프레젠테이션 생성
# ─────────────────────────────────────────────
def create_blank_presentation() -> Presentation:
    """
    A4 세로 크기의 빈 프레젠테이션 객체를 만들어 반환합니다.

    슬라이드 크기: 가로 21cm × 세로 29.7cm (constants.py에서 설정)
    슬라이드는 0장 — page_builders.py에서 슬라이드를 추가합니다.
    """
    prs = Presentation()

    # A4 크기 강제 설정
    # 여기를 수정하면 PPT 페이지 크기가 바뀝니다 (constants.py의 값을 바꾸는 것을 권장)
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    return prs


# ─────────────────────────────────────────────
# (b) 텍스트 박스 추가
# ─────────────────────────────────────────────
def add_text_box(
    slide,
    text: str,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    font_name: str   = None,
    font_size_pt: float = 10.5,
    font_bold: bool  = False,
    font_color: RGBColor = None,
    align: str       = "left",
    line_spacing_pt: float = None,
):
    """
    슬라이드에 텍스트 박스를 추가합니다.

    Parameters
    ----------
    slide        : 추가할 슬라이드 객체
    text         : 표시할 텍스트 (\n으로 줄바꿈 가능)
    left_cm      : 왼쪽 여백 (cm)
    top_cm       : 위쪽 여백 (cm)
    width_cm     : 텍스트 박스 가로 크기 (cm)
    height_cm    : 텍스트 박스 세로 크기 (cm)
    font_name    : 폰트 이름 (None이면 FONT_LIGHT 사용)
    font_size_pt : 글자 크기 (pt)
    font_bold    : 굵게 여부
    font_color   : RGBColor 객체 (None이면 검정)
    align        : 정렬 "left" / "center" / "right"
    line_spacing_pt : 줄 간격 (pt), None이면 기본값
    """
    # 입력값 검증
    if width_cm <= 0 or height_cm <= 0:
        raise ValueError(f"텍스트 박스 크기가 0 이하입니다: width={width_cm}, height={height_cm}")
    if not text:
        text = ""

    # 폰트 기본값 설정
    if font_name is None:
        font_name = FONT_LIGHT  # 여기를 수정하면 기본 폰트가 바뀝니다

    txBox = slide.shapes.add_textbox(
        _cm(left_cm), _cm(top_cm),
        _cm(width_cm), _cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # 줄바꿈 처리 (\n 기준으로 분리)
    lines = text.split("\n")
    for line_idx, line in enumerate(lines):
        if line_idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.alignment = _align_enum(align)

        # 줄 간격 설정
        if line_spacing_pt is not None:
            para.line_spacing = _pt(line_spacing_pt)

        run = para.add_run()
        run.text = line

        # 폰트 적용
        font = run.font
        font.name = font_name
        font.size = _pt(font_size_pt)
        font.bold = font_bold
        if font_color is not None:
            font.color.rgb = font_color

    return txBox


# ─────────────────────────────────────────────
# (c) 이미지 추가
# ─────────────────────────────────────────────
def add_image(
    slide,
    image_source,           # 파일 경로(str) 또는 bytes/BytesIO
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    border_color: RGBColor = None,
    border_width_pt: float = None,
    fit_mode: str = "contain",   # "cover" 또는 "contain"
):
    """
    슬라이드에 이미지를 삽입합니다.

    Parameters
    ----------
    image_source   : 이미지 파일 경로(str) 또는 bytes/BytesIO 객체
    left_cm        : 왼쪽 여백 (cm)
    top_cm         : 위쪽 여백 (cm)
    width_cm       : 가로 크기 (cm)
    height_cm      : 세로 크기 (cm)
    border_color   : 테두리 색 (RGBColor, None이면 테두리 없음)
    border_width_pt: 테두리 두께 (pt)
    fit_mode       : "contain" = 비율 유지하며 박스 안에 맞춤
                     "cover"   = 비율 유지하며 박스를 꽉 채움 (잘림 발생)
    """
    if width_cm <= 0 or height_cm <= 0:
        raise ValueError(f"이미지 크기가 0 이하입니다: width={width_cm}, height={height_cm}")

    # bytes/BytesIO로 통일
    if isinstance(image_source, (bytes, bytearray)):
        img_stream = io.BytesIO(image_source)
    elif isinstance(image_source, io.BytesIO):
        image_source.seek(0)
        img_stream = image_source
    elif isinstance(image_source, str):
        if not os.path.exists(image_source):
            raise FileNotFoundError(f"이미지 파일을 찾을 수 없습니다: {image_source}")
        with open(image_source, "rb") as f:
            img_stream = io.BytesIO(f.read())
    else:
        raise TypeError(f"image_source 타입이 올바르지 않습니다: {type(image_source)}")

    # fit_mode 적용 (PIL로 전처리)
    img_stream = _fit_image(img_stream, width_cm, height_cm, fit_mode)

    pic = slide.shapes.add_picture(
        img_stream,
        _cm(left_cm), _cm(top_cm),
        _cm(width_cm), _cm(height_cm)
    )

    # 테두리 추가
    if border_color is not None and border_width_pt is not None:
        _add_border_to_shape(pic, border_color, border_width_pt)

    return pic


def _fit_image(img_stream: io.BytesIO, width_cm: float, height_cm: float, mode: str) -> io.BytesIO:
    """PIL로 이미지를 지정 비율에 맞게 조정합니다. PNG 투명도를 보존합니다."""
    try:
        img = Image.open(img_stream)

        # PNG 투명도(알파 채널) 보존 — 투명 영역이 검정으로 바뀌는 문제 방지
        # 여기를 수정하면 투명도 처리 방식이 바뀝니다
        has_alpha = img.mode in ("RGBA", "LA", "PA") or (
            img.mode == "P" and "transparency" in img.info
        )
        img = img.convert("RGBA") if has_alpha else img.convert("RGB")

        target_ratio = width_cm / height_cm
        img_ratio    = img.width / img.height

        if mode == "cover":
            # 박스를 꽉 채움 (넘치는 부분 크롭)
            if img_ratio > target_ratio:
                new_h = img.height
                new_w = int(new_h * target_ratio)
                left  = (img.width - new_w) // 2
                img   = img.crop((left, 0, left + new_w, new_h))
            else:
                new_w = img.width
                new_h = int(new_w / target_ratio)
                top   = (img.height - new_h) // 2
                img   = img.crop((0, top, new_w, top + new_h))
        # contain 모드는 python-pptx가 자동으로 비율 유지

        # PPT가 이미지를 늘릴 때 흐릿해지지 않도록 최소 해상도 보장
        # 긴 쪽이 MIN_PX 미만이면 LANCZOS로 업스케일
        # 여기를 수정하면 최소 해상도 기준이 바뀝니다
        MIN_PX = 1800
        longer = max(img.width, img.height)
        if longer < MIN_PX:
            scale   = MIN_PX / longer
            new_w   = int(img.width  * scale)
            new_h   = int(img.height * scale)
            img     = img.resize((new_w, new_h), Image.LANCZOS)

        output = io.BytesIO()
        img.save(output, format="PNG")
        output.seek(0)
        return output
    except Exception:
        img_stream.seek(0)
        return img_stream


def _add_border_to_shape(shape, color: RGBColor, width_pt: float):
    """도형 XML을 직접 수정해서 테두리를 추가합니다."""
    try:
        sp_pr = shape.element.spPr
        ln = etree.SubElement(sp_pr, qn("a:ln"))
        ln.set("w", str(int(Pt(width_pt))))
        solid_fill = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srgb.set("val", "{:02X}{:02X}{:02X}".format(color.r, color.g, color.b))
    except Exception:
        pass


# ─────────────────────────────────────────────
# (d) 원형 이미지 추가
# ─────────────────────────────────────────────
def add_oval_with_image(
    slide,
    image_bytes: bytes,
    left_cm: float,
    top_cm: float,
    diameter_cm: float,
    border_color: RGBColor = None,
    border_width_pt: float = 2.0,
):
    """
    PIL로 이미지를 원형으로 크롭한 뒤 슬라이드에 삽입합니다.
    레이아웃의 연두색 동그라미 이미지 자리에 사용합니다.

    Parameters
    ----------
    image_bytes    : 원형으로 자를 이미지 bytes
    left_cm        : 원 왼쪽 여백 (cm)
    top_cm         : 원 위쪽 여백 (cm)
    diameter_cm    : 원 지름 (cm) — 가로=세로 동일
    border_color   : 테두리 색 (RGBColor)
    border_width_pt: 테두리 두께 (pt)
    """
    if not image_bytes:
        raise ValueError("이미지 데이터가 없습니다.")
    if diameter_cm <= 0:
        raise ValueError(f"원 지름이 0 이하입니다: {diameter_cm}")

    # PIL로 원형 크롭
    cropped_bytes = _crop_circle(image_bytes)

    return add_image(
        slide, cropped_bytes,
        left_cm, top_cm,
        diameter_cm, diameter_cm,
        border_color=border_color,
        border_width_pt=border_width_pt,
    )


def _crop_circle(image_bytes: bytes) -> bytes:
    """PIL로 이미지를 정사각형으로 크롭 후 원형 마스크 적용 (PNG 투명도 사용)."""
    img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")

    # 정사각형으로 크롭
    size = min(img.width, img.height)
    left = (img.width  - size) // 2
    top  = (img.height - size) // 2
    img  = img.crop((left, top, left + size, top + size))

    # 원형 마스크 생성
    mask = Image.new("L", (size, size), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, size, size), fill=255)

    # 마스크 적용
    output_img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    output_img.paste(img, (0, 0), mask)

    buf = io.BytesIO()
    output_img.save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


def make_circular_image_png(
    image_bytes: bytes,
    output_size: int = 512,
    border_color_rgb: tuple = (255, 255, 255),
    border_width_px: int = 8,
) -> bytes:
    """
    이미지를 원형으로 크롭하고 완전 불투명 흰색 테두리를 그린 PNG bytes를 반환합니다.
    2× 오버샘플링으로 안티앨리어싱 적용.

    Parameters
    ----------
    image_bytes      : 원본 이미지 bytes
    output_size      : 출력 이미지 크기 (px, 정사각형)
    border_color_rgb : 테두리 색상 (R, G, B) — alpha는 항상 255
    border_width_px  : 테두리 두께 (px) — 기본 8px = inner oval 선 6pt와 동일
    """
    img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")

    # 정사각형 중앙 크롭 후 리사이즈
    sq = min(img.width, img.height)
    lx = (img.width  - sq) // 2
    ty = (img.height - sq) // 2
    img = img.crop((lx, ty, lx + sq, ty + sq)).resize(
        (output_size, output_size), Image.LANCZOS
    )

    r, g, b = border_color_rgb
    aa = output_size * 2          # 2× 오버샘플링 크기
    pi = border_width_px * 2      # 오버샘플 공간에서의 inset

    # ── 1) 테두리 원 레이어: filled 원 전체 → 안티앨리어싱 후 출력 크기로 축소
    border_aa = Image.new("RGBA", (aa, aa), (0, 0, 0, 0))
    ImageDraw.Draw(border_aa).ellipse((0, 0, aa - 1, aa - 1), fill=(r, g, b, 255))
    border_layer = border_aa.resize((output_size, output_size), Image.LANCZOS)

    # ── 2) 사진 원형 마스크: inset 안쪽 circle → 안티앨리어싱 후 출력 크기로 축소
    mask_aa = Image.new("L", (aa, aa), 0)
    ImageDraw.Draw(mask_aa).ellipse((pi, pi, aa - 1 - pi, aa - 1 - pi), fill=255)
    photo_mask = mask_aa.resize((output_size, output_size), Image.LANCZOS)

    # ── 3) border_layer 위에 사진을 photo_mask 영역에 붙이기
    canvas = border_layer.copy()
    canvas.paste(img, (0, 0), photo_mask)

    buf = io.BytesIO()
    canvas.save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# (e) PPT 파일 저장
# ─────────────────────────────────────────────
def save_presentation(prs: Presentation, output_path: str) -> str:
    """
    PPT 파일을 저장합니다.
    파일명에 사용 불가 문자가 있으면 자동으로 '_'로 치환합니다.

    Parameters
    ----------
    prs         : 저장할 Presentation 객체
    output_path : 저장 경로 (폴더 포함 전체 경로)

    Returns
    -------
    실제 저장된 파일 경로 (문자 치환 후)
    """
    # 폴더 부분과 파일명 분리
    folder    = os.path.dirname(output_path)
    filename  = os.path.basename(output_path)

    # 파일명 특수문자 치환 (폴더 경로는 건드리지 않음)
    safe_name = _sanitize_filename(filename)
    safe_path = os.path.join(folder, safe_name) if folder else safe_name

    # 폴더 없으면 자동 생성
    if folder:
        os.makedirs(folder, exist_ok=True)

    prs.save(safe_path)
    return safe_path
