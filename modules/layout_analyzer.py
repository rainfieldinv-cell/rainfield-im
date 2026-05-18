"""
layout_analyzer.py
─────────────────────────────────────────────────────────
templates/레이아웃.pptx 파일을 분석해서 각 슬라이드의
도형·텍스트·색상·폰트 정보를 JSON으로 추출·저장합니다.

사용법 (터미널에서 한 번만 실행):
  python -c "from modules.layout_analyzer import run_analysis; run_analysis()"

결과 파일: templates/layout_data.json
─────────────────────────────────────────────────────────
"""

import json
import re
import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from modules.constants import LAYOUT_PPTX_PATH, LAYOUT_JSON_PATH


# ─────────────────────────────────────────────
# [단위 변환 헬퍼]
# EMU(English Metric Units) → cm 변환
# python-pptx는 내부적으로 EMU 단위 사용
# ─────────────────────────────────────────────
def _emu_to_cm(emu: int) -> float:
    """EMU → cm 변환 (소수점 4자리)"""
    return round(emu / 360000, 4)

def _emu_to_pt(emu: int) -> float:
    """EMU → pt 변환"""
    return round(emu / 12700, 2)

def _rgb_to_hex(rgb) -> str:
    """RGBColor 객체 → '#RRGGBB' 문자열 변환"""
    try:
        return "#{:02X}{:02X}{:02X}".format(rgb.r, rgb.g, rgb.b)
    except Exception:
        return None


def _extract_run_color(run) -> tuple:
    """
    런(run) XML에서 색상을 추출합니다.
    반환값: (rgb_hex, theme_name)
      - rgb_hex  : '#RRGGBB' 문자열 (명시적 RGB일 때) 또는 None
      - theme_name: 'dk1', 'bg1' 등 테마 컬러 이름 또는 None
    여기를 수정하면 색상 추출 방식이 바뀝니다
    """
    try:
        rpr = run._r.find(qn('a:rPr'))
        if rpr is None:
            return None, None
        solid = rpr.find(qn('a:solidFill'))
        if solid is None:
            return None, None
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            val = srgb.get('val', '')
            return (f"#{val.upper()}" if val else None), None
        scheme = solid.find(qn('a:schemeClr'))
        if scheme is not None:
            return None, scheme.get('val')
    except Exception:
        pass
    return None, None


# ─────────────────────────────────────────────
# [도형 타입 판별]
# ─────────────────────────────────────────────
def _detect_shape_type(shape) -> str:
    """
    python-pptx 도형 객체를 보고 종류를 문자열로 반환합니다.
    반환값: "text_box" / "image" / "rectangle" / "oval" / "line" / "group" / "other"
    """
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        if st == MSO_SHAPE_TYPE.LINE:
            return "line"
        if st == MSO_SHAPE_TYPE.GROUP:
            return "group"
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "text_box"
        # 자동 도형(사각형·타원 등) 구분
        if hasattr(shape, "auto_shape_type"):
            auto = str(shape.auto_shape_type)
            if "OVAL" in auto or "ELLIPSE" in auto:
                return "oval"
            if "RECT" in auto or "ROUNDED" in auto:
                return "rectangle"
        # 텍스트가 있으면 text_box로 간주
        if shape.has_text_frame:
            return "text_box"
    except Exception:
        pass
    return "other"


# ─────────────────────────────────────────────
# [텍스트 정렬 판별]
# ─────────────────────────────────────────────
def _detect_align(paragraph) -> str:
    """단락의 텍스트 정렬 방향을 문자열로 반환"""
    try:
        align = paragraph.alignment
        if align == PP_ALIGN.CENTER:
            return "center"
        if align == PP_ALIGN.RIGHT:
            return "right"
    except Exception:
        pass
    return "left"


# ─────────────────────────────────────────────
# [placeholder_hint 자동 감지]
# 텍스트 내용으로 이 도형이 어떤 역할인지 추론
# ─────────────────────────────────────────────
def _detect_placeholder_hint(text: str) -> str:
    """
    텍스트 내용을 보고 도형의 역할(placeholder 종류)을 추론합니다.
    여기에 패턴을 추가하면 더 많은 종류를 인식합니다.
    """
    if not text:
        return None
    t = text.strip()

    # 패턴 매핑 (텍스트 포함 여부로 판별)
    # 여기를 수정하면 placeholder 감지 기준이 바뀝니다
    hint_map = [
        (["사업명"],                         "business_name"),
        (["그해 년도", "년도", "YYYY"],       "year"),
        (["몇월", "영어 표시", "Month"],      "month_en"),
        (["부제목", "Sub"],                   "subtitle"),
        (["제목", "Title"],                   "title"),
        (["내용", "Content"],                 "content"),
        (["Disclaimer", "면책", "본 자료"],   "disclaimer"),
        (["CONTENTS", "목차"],               "toc"),
        (["연락처", "Contact"],               "contact"),
        (["Information Memorandum", "IM"],   "im_label"),
    ]

    for keywords, hint in hint_map:
        for kw in keywords:
            if kw in t:
                return hint

    # "00", "XX", "OO" 같은 숫자 자리표시자
    if re.fullmatch(r"[0O]{2,4}|XX+|00+|\d{2}", t):
        return "number_placeholder"

    return None


# ─────────────────────────────────────────────
# [슬라이드 타입 자동 추론]
# ─────────────────────────────────────────────
def _detect_slide_type(slide, slide_index: int, total_slides: int) -> str:
    """
    슬라이드 내 텍스트 패턴을 보고 슬라이드 종류를 추론합니다.
    반환값: "cover" / "toc" / "section_divider" / "contact" / "content"
    """
    # 슬라이드 전체 텍스트 수집
    all_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                all_text += para.text + "\n"

    # 1번 슬라이드 = 표지
    if slide_index == 0:
        return "cover"

    # "CONTENTS" 또는 "목차" 텍스트 포함 = 목차
    if "CONTENTS" in all_text or "목차" in all_text:
        return "toc"

    # 마지막 슬라이드에 연락처 관련 텍스트 = 연락처 페이지
    if slide_index >= total_slides - 2:
        if "연락처" in all_text or "본 자료와 관련하여" in all_text or "Contact" in all_text:
            return "contact"

    # "01", "02" 같은 섹션 번호 + 사업 관련 키워드 = 섹션 구분 페이지
    if re.search(r"\b0[1-9]\b", all_text):
        return "section_divider"

    return "content"


# ─────────────────────────────────────────────
# [도형 1개 분석]
# ─────────────────────────────────────────────
def _analyze_shape(shape, shape_index: int) -> dict:
    """도형 1개의 모든 정보를 딕셔너리로 추출합니다."""
    result = {
        "shape_index": shape_index,
        "shape_name":  shape.name,
        "shape_type":  _detect_shape_type(shape),
        "left_cm":     _emu_to_cm(shape.left)   if shape.left   is not None else 0,
        "top_cm":      _emu_to_cm(shape.top)    if shape.top    is not None else 0,
        "width_cm":    _emu_to_cm(shape.width)  if shape.width  is not None else 0,
        "height_cm":   _emu_to_cm(shape.height) if shape.height is not None else 0,
        "fill_color":  None,
        "line_color":  None,
        "line_width_pt": None,
        "text":        "",
        "font_name":   None,
        "font_size_pt": None,
        "font_bold":   None,
        "font_color":       None,   # 명시적 RGB색상 '#RRGGBB'
        "font_color_theme": None,   # 테마 컬러 이름 (예: 'dk1', 'bg1')
        "text_align":  None,
        "placeholder_hint": None,
    }

    # ── 채우기 색상 ──
    try:
        fill = shape.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            result["fill_color"] = _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass

    # ── 테두리 색상·두께 ──
    try:
        line = shape.line
        if line.color and line.color.rgb:
            result["line_color"]    = _rgb_to_hex(line.color.rgb)
            result["line_width_pt"] = _emu_to_pt(line.width) if line.width else None
    except Exception:
        pass

    # ── 텍스트 및 폰트 ──
    if shape.has_text_frame:
        paragraphs = shape.text_frame.paragraphs
        full_text  = "\n".join(p.text for p in paragraphs)
        result["text"] = full_text

        # 첫 번째 런(run)에서 폰트 정보 추출
        for para in paragraphs:
            result["text_align"] = _detect_align(para)
            for run in para.runs:
                font = run.font
                try:
                    result["font_name"] = font.name
                except Exception:
                    pass
                try:
                    result["font_size_pt"] = _emu_to_pt(font.size) if font.size else None
                except Exception:
                    pass
                try:
                    result["font_bold"] = font.bold
                except Exception:
                    pass
                # RGB → 실패 시 XML에서 테마 컬러 이름까지 시도
                try:
                    result["font_color"] = _rgb_to_hex(font.color.rgb)
                except Exception:
                    pass
                if result["font_color"] is None:
                    rgb_hex, theme_name = _extract_run_color(run)
                    result["font_color"]       = rgb_hex
                    result["font_color_theme"] = theme_name
                break  # 첫 번째 런만 대표로 사용
            if result["font_name"]:
                break

        result["placeholder_hint"] = _detect_placeholder_hint(full_text)

    return result


# ─────────────────────────────────────────────
# [메인 분석 함수]
# ─────────────────────────────────────────────
def analyze_layout_pptx(pptx_path: str) -> dict:
    """
    레이아웃.pptx 파일을 열어 모든 슬라이드의 도형 정보를 추출합니다.

    반환값:
    {
        "slides": [ { "slide_index": 0, "slide_type": "cover", "shapes": [...] }, ... ],
        "color_palette": ["#XXXXXX", ...],
        "fonts_used": ["피플폰트 Bold", ...]
    }
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"레이아웃 파일을 찾을 수 없습니다: {pptx_path}")

    prs          = Presentation(pptx_path)
    total_slides = len(prs.slides)
    slides_data  = []
    all_colors   = []
    all_fonts    = set()

    for slide_idx, slide in enumerate(prs.slides):
        slide_type = _detect_slide_type(slide, slide_idx, total_slides)
        shapes_data = []

        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = _analyze_shape(shape, shape_idx)
            shapes_data.append(shape_info)

            # 색상 수집
            for color_key in ("fill_color", "line_color", "font_color"):
                c = shape_info.get(color_key)
                if c and c not in all_colors:
                    all_colors.append(c)

            # 폰트 수집
            if shape_info.get("font_name"):
                all_fonts.add(shape_info["font_name"])

        slides_data.append({
            "slide_index": slide_idx,
            "slide_type":  slide_type,
            "shapes":      shapes_data,
        })

    return {
        "slides":        slides_data,
        "color_palette": all_colors,
        "fonts_used":    sorted(all_fonts),
    }


# ─────────────────────────────────────────────
# [JSON 저장 함수]
# ─────────────────────────────────────────────
def save_layout_json(layout_data: dict, output_path: str):
    """
    분석 결과를 JSON 파일로 저장합니다.
    들여쓰기 2칸 적용, UTF-8 인코딩으로 한글 정상 저장.
    """
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(layout_data, f, ensure_ascii=False, indent=2)
    print(f"[완료] 레이아웃 JSON 저장: {output_path}")


# ─────────────────────────────────────────────
# [분석 실행 + 요약 출력 함수]
# 터미널에서 python -c "from modules.layout_analyzer import run_analysis; run_analysis()"
# ─────────────────────────────────────────────
def run_analysis():
    """레이아웃.pptx를 분석하고 JSON 저장 후 요약을 출력합니다."""
    print(f"[분석 시작] {LAYOUT_PPTX_PATH}")
    data = analyze_layout_pptx(LAYOUT_PPTX_PATH)
    save_layout_json(data, LAYOUT_JSON_PATH)

    # ── 요약 출력 ──
    print("\n" + "="*50)
    print(f"총 슬라이드 수  : {len(data['slides'])}장")
    print(f"감지된 폰트     : {data['fonts_used']}")
    print(f"추출된 색상 수  : {len(data['color_palette'])}개")
    print(f"색상 팔레트     : {data['color_palette'][:10]}")  # 처음 10개만 출력
    print("="*50)

    print("\n[슬라이드별 요약]")
    for slide in data["slides"]:
        shape_count = len(slide["shapes"])
        text_shapes = [s for s in slide["shapes"] if s["text"]]
        print(f"  슬라이드 {slide['slide_index']+1:2d}  |  타입: {slide['slide_type']:18s}  "
              f"|  도형 {shape_count:2d}개  |  텍스트 도형 {len(text_shapes):2d}개")

    print("\n[표지(cover) 슬라이드 상세 - 좌표/텍스트 확인용]")
    for slide in data["slides"]:
        if slide["slide_type"] == "cover":
            for s in slide["shapes"]:
                if s["text"]:
                    print(f"  [{s['shape_type']:12s}] "
                          f"위치:({s['left_cm']:5.1f}, {s['top_cm']:5.1f})cm  "
                          f"크기:({s['width_cm']:5.1f} × {s['height_cm']:5.1f})cm  "
                          f"폰트:{s['font_name']}  크기:{s['font_size_pt']}pt  "
                          f"힌트:{s['placeholder_hint']}  "
                          f'텍스트:"{s["text"][:30]}"')
            break

    return data
