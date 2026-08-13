# -*- coding: utf-8 -*-
"""금융구조도 슬라이드 — 기본 '내용 페이지' 위에 도형으로 직접 작도.
   기존 2.1 금융 구조도(템플릿 asis) '다음 순서'에 추가되는 새 구조도 페이지."""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

_FONT = "피플폰트 Bold"
NAVY = RGBColor(0x08, 0x37, 0x7C)
BLUE = RGBColor(0x3E, 0x95, 0xBE)
RED = RGBColor(0xC0, 0x00, 0x00)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)


def _box(slide, x, y, w, h, text, fill=WHITE, line=NAVY, txt=DARK, size=9):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.01)
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txt; r.font.name = _FONT


def _arrow(slide, x1, y1, x2, y2, label, lx, ly, color=NAVY, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(1.5)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dash:
        ln.insert(0, ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(1.8), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02); tf.margin_top = tf.margin_bottom = Inches(0)
    for i, s in enumerate(str(label).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = s
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = color; r.font.name = _FONT


def build_finance_diagram_slide(prs, fin: dict, *, business_name="",
                                section_label="02  금융개요", subtitle="금융구조도"):
    """금융구조도 페이지를 prs에 추가. fin = {차주, 신탁, 본건{금액,LTV,금리,담보}, 브릿지, 연대보증}."""
    from modules.page_builders import clone_slide_layout
    from modules.frame_builders import _fill_header, _add_combined_footer

    fin = fin or {}
    bg = fin.get("본건") or {}
    slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)
    _fill_header(slide, section_label, subtitle,
                 "본건 담보대출을 기초자산으로 신규 유동화 SPC가 사모사채를 발행하는 구조입니다.")

    차주 = fin.get("차주") or "차주"
    신탁 = fin.get("신탁") or "부동산담보신탁\n(신탁사)"
    연대 = fin.get("연대보증") or ""
    브릿지 = fin.get("브릿지") or ""
    담보 = bg.get("담보") or ""
    본건줄 = "본건 담보대출  " + (bg.get("금액") or "")
    조건줄 = " · ".join([x for x in [
        ("LTV " + bg["LTV"]) if bg.get("LTV") else "",
        ("금리 " + bg["금리"]) if bg.get("금리") else ""] if x])

    _box(slide, 0.55, 1.85, 1.9, 0.6, str(신탁), fill=LGRAY, size=9)
    _box(slide, 0.55, 3.4, 1.9, 0.75, "차주\n" + str(차주), fill=WHITE, size=9)
    if 연대:
        _box(slide, 0.55, 5.45, 1.9, 0.55, str(연대), fill=LGRAY, size=9)

    _box(slide, 3.7, 3.3, 2.7, 0.95,
         "\n".join([x for x in [본건줄, ("(" + 조건줄 + ")") if 조건줄 else "", 담보] if x]),
         fill=RGBColor(0xFD, 0xE7, 0xE7), line=RED, txt=RED, size=9)

    _box(slide, 8.0, 2.6, 2.3, 0.7, "신규 유동화\nSPC (TBD)", fill=NAVY, txt=WHITE, size=10.5)
    _box(slide, 8.0, 4.45, 2.3, 0.65, "투자자\n(사채권자)", fill=BLUE, txt=WHITE, size=10.5)

    if 브릿지:
        _box(slide, 3.7, 5.5, 3.2, 0.75,
             "브릿지 대출  " + str(브릿지) + "\n→ 본건 대출 상환재원", fill=LGRAY, size=9, txt=DARK)

    _arrow(slide, 1.5, 3.4, 1.5, 2.45, "담보신탁계약", 1.62, 2.75)
    _arrow(slide, 2.45, 3.75, 3.7, 3.75, "담보대출", 2.55, 3.45)
    if 연대:
        _arrow(slide, 1.5, 5.45, 1.5, 4.15, "연대보증", 1.62, 4.7)
    _arrow(slide, 6.4, 3.6, 8.0, 3.0, "대출채권\n(기초자산)", 6.35, 2.8)
    _arrow(slide, 9.15, 3.3, 9.15, 4.45, "사모사채\n발행·인수", 9.3, 3.65)
    if 브릿지:
        _arrow(slide, 5.3, 5.5, 5.3, 4.25, "상환", 5.42, 4.75, dash=True)

    try:
        _add_combined_footer(slide, business_name)
    except Exception:
        pass
    return slide


# ── 금융구조도용 데이터 추출(LLM) ────────────────────
_FIN_SYS = """부동산 금융 IM 원문에서 '금융구조도'에 쓸 항목만 뽑아 JSON만 출력한다.
값은 원문 그대로. 없으면 null. 대괄호 표기([30]억원 등)도 원문대로.
{"차주":..,"신탁":"부동산담보신탁 등",
 "본건":{"금액":"[30]억원","LTV":"43.98%","금리":"10%","담보":"주택 9개 호실 1순위 담보신탁"},
 "브릿지":"Tr.A 1,450억 / Tr.B 200억 (있으면)","연대보증":"조합장 연대보증 등(있으면)"}
JSON 외 금지."""


def extract_finance_diagram_data(full_text: str) -> dict:
    from modules.claude_api import call_claude
    res = call_claude(_FIN_SYS, f"[IM 원문]\n{full_text}", slide_num=802,
                      pdf_context=full_text, prompt_version="fin_diagram_v1")
    return res.get("data") or {} if res.get("ok") else {}
