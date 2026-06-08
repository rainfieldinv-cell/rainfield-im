"""
ai_slide_builders.py
─────────────────────────────────────────────────────────
Claude API 기반 슬라이드 데이터 생성 + PPT 빌더.

이 파일은 modules/page_builders.py 를 수정하지 않고
AI 처리가 필요한 슬라이드만 별도로 처리합니다.

현재 구현된 슬라이드 (그룹 A):
  - 슬라이드 2: Executive Summary
  - 슬라이드 5: 1.1 본건 사모사채 개요
  - 슬라이드 7: 2.1 투자구조도
─────────────────────────────────────────────────────────
"""

import copy

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn as _qn
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from modules.claude_api import call_claude, verify_numbers_in_pdf
from modules.page_builders import (
    clone_slide_layout,
    _replace_text_frame_content as _replace_tf_content,
    _replace_footer_business_name as _replace_footer,
)

# 슬라이드 크기 (cm)
_SLIDE_W = 27.517
_SLIDE_H = 19.05

# Rainfield 다크 네이비 + 포인트 그린
_C_DARK   = RGBColor(0x00, 0x20, 0x60)   # 002060
_C_GREEN  = RGBColor(0x70, 0xAD, 0x47)   # 70AD47
_C_GRAY   = RGBColor(0x7F, 0x7F, 0x7F)   # 7F7F7F
_C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_C_LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)   # F2F2F2

# ════════════════════════════════════════════════════════
# 신도림 팔레트 (전체 PPT 공통 — 앞으로 모든 슬라이드에서 이 색만 사용)
# ════════════════════════════════════════════════════════
PALETTE = {
    "navy_dark": RGBColor(0x08, 0x37, 0x7C),  # 08377C 진남색 (표 헤더 주력 / 차주)
    "blue":      RGBColor(0x00, 0x63, 0xA1),  # 0063A1 파랑 (신탁사·시공사·AMC)
    "steel":     RGBColor(0x3E, 0x95, 0xBE),  # 3E95BE 스틸블루 (강조, 보통 투명 35%)
    "maroon":    RGBColor(0x8C, 0x4A, 0x59),  # 8C4A59 자주 (대주)
    "gray":      RGBColor(0xD9, 0xD9, 0xD9),  # D9D9D9 회색
    "label_gray": RGBColor(0xF2, 0xF2, 0xF2),  # F2F2F2 구분열 (흰색 배경1, 5% 더 어둡게)
    "red":       RGBColor(0xC0, 0x00, 0x00),  # C00000 빨강 강조
    "gray_text": RGBColor(0x80, 0x80, 0x80),  # 808080 푸터/보조 회색
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "black":     RGBColor(0x00, 0x00, 0x00),
}

# ════════════════════════════════════════════════════════
# 슬라이드 2: Executive Summary
# ════════════════════════════════════════════════════════

SLIDE_2_SYSTEM_PROMPT = """당신은 부동산 PF 투자제안서 작성 전문가입니다.
주어진 PDF 원문 전체를 읽고, Executive Summary 슬라이드에 들어갈 핵심 메시지 5개를 추출합니다.

[엄격한 규칙]
1. 모든 숫자(금액/비율/날짜/평수/세대수)는 PDF 원문에 있는 값만 사용하라. 절대 새로 만들거나 추정하지 마라.
2. 출력은 반드시 아래 JSON 스키마를 따라야 한다. 키를 추가하거나 변경하지 마라.
3. key_points 는 정확히 5개여야 한다.
4. 각 항목의 title 은 12자 이내, description 은 80자 이내로 간결하게.
5. 한국어로 작성하되, 투자자가 한눈에 핵심을 파악할 수 있도록 명확하게.

[5개 항목 고정 순서 — title 은 이 순서대로]
1. "낮은 인허가 리스크"
2. "시공사 리스크"
3. "낮은 토지확보 리스크"
4. "높은 분양성"
5. "만기"

[JSON 출력 스키마]
{
  "deal_title": "본 건 거래 한 줄 요약 (예: 천안 부성2지구 도시개발사업 토지담보대출 사모사채)",
  "deal_summary": "본 건 개요 2~3 문장. 차주, 사업지, 금액, 만기 포함.",
  "key_points": [
    {"title": "낮은 인허가 리스크", "description": "PDF 원문 근거 포함 80자 이내"},
    {"title": "시공사 리스크", "description": "PDF 원문 근거 포함 80자 이내"},
    {"title": "낮은 토지확보 리스크", "description": "PDF 원문 근거 포함 80자 이내"},
    {"title": "높은 분양성", "description": "PDF 원문 근거 포함 80자 이내"},
    {"title": "만기", "description": "PDF 원문 근거 포함 80자 이내"}
  ]
}"""

SLIDE_2_USER_TEMPLATE = """[PDF 원문]
{pdf_text}

위 PDF 를 분석해 Executive Summary JSON 을 출력하라.
key_points 는 반드시 5개, title 은 지정된 순서와 동일하게."""


def generate_executive_summary(pdf_text: str) -> dict:
    """
    슬라이드 2 (Executive Summary) 데이터를 Claude 로 생성합니다.

    Returns
    -------
    call_claude() 반환값 dict
    {"ok": bool, "data": {...}, "usage": {...}, "cached": bool, ...}
    """
    print("=" * 60)
    print(f"[FORCE-DEBUG-AI] generate_executive_summary 호출됨 - PDF 길이={len(pdf_text)}")
    print("=" * 60)

    result = call_claude(
        system_prompt=SLIDE_2_SYSTEM_PROMPT,
        user_prompt=SLIDE_2_USER_TEMPLATE.format(pdf_text=pdf_text),
        slide_num=2,
        pdf_context=pdf_text,
        prompt_version="v1",
    )

    if result["ok"]:
        v = verify_numbers_in_pdf(result["data"], pdf_text)
        if not v["ok"]:
            print(f"[경고] 슬라이드 2 환각 의심 숫자: {v['hallucinated_numbers']}")
        else:
            print(f"[슬라이드 2] 숫자 검증 통과 ({v['verified_count']}개)")

    return result


# ════════════════════════════════════════════════════════
# 슬라이드 5: 1.1 본건 사모사채 개요
# ════════════════════════════════════════════════════════

SLIDE_5_SYSTEM_PROMPT = """당신은 사모사채 발행 전문가입니다.
주어진 PDF 의 Bridge Loan 정보를 바탕으로, 그 대출의 Tr.B 를 기초자산으로 발행하는 사모사채의 조건을 정리합니다.

[엄격한 규칙]
1. 차주, 사업명, 금액, 만기 등 PDF 에 있는 값은 그대로 사용하라.
2. PDF 에 없는 항목(채권 금리, 인수수수료 등)은 반드시 "[ TBD ]" 로 표시하라.
3. 발행금액은 PDF 의 Tr.B 금액과 동일하게.
4. 발행일은 PDF 의 대출 인출 예정일(또는 약정 예정일)과 동일.
5. 만기일은 발행일 기준 6개월 후.
6. fields 배열은 정확히 9개. 순서와 label 은 스키마와 동일하게.
7. 출력은 JSON 만. 다른 텍스트 금지.

[JSON 출력 스키마]
{
  "intro_paragraph": "본 건 사모사채 소개 1~2문장 (차주명, 사업명, Tr.B 금액, 만기 포함)",
  "fields": [
    {"label": "사모사채명", "value": "TBD(신규 유동화 SPC) 제1회 무기명식 무보증 사모사채"},
    {"label": "사채 유형", "value": "국내 전자등록 또는 실물 발행, 무기명식 무보증 사모사채"},
    {"label": "발행인", "value": "TBD(신규 유동화 SPC)"},
    {"label": "기초자산", "value": "[사업명] 토지담보대출 Tr.B 대출채권 — PDF 값 사용"},
    {"label": "발행금액", "value": "[Tr.B 금액] — PDF 값 사용"},
    {"label": "발행일", "value": "[YYYY년 MM월] (예정) — PDF 값 사용"},
    {"label": "만기일", "value": "[YYYY년 MM월 00일] (발행일로부터 약 6개월) — PDF 값 사용"},
    {"label": "금융조건", "value": "All-in [ TBD ]%, 채권 금리: 연 [ TBD ]%(고정, 세전), 인수수수료: 발행금액의 [ TBD ]%"},
    {"label": "이자지급주기", "value": "3개월 단위 후취 또는 만기 전액 후취"}
  ]
}"""

SLIDE_5_USER_TEMPLATE = """[PDF 원문]
{pdf_text}

위 PDF 에서 Tr.B 금액, 발행일, 만기일, 차주명, 사업명을 찾아 사모사채 개요 JSON 을 출력하라.
PDF 에 없는 항목은 반드시 [ TBD ] 로 표시하고, fields 는 정확히 9개."""


def generate_sasae_overview(pdf_text: str) -> dict:
    """
    슬라이드 5 (1.1 본건 사모사채 개요) 데이터를 Claude 로 생성합니다.

    Returns
    -------
    call_claude() 반환값 dict
    """
    print("=" * 60)
    print(f"[FORCE-DEBUG-AI] generate_sasae_overview 호출됨 - PDF 길이={len(pdf_text)}")
    print("=" * 60)

    result = call_claude(
        system_prompt=SLIDE_5_SYSTEM_PROMPT,
        user_prompt=SLIDE_5_USER_TEMPLATE.format(pdf_text=pdf_text),
        slide_num=5,
        pdf_context=pdf_text,
        prompt_version="v1",
    )

    if result["ok"]:
        v = verify_numbers_in_pdf(result["data"], pdf_text)
        if not v["ok"]:
            print(f"[경고] 슬라이드 5 환각 의심 숫자: {v['hallucinated_numbers']}")
        else:
            print(f"[슬라이드 5] 숫자 검증 통과 ({v['verified_count']}개)")

    return result


# ════════════════════════════════════════════════════════
# 슬라이드 7: 2.1 투자구조도
# ════════════════════════════════════════════════════════

SLIDE_7_SYSTEM_PROMPT = """당신은 부동산 금융 투자구조 분석 전문가입니다.
PDF 에서 투자 참여기관을 추출하여 투자구조도용 데이터를 만듭니다.

[엄격한 규칙]
1. 차주, 시행사, 시공사, 신탁사, 대주, 사채권자는 PDF 에 있는 이름 그대로 사용.
2. 금액(담보대출 총액, Tr.A, Tr.B)은 PDF 원문 그대로 사용. 추정 금지.
3. SPC 명은 "신규 유동화 SPC T.B.D." 로 표시.
4. 사채권자는 "사채권자 T.B.D." 로 표시.
5. relationships 는 PDF 에서 파악 가능한 계약관계만 기재.
6. 출력은 JSON 만.

[JSON 출력 스키마]
{
  "intro_paragraph": "투자구조 소개 1~2문장 (총 대출금액, Tr.A/B 구성 포함)",
  "total_loan_amount": "PDF 원문 값",
  "tranches": [
    {"name": "Tranche A", "amount": "PDF 원문 값"},
    {"name": "Tranche B", "amount": "PDF 원문 값"}
  ],
  "entities": {
    "borrower": "차주명 — PDF 값",
    "constructor": "시공사명 — PDF 값",
    "trustee": "신탁사명 — PDF 값",
    "tr_a_lenders": "Tr.A 대주 (PDF 에 없으면 'T.B.D.')",
    "spc": "신규 유동화 SPC T.B.D.",
    "bondholders": "사채권자 T.B.D."
  },
  "relationships": [
    "차주 - 시공사: 공사도급계약",
    "차주 - 대주: 대출약정",
    "차주 - 신탁사: 담보신탁계약",
    "SPC - 사채권자: 사모사채 발행/인수"
  ]
}"""

SLIDE_7_USER_TEMPLATE = """[PDF 원문]
{pdf_text}

위 PDF 에서 투자 참여기관과 금액을 찾아 투자구조도 JSON 을 출력하라.
PDF 에 없는 기관명은 T.B.D. 로 표시하고, 금액은 절대 추정하지 마라."""


def generate_investment_structure(pdf_text: str) -> dict:
    """
    슬라이드 7 (2.1 투자구조도) 데이터를 Claude 로 생성합니다.

    Returns
    -------
    call_claude() 반환값 dict
    """
    print("=" * 60)
    print(f"[FORCE-DEBUG-AI] generate_investment_structure 호출됨 - PDF 길이={len(pdf_text)}")
    print("=" * 60)

    result = call_claude(
        system_prompt=SLIDE_7_SYSTEM_PROMPT,
        user_prompt=SLIDE_7_USER_TEMPLATE.format(pdf_text=pdf_text),
        slide_num=7,
        pdf_context=pdf_text,
        prompt_version="v1",
    )

    if result["ok"]:
        v = verify_numbers_in_pdf(result["data"], pdf_text)
        if not v["ok"]:
            print(f"[경고] 슬라이드 7 환각 의심 숫자: {v['hallucinated_numbers']}")
        else:
            print(f"[슬라이드 7] 숫자 검증 통과 ({v['verified_count']}개)")

    return result


# ════════════════════════════════════════════════════════
# 내부 헬퍼
# ════════════════════════════════════════════════════════

def _replace_text_keep_runs(tf, new_text: str):
    """기존 단락/런의 XML 구조(폰트·크기·색상)를 보존하면서 텍스트만 교체합니다.

    _replace_text_frame_content 와 달리 모든 <a:p>/<a:r>/<a:rPr> 를 제거하지 않고
    기존 단락을 재사용하므로 템플릿에서 지정한 폰트가 그대로 유지됩니다.
    줄 수가 기존 단락보다 많으면 마지막 단락 스타일을 복사해 추가합니다.
    """
    lines = (new_text or '').split('\n')
    txBody = tf._txBody
    all_paras = txBody.findall(_qn('a:p'))

    if not all_paras:
        _replace_tf_content(tf, new_text)
        return

    last_p = all_paras[-1]

    for i, para in enumerate(all_paras):
        if i < len(lines):
            runs = para.findall(_qn('a:r'))
            if runs:
                t_elem = runs[0].find(_qn('a:t'))
                if t_elem is not None:
                    t_elem.text = lines[i]
                else:
                    t_elem = etree.SubElement(runs[0], _qn('a:t'))
                    t_elem.text = lines[i]
                for r in runs[1:]:
                    para.remove(r)
            else:
                new_r = etree.SubElement(para, _qn('a:r'))
                last_runs = last_p.findall(_qn('a:r'))
                if last_runs:
                    old_rPr = last_runs[0].find(_qn('a:rPr'))
                    if old_rPr is not None:
                        new_r.insert(0, copy.deepcopy(old_rPr))
                t_elem = etree.SubElement(new_r, _qn('a:t'))
                t_elem.text = lines[i]
        else:
            runs = para.findall(_qn('a:r'))
            if runs:
                t_elem = runs[0].find(_qn('a:t'))
                if t_elem is not None:
                    t_elem.text = ''
                for r in runs[1:]:
                    para.remove(r)

    for i in range(len(all_paras), len(lines)):
        new_p = copy.deepcopy(last_p)
        new_runs = new_p.findall(_qn('a:r'))
        if new_runs:
            t_elem = new_runs[0].find(_qn('a:t'))
            if t_elem is not None:
                t_elem.text = lines[i]
            for r in new_runs[1:]:
                new_p.remove(r)
        else:
            new_r = etree.SubElement(new_p, _qn('a:r'))
            t_elem = etree.SubElement(new_r, _qn('a:t'))
            t_elem.text = lines[i]
        txBody.append(new_p)


def _find_shape_by_pos(slide, left_cm: float, top_cm: float, tol_cm: float = 0.35):
    """좌표(cm)로 슬라이드에서 텍스트 프레임 도형을 찾습니다."""
    tol = int(tol_cm * 360000)
    lx  = int(left_cm * 360000)
    tx  = int(top_cm  * 360000)
    best, best_dist = None, float('inf')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        dl = abs(shape.left - lx)
        dt = abs(shape.top  - tx)
        if dl <= tol and dt <= tol:
            d = dl + dt
            if d < best_dist:
                best_dist, best = d, shape
    return best


def _find_table_by_pos(slide, left_cm: float, top_cm: float, tol_cm: float = 0.35):
    """좌표(cm)로 슬라이드에서 TABLE shape 를 찾습니다."""
    tol = int(tol_cm * 360000)
    lx  = int(left_cm * 360000)
    tx  = int(top_cm  * 360000)
    best, best_dist = None, float('inf')
    for shape in slide.shapes:
        if shape.shape_type != 19:   # 19 = TABLE
            continue
        dl = abs(shape.left - lx)
        dt = abs(shape.top  - tx)
        if dl <= tol and dt <= tol:
            d = dl + dt
            if d < best_dist:
                best_dist, best = d, shape
    return best


def _fill_group_label(group_shape, text: str):
    """GROUP 내부 첫 번째 TEXT_BOX 의 텍스트를 교체합니다 (섹션 헤더 레이블)."""
    try:
        for child in group_shape.shapes:
            if child.has_text_frame:
                _replace_text_keep_runs(child.text_frame, text)
                return
    except Exception:
        pass


def _remove_non_pdf_shapes(slide):
    """PDF 원본 구조도에 없는 shape 를 XML 에서 완전히 제거합니다.

    삭제 대상 (좌표 기반 식별):
      - shape[10] TEXT_BOX  (L=7.41, T=10.72)  원래 "사업 시행"
      - shape[11] TABLE 2r  (L=13.18, T=14.26) 원래 자산관리자
      - shape[13] LINE      (L=14.44, T=11.77) 자산관리자→차주 세로 화살표
      - shape[14] TEXT_BOX  (L=13.87, T=12.65) 원래 "자산관리"
      - shape[17] TABLE 7r  (L=18.87, T=10.96) 원래 투자자 목록
      - shape[18] LINE      (L=15.70, T=11.07) 투자자→차주 가로 화살표
      - shape[22] TEXT_BOX  (L=15.52, T=11.07) 원래 "Equity 825억원"
      - shape[25] TEXT_BOX  (L=8.26, T=12.74)  원래 "책임준공 확약"
      - shape[27] PICTURE   (shape_type=13)     사업지 이미지
      - shape[29] TEXT_BOX  (L=0.00, T=18.35)  사업명 푸터
    """
    # 좌표로 삭제할 shape 를 식별
    # (left_cm, top_cm, tolerance_cm, shape_type_filter)
    # shape_type_filter: None=모든 타입, 19=TABLE, 17=TEXT_BOX, 13=PICTURE, 9=LINE
    delete_specs = [
        (7.41,  10.72, 0.35, 17),   # [10] "사업 시행"
        (13.18, 14.26, 0.35, 19),   # [11] 자산관리자 TABLE
        (14.44, 11.77, 0.35, 9),    # [13] 자산관리자 세로 LINE
        (13.87, 12.65, 0.35, 17),   # [14] "자산관리"
        (18.87, 10.96, 0.35, 19),   # [17] 투자자 TABLE
        (15.70, 11.07, 0.20, 9),    # [18] 투자자 가로 LINE
        (15.52, 11.07, 0.35, 17),   # [22] "Equity"
        (8.26,  12.74, 0.35, 17),   # [25] "책임준공 확약"
        (0.00,  18.35, 0.20, 17),   # [29] 사업명 푸터
    ]

    to_remove = []
    for shape in slide.shapes:
        l_cm = shape.left / 360000 if shape.left else 0
        t_cm = shape.top  / 360000 if shape.top  else 0
        st   = shape.shape_type

        # PICTURE (사업지 이미지) — shape_type == 13
        if st == 13:
            to_remove.append(shape)
            continue

        for spec_l, spec_t, tol, spec_type in delete_specs:
            if spec_type is not None and st != spec_type:
                continue
            if abs(l_cm - spec_l) <= tol and abs(t_cm - spec_t) <= tol:
                to_remove.append(shape)
                break

    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    if to_remove:
        print(f"[build_slide_7] PDF에 없는 shape {len(to_remove)}개 삭제 완료")


def _fill_sasae_table(slide, fields: list):
    """슬라이드 6 템플릿의 9r×2c TABLE 에 value 값을 씁니다.
    레이블 열(col 0)은 그대로 유지하고, 값 열(col 1) 만 교체합니다."""
    sh = _find_table_by_pos(slide, left_cm=1.09, top_cm=3.84, tol_cm=0.40)
    if sh is None:
        print("[경고] _fill_sasae_table: 9r×2c TABLE 못 찾음")
        return
    tbl = sh.table
    for i, f in enumerate(fields):
        if i >= len(tbl.rows):
            break
        _replace_text_keep_runs(tbl.cell(i, 1).text_frame, f.get("value", ""))


def _fill_entity_tables(slide, entities: dict, tranches: list, total_loan: str):
    """슬라이드 8 투자구조도의 entity TABLE 셀에 실제 기관명/금액을 씁니다.

    PDF 원본 구조도 기준:
      - 신탁사·차주·시공사 3개 entity 박스(2r×1c) 기관명 교체
      - 대주 박스(4r×1c): "Tr.A" / "Tr.B" 만 표기 (금액 제거)
      - 자산관리자·투자자 박스는 _remove_non_pdf_shapes() 에서 삭제
    """
    # ── 3개 entity 박스 (2r×1c): row[0]=레이블 유지, row[1]=기관명 교체 ──
    entity_map = [
        (10.47, 4.98,  entities.get("trustee",     "")),   # 신탁사
        (10.48, 9.55,  entities.get("borrower",    "")),   # 차주
        (10.20, 14.26, entities.get("constructor", "")),   # 시공사
    ]
    for l, t, val in entity_map:
        sh = _find_table_by_pos(slide, l, t)
        if sh:
            _replace_text_keep_runs(sh.table.cell(1, 0).text_frame, val)
        else:
            print(f"[경고] _fill_entity_tables: entity TABLE 못 찾음 ({l}, {t})")

    # ── 대주 박스 (shape[09], 4r×1c, L=18.87 T=5.44) — 금액 제거, Tr 라벨만 ──
    loan_sh = _find_table_by_pos(slide, 18.87, 5.44)
    if loan_sh:
        tbl = loan_sh.table
        # row[0]: 헤더 → 비움 (Bridge Loan 총액 제거)
        _replace_text_keep_runs(tbl.cell(0, 0).text_frame, "")
        # row[1..]: tranche 이름만 (금액 없이)
        for i, tr in enumerate(tranches):
            row_idx = i + 1
            if row_idx < len(tbl.rows):
                _replace_text_keep_runs(
                    tbl.cell(row_idx, 0).text_frame,
                    tr.get('name', '')
                )
        for i in range(len(tranches) + 1, len(tbl.rows)):
            _replace_text_keep_runs(tbl.cell(i, 0).text_frame, "")
    else:
        print("[경고] _fill_entity_tables: 대주 TABLE 못 찾음")


# ════════════════════════════════════════════════════════
# 슬라이드 2 빌더: Executive Summary
# ════════════════════════════════════════════════════════

def _relayout_exec_summary(slide):
    """슬라이드 1 Executive Summary 동적 섹션 레이아웃 (이름 기반).

    - 육각형(<>) 좌우폭/left 통일 (바깥/안쪽 각각 동일값, 세로만 조절)
    - 본문 실제 텍스트의 줄 수를 추정 → 본문 높이 산출 →
      꺾쇠 H = 본문높이 + 상하 0.15" 여백 (본문은 꺾쇠 안 세로중앙 정렬)
    - [제목(+부제)+꺾쇠] 블록 단위로 세 섹션을 균등 간격으로 세로 중앙 분포
    - 바깥/안쪽 육각형의 상하 offset(델타)은 원본 비율 보존
    좌표 단위: 인치(EMU = 인치×914400). 모든 대상은 고유 name 으로 매칭.
    """
    IN = 914400
    by = {sh.name: sh for sh in slide.shapes}

    O_L, O_W = int(1.30 * IN), int(8.24 * IN)   # 바깥 육각형 좌/폭
    I_L, I_W = int(1.45 * IN), int(7.94 * IN)   # 안쪽 육각형 좌/폭
    MARGIN = 0.15            # 꺾쇠 내부 상하 여백(인치)
    TITLE_H = 0.40           # 제목 GROUP 높이
    SUB_BLOCK = 0.71         # 섹션1: 제목+부제+여백 (제목top→꺾쇠top)
    TITLE_GAP = 0.45         # 섹션2·3: 제목top→꺾쇠top
    TOP, BOT = 0.84, 7.00    # 사용 가능 세로 범위(로고 아래 ~ 푸터 위)

    def _strip_empty_paras(tb):
        """본문 텍스트프레임의 빈 단락(<a:p> 텍스트 없음)을 제거. 최소 1개는 유지.
        (_replace_text_keep_runs 가 템플릿의 잔여 빈 단락을 남겨 줄 수를 부풀리는 문제 해결)"""
        txbody = tb.text_frame._txBody
        ps = txbody.findall(_qn('a:p'))
        nonempty = [p for p in ps
                    if "".join(t.text or "" for t in p.findall('.//' + _qn('a:t'))).strip()]
        if not nonempty:
            return
        for p in ps:
            txt = "".join(t.text or "" for t in p.findall('.//' + _qn('a:t'))).strip()
            if not txt:
                txbody.remove(p)

    def _est_body(tb):
        """본문 줄 수·행높이·예상높이(인치) 추정 (빈 단락 제거 후)."""
        tf = tb.text_frame
        width_in = (tb.width or 0) / IN
        fpt = 11.0
        done = False
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size:
                    fpt = r.font.size.pt
                    done = True
                    break
            if done:
                break
        # 한글은 대체로 전각(글자폭≈fontpt) → 안전계수 없이 fontpt 기준
        cpl = max(1, int(width_in * 72.0 / fpt))
        lines = 0
        for p in tf.paragraphs:
            s = "".join(r.text for r in p.runs).strip()
            if s:
                lines += (len(s) + cpl - 1) // cpl   # ceil
        lines = max(1, lines)
        line_h = fpt * 1.25 / 72.0
        return lines, round(line_h, 3), round(lines * line_h, 3)

    # (제목GROUP, 부제, 바깥육각형, 안쪽육각형, 본문, has_sub)
    secs = [
        # 섹션1 부제(TextBox 55) 삭제 → ns=None, has_sub=False (간격 압축 = 높이 넘침 방지)
        ("그룹 1",  None,         "육각형 5",  "육각형 17", "TextBox 42", False),
        ("그룹 2",  None,         "육각형 24", "육각형 25", "TextBox 26", False),
        ("그룹 51", None,         "육각형 48", "육각형 49", "TextBox 50", False),
    ]

    # 1) 본문 빈 단락 제거 후 추정 → 꺾쇠 H, 블록 H
    est, hexH, blockH = [], [], []
    for (nt, ns, nho, nhi, nb, hassub) in secs:
        tb = by.get(nb)
        if tb is not None:
            _strip_empty_paras(tb)
        e = _est_body(tb) if tb is not None else (1, 0.2, 0.2)
        est.append(e)
        h = e[2] + 2 * MARGIN
        hexH.append(h)
        blockH.append((SUB_BLOCK if hassub else TITLE_GAP) + h)

    # 2) 균등 간격 산출 (세로 중앙 분포)
    content = sum(blockH)
    gap = max(0.30, (BOT - TOP - content) / 2)

    # 3) 배치
    report = []
    cur = TOP
    for i, (nt, ns, nho, nhi, nb, hassub) in enumerate(secs):
        ttop = cur
        if hassub:
            stop = ttop + TITLE_H
            hexT = ttop + SUB_BLOCK
        else:
            stop = None
            hexT = ttop + TITLE_GAP
        if by.get(nt) is not None:
            by[nt].top = int(ttop * IN)
        if ns and by.get(ns) is not None:
            by[ns].top = int(stop * IN)
        ho, hi = by.get(nho), by.get(nhi)
        if ho is not None and hi is not None:
            dT = (hi.top or 0) - (ho.top or 0)
            dH = (hi.height or 0) - (ho.height or 0)
            before_h = round((ho.height or 0) / IN, 2)
            ho.left, ho.width, ho.top, ho.height = O_L, O_W, int(hexT * IN), int(hexH[i] * IN)
            hi.left, hi.width = I_L, I_W
            hi.top, hi.height = ho.top + dT, ho.height + dH
        b = by.get(nb)
        if b is not None:
            b.top = int((hexT + MARGIN) * IN)
            b.height = int((hexH[i] - 2 * MARGIN) * IN)
            try:
                b.text_frame.word_wrap = True
                b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
        report.append({
            "sec": nt, "title_top": round(ttop, 2), "hex_top": round(hexT, 2),
            "lines": est[i][0], "body_h": est[i][2],
            "hexH_before": before_h if ho is not None else None,
            "hexH_after": round(hexH[i], 2),
        })
        cur = hexT + hexH[i] + gap

    print(f"[build_slide_2] 동적 레이아웃 gap={round(gap, 2)}\":")
    for r in report:
        print(f"   {r}")
    return report


def style_table(sh, *, nested=False, has_header=False, label_cols=(0,),
                header_fill=None, value_fill=None, label_fill=None):
    """표 셀 서식 표준 적용 (재사용 헬퍼).

    - 여백: 기본 0 / 좌측정렬→marL=0.3cm / 우측정렬→marR=0.3cm / 가운데→전부 0
    - 글씨: 일반표 10.5pt, 중첩표(nested=True) 10pt
    - Bold/Light: 헤더행(has_header & row0)·구분열(label_cols)=피플폰트 Bold,
                  세부내용=피플폰트 Light  (font.bold 속성은 쓰지 않고 폰트명으로 굵기 표현)
    - 색: header_fill(헤더행)·value_fill(값셀) 지정 시 적용, 미지정이면 기존 유지
          (헤더 08377C, 구분/값 회색 D9D9D9, 강조 3E95BE 등은 호출부에서 지정)
    - 정렬: 가로 기본 가운데(미지정 시), 세로 항상 가운데(MIDDLE)
    - 테두리: 모든 셀 4변 전체, 0.5pt(6350 EMU), A5A5A5(회색 강조3)
    - 행높이: has_header 면 헤더행 0.6cm(중첩 0.55cm)
    """
    tbl = sh.table
    # 자동 밴딩/특수행 스타일 제거 — 값 셀에 의도치 않은 음영(연파랑 줄무늬)이 깔리는 것 방지.
    #   셀 색은 아래에서 직접 지정(헤더 네이비/구분 F2F2F2/값 투명)하므로 자동 스타일은 끈다.
    for attr in ("first_row", "last_row", "first_col", "last_col", "horz_banding", "vert_banding"):
        try:
            setattr(tbl, attr, False)
        except Exception:
            pass
    fpt = 10 if nested else 10.5
    M = Cm(0.3)
    Z = Inches(0)
    BORDER_HEX = "A5A5A5"          # 회색, 강조3 (테마 accent3) — 레이아웃 표 실측값
    BORDER_W = 6350               # 0.5pt (EMU)

    def _set_borders(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for el in tcPr.findall(_qn(tag)):
                tcPr.remove(el)
        for i, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
            ln = tcPr.makeelement(_qn(tag),
                                  {"w": str(BORDER_W), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
            sf = ln.makeelement(_qn("a:solidFill"), {})
            sf.append(sf.makeelement(_qn("a:srgbClr"), {"val": BORDER_HEX}))
            ln.append(sf)
            ln.append(ln.makeelement(_qn("a:prstDash"), {"val": "solid"}))
            tcPr.insert(i, ln)          # 테두리는 tcPr 최상단(스키마 순서)

    for ri in range(len(tbl.rows)):
        for ci in range(len(tbl.columns)):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            # 가로 정렬: 미지정이면 가운데로 통일
            for p in tf.paragraphs:
                if p.alignment is None:
                    p.alignment = PP_ALIGN.CENTER
            algn = str(tf.paragraphs[0].alignment or "")
            cell.margin_left = M if algn.startswith("LEFT") else Z
            cell.margin_right = M if algn.startswith("RIGHT") else Z
            cell.margin_top = Z
            cell.margin_bottom = Z
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE       # 세로 항상 가운데
            is_header_cell = has_header and ri == 0
            is_bold = is_header_cell or (ci in label_cols)
            fname = "피플폰트 Bold" if is_bold else "피플폰트 Light"
            # 헤더(색칠된 셀) 글자=흰색, 그 외(구분·내용)=검정
            fcolor = _C_WHITE if is_header_cell else RGBColor(0, 0, 0)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fpt)
                    r.font.name = fname
                    r.font.color.rgb = fcolor
            if has_header and ri == 0 and header_fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif ci in label_cols and label_fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = label_fill
            elif value_fill is not None and ci not in label_cols:
                cell.fill.solid()
                cell.fill.fore_color.rgb = value_fill
            else:
                # 색 미지정 셀(내용/값 셀) = 흰색으로 확실히 칠함 (표 스타일 음영 비침 방지)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            _set_borders(cell)                              # 전체 테두리 A5A5A5 0.5pt
    if has_header and len(tbl.rows) > 0:
        tbl.rows[0].height = Cm(0.55 if nested else 0.6)


def clean_bullet(tf, marL=171450, indent=-171450):
    """텍스트프레임 모든 단락의 자동 글머리/레벨/들여쓰기를 정리한다.

    - a:buChar / a:buAutoNum / a:buFont 제거 후 a:buNone 삽입 (자동 글머리 제거)
    - lvl 속성 제거 → 모든 단락 level 0 통일
    - marL/indent 를 동일 값으로 통일 (기본: marL=171450, indent=-171450)
    텍스트·런·폰트·색은 건드리지 않는다.
    """
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _q(tag):
        return "{%s}%s" % (A, tag)

    for p in tf.paragraphs:
        pPr = p._p.get_or_add_pPr()
        # 들여쓰기/레벨 통일
        pPr.set("marL", str(int(marL)))
        pPr.set("indent", str(int(indent)))
        if "lvl" in pPr.attrib:
            del pPr.attrib["lvl"]
        # 기존 글머리/글머리폰트/자동번호 제거
        for tag in ("buChar", "buAutoNum", "buFont", "buNone"):
            for el in pPr.findall(_q(tag)):
                pPr.remove(el)
        # buNone 삽입 (스키마상 buColor/buSzx/buFont 다음, 단 여기선 다 제거했으므로 끝에 append)
        pPr.append(pPr.makeelement(_q("buNone"), {}))


# 전 슬라이드 공통 푸터 사업명 (딜별 기본값 — business_name 인자로 덮어쓸 수 있음)
_FOOTER_BIZ = "천안 부성2지구 도시개발사업"
_FOOTER_GRAY = RGBColor(0x80, 0x80, 0x80)
_FONT_LIGHT = "피플폰트 Light"


def add_footer(slide, page_num, source=None, business_name=None):
    """전 슬라이드 공통 푸터 (형식 고정).

    - 오른쪽 아래: "{사업명}  |  {page_num}"  (피플폰트 Light 8pt, 808080, 우측정렬)
    - 왼쪽 아래(source 있을 때만): "출처: {source}"  (피플폰트 Light 9pt, 808080, 좌측정렬)
    - page_num 은 PPT 슬라이드 순서(slides 인덱스+1) 권장
    - 두 텍스트박스 모두 내부 여백 0
    기존 푸터성 텍스트박스(name 에 "슬라이드 번호" 포함)는 먼저 삭제 후 새로 생성.
    """
    biz = business_name or _FOOTER_BIZ

    # 1) 기존 푸터성(페이지번호/사업명) 텍스트박스 제거
    for sh in list(slide.shapes):
        if "슬라이드 번호" in (sh.name or ""):
            sh._element.getparent().remove(sh._element)

    def _mk(left, width, text, size_pt, align):
        tb = slide.shapes.add_textbox(Inches(left), Inches(7.16), Inches(width), Inches(0.26))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Inches(0)
        tf.margin_top = tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = _FONT_LIGHT
        run.font.size = Pt(size_pt)
        run.font.color.rgb = _FOOTER_GRAY
        return tb

    # 2) 오른쪽 푸터 (사업명 | 순번) — 우측 끝 ≈ 10.50"
    _mk(4.83, 5.67, f"{biz}  |  {page_num}", 8, PP_ALIGN.RIGHT)

    # 3) 왼쪽 출처 (source 지정 시에만)
    if source:
        _mk(0.33, 5.50, f"출처: {source}", 9, PP_ALIGN.LEFT)


def build_slide_2_executive_summary(prs, data: dict,
                                     business_name: str = "",
                                     page_num: int = 2):
    """
    슬라이드 2 (Executive Summary) PPT 슬라이드를 생성합니다.
    레이아웃.pptx index=1 템플릿을 복제해 텍스트만 교체합니다.

    Parameters
    ----------
    prs           : create_presentation_from_template() 로 만든 Presentation 객체
    data          : generate_executive_summary() 반환 data dict
    business_name : 하단 사업명 푸터
    page_num      : 슬라이드 번호
    """
    slide = clone_slide_layout(prs, "executive_summary")

    deal_title   = data.get("deal_title", "")
    deal_summary = data.get("deal_summary", "")
    key_points   = data.get("key_points", [])

    kp5 = key_points[4] if len(key_points) > 4 else {}

    # ── 본문/부제 텍스트 (★데이터 기반 — '지금 변환 중인 PDF' 의 Executive Summary 에서 생성) ──
    #   하드코딩 절대 금지. deal_summary / key_points 는 generate_executive_summary 가
    #   현재 PDF 의 Executive Summary 페이지만 보고 만든 값이다(천안 잔재 X).
    _SEC1_BODY = (deal_summary or "").strip()
    _pts = [p for p in key_points[:4] if (p.get("description") or "").strip()]
    # ★섹션2도 나머지 2개 섹션처럼 화살표(→) 스타일, '·'·':' 사용 금지(사용자 지시)
    _SEC2_BODY = "\n".join(
        f"→ {p.get('title','')} {p.get('description','')}".strip() for p in _pts
    )
    _SEC3_BODY = (kp5.get("description", "") or "").strip()

    def _s(l, t, text):
        sh = _find_shape_by_pos(slide, l, t)
        if sh:
            _replace_text_keep_runs(sh.text_frame, text)
        else:
            print(f"[경고] build_slide_2: TEXT_BOX 못 찾음 ({l:.2f}, {t:.2f})")

    # 텍스트 박스 교체 (좌표: 원본 템플릿 cm 기준 — relayout 전)
    _s(6.80, 3.14,  deal_title)      # TextBox 55 — 섹션1 부제(Bridge Loan 총 1,640억) 유지
    _s(4.49, 4.07,  _SEC1_BODY)      # TextBox 42 — 섹션1 본문(압축)
    _s(4.91, 8.22,  _SEC2_BODY)      # TextBox 26 — 섹션2 본문(4줄)
    _s(4.50, 16.31, _SEC3_BODY)      # TextBox 50 — 섹션3 본문(압축)
    # 부제 TextBox 56("핵심 투자 포인트")·TextBox 54("만기")는 아래에서 삭제하므로 채우지 않음

    # GROUP 제목 교체:
    #   섹션1 "Executive Summary"(고정 라벨) → 내용 기반 제목으로 교체
    #     (★요청: 항목 제목에 'Executive Summary' 단어 금지. 내용을 보고 제목을 쓴다.)
    #   섹션2 → "핵심 투자 포인트", 섹션3 → kp5 제목("만기")
    def _grp_text(grp):
        try:
            return " ".join(c.text_frame.text for c in grp.shapes
                            if c.has_text_frame).strip()
        except Exception:
            return ""

    for shape in slide.shapes:
        if shape.shape_type != 6:   # GROUP
            continue
        cur  = _grp_text(shape)
        t_cm = shape.top / 360000
        if "Executive" in cur:          # 섹션1 제목 (고정 'Executive Summary')
            _fill_group_label(shape, "거래 개요")
        elif abs(t_cm - 6.39) < 0.40:     # 그룹 2 (섹션2 제목)
            _fill_group_label(shape, "핵심 투자 포인트")
        elif abs(t_cm - 13.95) < 0.40:  # 그룹 51 (섹션3 제목)
            _fill_group_label(shape, kp5.get("title", "만기"))

    # ── 수정 2: 섹션 간격 압축 + 육각형 컨테이너 폭 통일 ──
    _relayout_exec_summary(slide)

    # ── 섹션2 본문(4줄) 글머리 정리: 자동 글머리 제거 + 레벨/들여쓰기 통일 ──
    #   (텍스트 앞의 "·" 가 유일한 글머리가 됨)
    s2_body = next((sh for sh in slide.shapes if sh.name == "TextBox 26"), None)
    if s2_body is not None and s2_body.has_text_frame:
        clean_bullet(s2_body.text_frame)

    # ── 사용자 요청: 하늘색 부제 '전부' 삭제 (제목 GROUP 3개는 모두 유지) ──
    #   섹션1 부제(TextBox 55)도 삭제 — "밑에 하늘색 글씨 빼"
    #   유지: 그룹 1/2/51(제목+체크마크), 본문, 육각형
    _DEL_NAMES = ["TextBox 55", "TextBox 56", "TextBox 54"]

    def _shape_text(sh):
        if sh.has_text_frame:
            return sh.text_frame.text.strip()
        try:
            return " ".join(c.text_frame.text for c in sh.shapes
                            if c.has_text_frame).strip()
        except Exception:
            return ""

    for nm in _DEL_NAMES:
        for sh in list(slide.shapes):
            if sh.name == nm:
                print(f"[build_slide_2] 삭제: name='{sh.name}' text='{_shape_text(sh)[:24]}'")
                sh._element.getparent().remove(sh._element)
                break

    # 공통 푸터 (사업명 | 순번)
    add_footer(slide, page_num, business_name=business_name or None)

    return slide


# ════════════════════════════════════════════════════════
# 슬라이드 5 빌더: 1.1 본건 사모사채 개요
# ════════════════════════════════════════════════════════

def build_slide_5_sasae_overview(prs, data: dict,
                                  business_name: str = "",
                                  page_num: int = 6):
    """
    슬라이드 5 (1.1 본건 사모사채 개요) PPT 슬라이드를 생성합니다.
    레이아웃.pptx index=5 템플릿의 9r×2c TABLE 을 그대로 활용합니다.

    Parameters
    ----------
    prs           : create_presentation_from_template() 로 만든 Presentation 객체
    data          : generate_sasae_overview() 반환 data dict
    business_name : 하단 사업명 푸터
    page_num      : 슬라이드 번호
    """
    # skip_graphic_frames=False → 9r×2c TABLE 포함 그대로 복제
    slide = clone_slide_layout(prs, "content", skip_graphic_frames=False)

    intro  = data.get("intro_paragraph", "")
    fields = data.get("fields", [])

    # 인트로 텍스트 (shape[03], T=2.29)
    intro_sh = _find_shape_by_pos(slide, 1.09, 2.29)
    if intro_sh:
        _replace_text_keep_runs(intro_sh.text_frame, intro)

    # 기존 9r×2c TABLE 의 값 열(col 1)만 교체
    _fill_sasae_table(slide, fields)

    # ── 표 서식 표준 적용 (일반표 10.5pt) ──
    #   구분열(col0)=Bold, 세부내용(col1)=Light, 가운데정렬셀 여백0·좌측정렬셀 marL=0.3cm
    #   (헤더 행 없는 레이블-값 표 → 색은 기존 무채움 유지)
    sasae_tbl = _find_table_by_pos(slide, 1.09, 3.84, tol_cm=0.40)
    if sasae_tbl is not None:
        style_table(sasae_tbl, nested=False, has_header=False, label_cols=(0,),
                    label_fill=PALETTE["label_gray"])   # 구분열(col0) 밝은회색 F2F2F2 (기초자산개요와 통일)

    # 공통 푸터 (사업명 | 순번)
    add_footer(slide, page_num, business_name=business_name or None)

    return slide


# ════════════════════════════════════════════════════════
# 슬라이드 7 빌더: 2.1 투자구조도
# ════════════════════════════════════════════════════════

def build_slide_7_investment_structure(prs, data: dict,
                                        business_name: str = "",
                                        page_num: int = 8):
    """
    슬라이드 7 (2.1 투자구조도) PPT 슬라이드를 생성합니다.

    [STEP 5-4-H 신규 방식]
    레이아웃 템플릿 다이어그램(박스/표/화살표/이미지)을 전부 버리고,
    헤더·푸터·배경·인트로만 남긴 "빈 베이스" 를 만든 뒤,
    PDF 4번 사진 좌표대로 박스·화살표를 코드로 직접 그립니다.

    STEP 1 (현재): 베이스만 구성 — 아래 6개 shape 만 유지하고 나머지 전부 삭제.
      유지:
        - [0]  "02 금융 개요"  상단 라벨        (L=1.09, T=1.11)
        - [3]  "2.1 투자구조도" 메인 타이틀      (L=1.09, T=1.64)
        - [16] 인트로 텍스트                    (L=1.09, T=2.29)
        - [1]  상단 구분선(배경)                (L=1.09, T=3.87)
        - [2]  하단 구분선(배경)                (L=1.09, T=18.16)
        - [28] 페이지번호(우하단)               (L=18.53, T=18.35)
      RAINFIELD 로고는 슬라이드 마스터/레이아웃 배경에 있어 자동 유지.

    STEP 2~3 (승인 후 예정): add_shape / add_connector 로 박스 4개 + 화살표 4개 직접 그림.

    Parameters
    ----------
    prs           : create_presentation_from_template() 로 만든 Presentation 객체
    data          : generate_investment_structure() 반환 data dict
    business_name : 하단 사업명 푸터 (이 슬라이드에서는 사용 안 함)
    page_num      : 슬라이드 번호
    """
    slide = clone_slide_layout(prs, "investment_structure")

    intro = data.get("intro_paragraph", "")

    # ── 1. 유지할 shape 좌표 화이트리스트 (cm) ──
    # (left_cm, top_cm, tol_cm)
    # 방식 A: 베이스 6개 + 레이아웃 entity 표 4개(신탁사/차주/시공사/대주)를 보존
    keep_specs = [
        (1.09,  1.11, 0.35),   # [0]  "02 금융 개요"
        (1.09,  1.64, 0.35),   # [3]  "2.1 투자구조도"
        (1.09,  2.29, 0.35),   # [16] 인트로 텍스트
        (1.09,  3.87, 0.35),   # [1]  상단 구분선
        (1.09, 18.16, 0.35),   # [2]  하단 구분선
        (18.53, 18.35, 0.35),  # [28] 페이지번호 (우하단)
        (10.47,  4.98, 0.35),  # [5]  신탁사 표 (2r×1c)
        (10.48,  9.55, 0.35),  # [6]  차주 표   (2r×1c)
        (10.20, 14.26, 0.35),  # [7]  시공사 표 (2r×1c)
        (18.87,  5.44, 0.35),  # [9]  대주 표   (4r×1c → Tr.C 삭제 후 3r)
    ]

    def _is_kept(shape):
        l_cm = shape.left / 360000 if shape.left is not None else -999
        t_cm = shape.top  / 360000 if shape.top  is not None else -999
        for spec_l, spec_t, tol in keep_specs:
            if abs(l_cm - spec_l) <= tol and abs(t_cm - spec_t) <= tol:
                return True
        return False

    # ── 2. 화이트리스트 외 shape 전부 XML 에서 삭제 ──
    to_remove = [sh for sh in slide.shapes if not _is_kept(sh)]
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    print(f"[build_slide_7] 베이스 정리: {len(to_remove)}개 삭제, "
          f"{len(slide.shapes)}개 유지")

    # ── 3. 인트로 텍스트 교체 ──
    intro_sh = _find_shape_by_pos(slide, 1.09, 2.29)
    if intro_sh:
        _replace_text_keep_runs(intro_sh.text_frame, intro)

    # ── 4. 공통 푸터 (사업명 | 순번) — 기존 페이지번호 placeholder 대체 ──
    add_footer(slide, page_num, business_name=business_name or None)

    # ── STEP 5-2 (방식 A): 레이아웃 entity 표 4개를 PDF p4 좌표로 재배치 ──
    entities = data.get("entities", {})

    # 박스별 헤더색 — 신도림 팔레트로 통일 (PALETTE 상수 사용)
    _H_BORROW  = PALETTE["navy_dark"]    # 08377C 차주
    _H_TRUSTEE = PALETTE["blue"]         # 0063A1 신탁사
    _H_CONSTR  = PALETTE["blue"]         # 0063A1 시공사
    _H_LENDER  = PALETTE["maroon"]       # 8C4A59 대주 (유지)
    _BODY_BLACK = RGBColor(0, 0, 0)

    def _set_cell_text_color(cell, rgb):
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = rgb

    def _place_table(orig_l_cm, orig_t_cm):
        return _find_table_by_pos(slide, orig_l_cm, orig_t_cm)

    # 원본 좌표로 4개 표를 먼저 확보 (재배치 전에)
    t_trustee = _place_table(10.47, 4.98)
    t_borrow  = _place_table(10.48, 9.55)
    t_constr  = _place_table(10.20, 14.26)
    t_lender  = _place_table(18.87, 5.44)

    def _layout_box(sh, left_in, top_in, w_in, row_h, header_rgb):
        """표를 PDF p4 좌표/크기로 재배치하고 헤더색·헤더 텍스트(흰색) 적용."""
        if sh is None:
            print("[build_slide_7] 경고: entity 표 못 찾음")
            return
        tbl = sh.table
        sh.left, sh.top = Inches(left_in), Inches(top_in)
        sh.width, sh.height = Inches(w_in), Inches(sum(row_h))
        tbl.columns[0].width = Inches(w_in)
        for i, rh in enumerate(row_h):
            if i < len(tbl.rows):
                tbl.rows[i].height = Inches(rh)
        hc = tbl.cell(0, 0)
        hc.fill.solid()
        hc.fill.fore_color.rgb = header_rgb
        _set_cell_text_color(hc, _C_WHITE)

    # 대주: Tr.C 행(4번째) 삭제 → 3r 유지
    if t_lender is not None:
        trs = t_lender.table._tbl.tr_lst
        if len(trs) >= 4:
            t_lender.table._tbl.remove(trs[3])

    # 재배치 (신탁사/시공사/차주 = 2r, 대주 = 3r)
    _layout_box(t_trustee, 4.5, 2.0, 1.5, [0.40, 0.60],        _H_TRUSTEE)
    _layout_box(t_constr,  1.0, 3.5, 1.5, [0.40, 0.60],        _H_CONSTR)
    _layout_box(t_borrow,  4.5, 3.5, 1.5, [0.40, 0.60],        _H_BORROW)
    _layout_box(t_lender,  8.0, 3.2, 1.5, [0.40, 0.50, 0.50],  _H_LENDER)

    # 본문 셀 텍스트 교체 (회사명, 검정)
    def _set_body(sh, row, text):
        if sh is None:
            return
        cell = sh.table.cell(row, 0)
        _replace_text_keep_runs(cell.text_frame, text)
        _set_cell_text_color(cell, _BODY_BLACK)

    _set_body(t_trustee, 1, entities.get("trustee", "신한자산신탁"))
    _set_body(t_borrow,  1, entities.get("borrower", "더함도시개발"))
    _set_body(t_constr,  1, entities.get("constructor", "포스코이앤씨"))
    # 대주: 헤더 "대주", Tr.A / Tr.B (금액 없음)
    if t_lender is not None:
        _replace_text_keep_runs(t_lender.table.cell(0, 0).text_frame, "대주")
        _set_cell_text_color(t_lender.table.cell(0, 0), _C_WHITE)
        _set_body(t_lender, 1, "Tr.A")
        _set_body(t_lender, 2, "Tr.B")

    # ── STEP 3: 화살표(연결선) + 라벨을 PDF p4 배치대로 그림 ──
    _FONT_BOLD = "피플폰트 Bold"   # 라벨 폰트 (레이아웃 따름)
    _C_ARROW = RGBColor(0x40, 0x40, 0x40)   # 진회색

    def _arrow(x1, y1, x2, y2, *, double, bent=False, adj1=None):
        """연결선. double=True 양 끝 화살촉, False 끝점만.
        bent=True 면 꺾인선(ELBOW, prst=bentConnector3), 기본은 직선(line).
        adj1(0~100000) 지정 시 꺾이는 지점을 조절(수평 구간 길이)."""
        cxn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW if bent else MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        cxn.line.color.rgb = _C_ARROW
        cxn.line.width = Pt(1.5)
        cxn.shadow.inherit = False
        ln = cxn.line._get_or_add_ln()
        if double:
            ln.append(ln.makeelement(
                _qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        ln.append(ln.makeelement(
            _qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        # 꺾임 지점 조절 (bentConnector3 의 adj1)
        if bent and adj1 is not None:
            prstGeom = cxn._element.find('.//' + _qn('a:prstGeom'))
            if prstGeom is not None:
                for av in prstGeom.findall(_qn('a:avLst')):
                    prstGeom.remove(av)
                avLst = etree.SubElement(prstGeom, _qn('a:avLst'))
                gd = etree.SubElement(avLst, _qn('a:gd'))
                gd.set('name', 'adj1')
                gd.set('fmla', 'val %d' % int(adj1))
        return cxn

    def _arrow_label(left, top, w, h, text):
        """화살표 옆 라벨 텍스트박스 (맑은 고딕 9pt, 검정, 배경 투명)."""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _FONT_BOLD      # 화살표 라벨도 "피플폰트 Bold" (레이아웃 따름)
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0, 0, 0)
        return tb

    def _freeform_arrow(points_in):
        """꺾는 점(인치 좌표 리스트)을 잇는 freeform 꺾은선 + 끝점 화살촉.
        선색 404040, 두께 1.5pt, 채움 없음, tailEnd=triangle."""
        pts = [(int(x * 914400), int(y * 914400)) for (x, y) in points_in]
        fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
        fb.add_line_segments(pts[1:], close=False)
        shp = fb.convert_to_shape()
        shp.fill.background()              # 채움 없음(열린 선)
        shp.line.color.rgb = _C_ARROW
        shp.line.width = Pt(1.5)
        shp.shadow.inherit = False
        ln = shp.line._get_or_add_ln()
        ln.append(ln.makeelement(          # 직각 모서리(round join 제거)
            _qn('a:miter'), {'lim': '800000'}))
        ln.append(ln.makeelement(          # 끝점(P4) 화살촉
            _qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        return shp

    # (1) 신탁사 ↔ 차주 : 세로 양방향 (중심 x=5.25, y 3.0→3.5) — 라벨 "신탁계약"만
    _arrow(5.25, 3.0, 5.25, 3.5, double=True)
    _arrow_label(3.30, 3.05, 1.10, 0.40, "신탁계약")              # 세로선 좌측

    # (2) 시공사 ↔ 차주 : 가로 양방향 (y=4.0, x 2.5→4.5)
    _arrow(2.5, 4.0, 4.5, 4.0, double=True)
    _arrow_label(2.75, 3.62, 1.50, 0.30, "공사도급계약")          # 화살표 위

    # (3) 차주 ↔ 대주 : 가로 양방향 (y=4.0, x 6.0→8.0)  ← L6.0/T4.0/W2.0/H0.0
    _arrow(6.0, 4.0, 8.0, 4.0, double=True)
    _arrow_label(6.25, 3.50, 1.50, 0.30, "대출약정")             # 화살표 위

    # (4) 신탁사 → 대주 : freeform 꺾은선(ㄱ자, miter 직각) — 꺾는 점 직접 고정
    #     P1(6.10,2.55) 신탁사 우변서 0.1" 띄움 → P2(7.50,2.55) 수평 →
    #     P3(7.50,3.55) 수직 꺾음 → P4(7.90,3.55) 대주 좌변서 0.1" 띄움(화살촉)
    #     마지막 수평구간 0.40", 양끝 박스와 0.10" 간격
    _freeform_arrow([(6.10, 2.55), (7.50, 2.55), (7.50, 3.55), (7.90, 3.55)])
    _arrow_label(6.30, 2.45, 1.95, 0.30, "담보신탁 우선수익권")    # 꺾은선 위 중앙

    # ── 박스 글씨 12~13pt + 셀 여백 0 (구조도 고정 규칙) ──
    #   신탁사/차주/시공사 = 13pt, 대주(Tr.A·Tr.B) = 12pt
    #   폰트명/색/정렬은 유지(크기만 변경), 여백 0 으로 공간 확보
    def _set_table_font(sh, size_pt):
        if sh is None:
            return
        tbl = sh.table
        for r in range(len(tbl.rows)):
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                cell.margin_left = cell.margin_right = Inches(0)
                cell.margin_top = cell.margin_bottom = Inches(0)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(size_pt)

    _set_table_font(t_trustee, 13)
    _set_table_font(t_borrow, 13)
    _set_table_font(t_constr, 13)
    _set_table_font(t_lender, 12)

    # ── 다이어그램 전체를 아래로 이동(세로 중앙 정렬) ──
    #   박스4 + 화살표4 + 라벨4 를 동일 offset(+DY)으로만 이동 → 상대위치·간격 유지
    #   헤더/인트로/구분선/로고/푸터는 이동하지 않음
    _DY = 1.0   # 인치 (헤더-푸터 사이 중앙 정렬용)
    _LBLS = {"신탁계약", "공사도급계약", "대출약정", "담보신탁 우선수익권"}
    for sh in slide.shapes:
        tp = str(sh.shape_type)
        is_diag = (sh.has_table
                   or "Connector" in (sh.name or "")
                   or "FREEFORM" in tp
                   or (sh.has_text_frame and sh.text_frame.text.strip() in _LBLS))
        if is_diag and sh.top is not None:
            sh.top = sh.top + int(_DY * 914400)

    return slide
