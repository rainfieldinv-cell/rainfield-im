"""
frame_builders.py
─────────────────────────────────────────────────────────
"틀 먼저 → 채움" 방식의 본문 슬라이드 빌더 + 페이지 레이아웃 자동 판별기.

build_full_presentation(_build_content_block)이 변환 중 페이지마다
classify_page()로 유형을 자동 판별하고, 우리가 구현한 유형(E)은
이 모듈의 전용 빌더로, 아직 미구현 유형은 기존 build_content_slide로
폴백(fallback)합니다.

유형 (rainfield-layout-patterns 참고):
  A. 전체폭 데이터 표              (미구현 → 폴백)
  B. 좌우 [이미지 표 | 텍스트 표]  (미구현 → 폴백)
  C. 좌우 [이미지 | 비교 표]       (미구현 → 폴백)
  D. 표/차트                       (미구현 → 폴백)
  E. 데이터 표 + 이미지 표         (구현 — build_frame_e_slide)

좌표는 사람이 만든 제안서(넷마블 G타워) 슬9 실측값에 맞춤:
  좌 데이터표 L0.43 W4.99 / 우 이미지표 L5.56 W4.85 / 본문 T1.83
"""

import io
import re
from collections import Counter

from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from modules import page_builders
from modules.page_builders import (
    clone_slide_layout,
    _replace_footer_business_name,
    _replace_text_frame_content,
    auto_resize_text_to_fit,
)
from modules.ai_slide_builders import (
    style_table, PALETTE, _replace_text_keep_runs,
)

# 이미지 행 라벨(데이터표에서 제외하고 이미지표 라벨로 사용)
IMG_LABELS = {"위치도", "조감도", "조 감 도", "위 치 도", "투시도", "전경", "건물전경", "건물 전경"}

# 유형 E 판별 임계값
_E_MIN_DATA_ROWS = 5       # label-value 표 최소 행 수
_E_MIN_IMG_AREA  = 8000    # 사진으로 인정할 최소 픽셀 면적(px*px)
_E_MAX_IMGS      = 3       # 이미지 표에 넣을 최대 사진 수


# ──────────────────────────────────────────────────────
# 표 정규화 / 판별 유틸
# ──────────────────────────────────────────────────────
def _to_label_value(rows):
    """2D 표 → [[라벨, 값], ...]. 행마다 비어있지 않은 첫 셀=라벨, 나머지 join=값."""
    out = []
    for row in rows or []:
        cells = [str(c or "").strip() for c in row]
        ne = [c for c in cells if c and c != "•"]
        if not ne:
            continue
        out.append([ne[0], " ".join(ne[1:]).replace("•", "").strip()])
    return out


def _is_label_value(rows):
    """행 대부분(60%+)이 비어있지않은 셀 2개 이하 → 구분|내용 형태로 간주."""
    if not rows:
        return False
    le2 = 0
    for r in rows:
        ne = [c for c in r if str(c or "").strip() and str(c or "").strip() != "•"]
        if len(ne) <= 2:
            le2 += 1
    return le2 >= max(2, len(rows) * 0.6)


def _big_images(images):
    """로고·아이콘 제외한 사진급 이미지만."""
    return [im for im in (images or [])
            if (im.get("width", 0) * im.get("height", 0)) >= _E_MIN_IMG_AREA]


def _main_table(tables):
    return max(tables, key=lambda t: len(t) * max((len(r) for r in t), default=0),
               default=None)


def classify_page(page: dict) -> str:
    """
    파싱된 페이지 dict(tables/images 포함) → 레이아웃 유형 문자열.

    현재 'E'만 전용 빌더가 있고 나머지는 'fallback'(기존 빌더)으로 보냅니다.
    """
    tables = page.get("tables") or []
    images = page.get("images") or []
    if not tables:
        return "fallback"
    main = _main_table(tables)
    if not main or len(main) < _E_MIN_DATA_ROWS:
        return "fallback"
    big = _big_images(images)
    # 유형 E: 큰 사진 1~3장 + 구분|내용(label-value) 데이터 표
    if 1 <= len(big) <= _E_MAX_IMGS and _is_label_value(main):
        return "E"
    return "fallback"


# ──────────────────────────────────────────────────────
# 헤더(제목/부제목/본문) 교체 — build_content_slide와 동일 규칙
# ──────────────────────────────────────────────────────
def _fill_header(slide, title, subtitle, intro):
    LEFT_MARGIN_CM = 1.09
    TOP_LIMIT_CM   = 4.0
    TOLERANCE      = Cm(0.5)
    header_shapes = sorted(
        [sh for sh in slide.shapes
         if sh.has_text_frame
         and abs(sh.left - Cm(LEFT_MARGIN_CM)) < TOLERANCE
         and sh.top / 360000 < TOP_LIMIT_CM],
        key=lambda s: s.top,
    )
    if len(header_shapes) >= 1:
        _replace_text_frame_content(header_shapes[0].text_frame, title)
    if len(header_shapes) >= 2:
        _replace_text_frame_content(header_shapes[1].text_frame, subtitle)
    if len(header_shapes) >= 3:
        if intro and intro.strip():
            _replace_text_frame_content(header_shapes[2].text_frame, intro.strip())
            auto_resize_text_to_fit(header_shapes[2].text_frame, max_size=10.5, min_size=8.0)
        else:
            _replace_text_frame_content(header_shapes[2].text_frame, "")


# ──────────────────────────────────────────────────────
# 유형 E 빌더 — [데이터 표 | 이미지 표]
# ──────────────────────────────────────────────────────
def build_frame_e_slide(prs, *, title: str, subtitle: str, intro: str,
                        data_rows: list, imgs: list, img_labels: list = None,
                        business_name: str = "", sub_label: str = None):
    """
    유형 E 슬라이드: 왼쪽 구분|내용 데이터 표, 오른쪽 구분|내용 이미지 표.

    Parameters
    ----------
    data_rows  : [[라벨, 값], ...]  (행 0 = 헤더 "구분"/"내용")
    imgs       : [{"data"|"bytes": bytes, "width"|"w": int, "height"|"h": int}, ...]
    img_labels : 이미지 표 각 행 라벨 (없으면 ["조감도","위치도",...])
    sub_label  : 표 위 소제목바 텍스트 (None이면 생략)
    """
    slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)
    _fill_header(slide, title, subtitle, intro)

    # ── 소제목바 (선택) ──
    if sub_label:
        sub = slide.shapes.add_textbox(Inches(0.43), Inches(1.5), Inches(3.0), Inches(0.3))
        run = sub.text_frame.paragraphs[0].add_run()
        run.text = sub_label
        run.font.name = "피플폰트 Bold"
        run.font.size = Pt(12)
        run.font.color.rgb = PALETTE["navy_dark"]

    # ── 왼쪽: 데이터 표 (L0.43 W4.99) ──
    rows = max(1, len(data_rows))
    gL = slide.shapes.add_table(rows, 2, Inches(0.43), Inches(1.83),
                                Inches(4.99), Inches(rows * 0.45))
    tL = gL.table
    tL.columns[0].width = Inches(1.40)
    tL.columns[1].width = Inches(3.59)
    for ri, pair in enumerate(data_rows):
        lab = pair[0] if len(pair) > 0 else ""
        val = pair[1] if len(pair) > 1 else ""
        _replace_text_keep_runs(tL.cell(ri, 0).text_frame, lab)
        _replace_text_keep_runs(tL.cell(ri, 1).text_frame, val)
        if ri >= 1:
            from pptx.enum.text import PP_ALIGN
            tL.cell(ri, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    style_table(gL, has_header=True, label_cols=(0,),
                header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])

    # ── 오른쪽: 이미지 표 (L5.56 W4.85) ──
    labels = img_labels or ["조감도", "위치도", "투시도"]
    img_rows = max(1, len(imgs))
    R_L, R_T, R_W = 5.56, 1.83, 4.85
    ROW_H = 2.30
    gR = slide.shapes.add_table(img_rows + 1, 2, Inches(R_L), Inches(R_T),
                                Inches(R_W), Inches(0.4 + img_rows * ROW_H))
    tR = gR.table
    tR.columns[0].width = Inches(0.95)
    tR.columns[1].width = Inches(3.90)
    _replace_text_keep_runs(tR.cell(0, 0).text_frame, "구 분")
    _replace_text_keep_runs(tR.cell(0, 1).text_frame, "내 용")
    tR.rows[0].height = Inches(0.4)
    for i in range(img_rows):
        _replace_text_keep_runs(tR.cell(i + 1, 0).text_frame,
                                labels[i] if i < len(labels) else "사진")
        tR.rows[i + 1].height = Inches(ROW_H)
    style_table(gR, has_header=True, label_cols=(0,),
                header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])

    # ── 이미지 삽입 (내용칸 위 비율유지·중앙) ──
    content_L = R_L + 0.95
    content_W = 3.90
    pad = 0.1
    for i, im in enumerate(imgs):
        data = im.get("data") or im.get("bytes")
        iw0  = im.get("width") or im.get("w") or 1
        ih0  = im.get("height") or im.get("h") or 1
        if not data:
            continue
        cell_T = R_T + 0.4 + i * ROW_H
        boxL, boxT = content_L + pad, cell_T + pad
        boxW, boxH = content_W - 2 * pad, ROW_H - 2 * pad
        scale = min(boxW / iw0, boxH / ih0)
        iw, ih = iw0 * scale, ih0 * scale
        px = boxL + (boxW - iw) / 2
        py = boxT + (boxH - ih) / 2
        slide.shapes.add_picture(io.BytesIO(data), Inches(px), Inches(py),
                                 Inches(iw), Inches(ih))

    _replace_footer_business_name(slide, business_name)
    return slide


# ──────────────────────────────────────────────────────
# 라우터 진입점 — _build_content_block에서 호출
# ──────────────────────────────────────────────────────
def build_page_auto(prs, page: dict, *, title: str, subtitle: str,
                    body_text: str, business_name: str) -> bool:
    """
    페이지 유형을 자동 판별해 틀-우선 빌더로 슬라이드를 만든다.
    LLM 토글(env RAINFIELD_LLM=1)이 켜져 있고 원문이 있으면 LLM 구조화 경로 우선.
    아니면 유형 E면 전용 E빌더, 그 외 모든 페이지는 범용 골격 빌더로 처리.
    처리했으면 True(항상), 예외 시에만 호출측이 기존 빌더로 폴백.
    """
    raw = page.get("raw_text", "") or ""
    sec = page.get("_sec_int")

    # ── 고정 페이지: 투자구조도(도형으로 직접 그림) ──
    if page.get("_invest_diagram") is not None:
        build_invest_diagram_slide(
            prs, page.get("_invest_diagram") or {},
            section_label=page.get("section_label"),
            subtitle=page.get("subtitle"),
            intro=page.get("_invest_intro", ""),
            business_name=business_name)
        return True

    # ── 고정 페이지: 투자구조도(원본 구조도 영역 이미지 캡처 — 비활성) ──
    if page.get("_invest_image"):
        build_invest_image_slide(
            prs,
            section_label=page.get("section_label"),
            subtitle=page.get("subtitle"),
            intro=page.get("_invest_intro", ""),
            img=page.get("_invest_image"),
            business_name=business_name)
        return True

    # ── 고정 페이지: [섹션1] 사모사채 개요 → 전용 고정 빌더(구분열 불변, 내용만 LLM) ──
    if sec == 1 and raw.strip():
        try:
            from modules.ai_slide_builders import (
                generate_sasae_overview, build_slide_5_sasae_overview)
            res = generate_sasae_overview(raw)
            if res.get("ok") and res.get("data"):
                sl = build_slide_5_sasae_overview(prs, res["data"], business_name=business_name)
                # ★사모사채개요 표 글꼴 전부 Bold 통일(뒤죽박죽 방지 — 사용자 지시)
                try:
                    for sh in sl.shapes:
                        if sh.has_table:
                            for row in sh.table.rows:
                                for cell in row.cells:
                                    for pp in cell.text_frame.paragraphs:
                                        _set_para_default_font(pp, 10.5, _FONT_BOLD)
                                        for rn in pp.runs:
                                            rn.font.name = _FONT_BOLD
                                            rn.font.size = Pt(10.5)
                except Exception:
                    pass
                return True
        except Exception as _exc:
            print(f"[build_page_auto] 사모사채개요 고정빌더 실패 → 구조화 폴백: {_exc}")

    # enrich_and_number()로 미리 구조화된 결과가 있으면 사용
    struct = page.get("_struct")
    if isinstance(struct, dict):
        build_structured_slide(
            prs, struct, business_name=business_name,
            section_label=page.get("section_label"),
            subtitle=page.get("subtitle"),
            images=page.get("images"),
            red_texts=page.get("_red_texts"),
            underline_texts=page.get("_underline_texts"),
            fill_texts=page.get("_filled_texts"))
        return True

    ptype = classify_page(page)

    if ptype == "E":
        tables = page.get("tables") or []
        main = _main_table(tables)
        norm = _to_label_value(main)

        def _norm_key(s):
            return s.replace(" ", "")
        img_label_set = {_norm_key(x) for x in IMG_LABELS}
        img_labels_found = [r[0] for r in norm if _norm_key(r[0]) in img_label_set]
        data_rows = [r for r in norm if _norm_key(r[0]) not in img_label_set]
        imgs = _big_images(page.get("images") or [])[: _E_MAX_IMGS]
        if data_rows and imgs:
            build_frame_e_slide(
                prs, title=title, subtitle=subtitle, intro=body_text,
                data_rows=data_rows, imgs=imgs,
                img_labels=img_labels_found or None,
                business_name=business_name,
            )
            return True

    # 그 외 모든 페이지 → 범용 골격 빌더
    build_universal_content_slide(
        prs, title=title, subtitle=subtitle, intro=body_text,
        tables=page.get("tables") or [], business_name=business_name,
    )
    return True


# ──────────────────────────────────────────────────────
# 범용 골격 빌더 — 대전 제안서 실측 골격 재현
#   섹션라벨/소제목(템플릿 자리) + 인트로(별도 글상자) + 표(정규화·스택) + 푸터/출처
#   ★본문은 절대 제목/소제목 자리에 넣지 않고 새 글상자로. 글상자 내부여백 0.
# ──────────────────────────────────────────────────────
_BLACK = PALETTE["black"]
_FONT_LIGHT = "피플폰트 Light"
_FONT_BOLD  = "피플폰트 Bold"

# 대전 실측 좌표(in)
_INTRO_L, _INTRO_T, _INTRO_W = 0.43, 1.01, 8.65
_BODY_TOP_START = 1.55      # 표/본문 영역 시작
_BODY_BOTTOM    = 6.85      # 각주/푸터 전 한계
_TBL_L, _TBL_W  = 0.43, 9.97
_SRC_L, _SRC_T, _SRC_W = 0.43, 7.19, 8.65


def _est_text_height(text, W, size):
    """글상자 높이(in) 추정 — 자동줄바꿈 고려. 빈 줄 군더더기 없이 내용에 딱 맞춤."""
    char_w = size * 0.78 / 72.0           # 한글 기준 글자폭 근사
    cpl = max(1, int(W / char_w))
    line_h = size * 1.32 / 72.0           # 한 줄 높이(여유 최소화)
    lines = 0
    for ln in str(text).split("\n"):
        lines += max(1, -(-len(ln) // cpl))   # ceil
    return max(line_h, lines * line_h) + 0.03


def _emph_segments(line, red_set, ul_set=None):
    """한 줄을 [(텍스트, is_red, is_underline)] 조각으로 분할.
       원본의 빨간 글씨/밑줄 부분만 강조(겹쳐도 문자 단위 마스크로 정확히 처리)."""
    n = len(line)
    if n == 0:
        return [("", False, False)]
    red_mask = [False] * n
    ul_mask = [False] * n

    def _mark(mask, frags, minlen):
        for f in sorted({str(s).strip() for s in (frags or []) if len(str(s).strip()) >= minlen},
                        key=len, reverse=True):
            start = 0
            while True:
                idx = line.find(f, start)
                if idx == -1:
                    break
                for k in range(idx, idx + len(f)):
                    mask[k] = True
                start = idx + len(f)

    _mark(red_mask, red_set, 4)   # 빨강은 4글자 이상(루시드 등 과다 방지)
    _mark(ul_mask, ul_set, 2)     # 밑줄은 2글자 이상
    if not any(red_mask) and not any(ul_mask):
        return [(line, False, False)]
    segs, i = [], 0
    while i < n:
        r, u = red_mask[i], ul_mask[i]
        j = i
        while j < n and red_mask[j] == r and ul_mask[j] == u:
            j += 1
        segs.append((line[i:j], r, u))
        i = j
    return segs


_BULLET_RE = re.compile(r"^\s*[•·▶◦●○■◆\-\*]\s+")


def _set_para_bullet(p, char="•", *, marL=0.17, font="Arial"):
    """문단에 '진짜' 글머리 기호(buChar) + 행잉 인덴트 적용 — 줄바꿈 시 텍스트가
       불릿 글자 아래가 아니라 첫 글자 위치에 정렬됨(원본처럼 깔끔). 글자 prefix 흉내 아님."""
    pPr = p._p.get_or_add_pPr()
    emu = int(marL * 914400)
    pPr.set("marL", str(emu))
    pPr.set("indent", str(-emu))
    for tag in ("a:buClrTx", "a:buClr", "a:buSzTx", "a:buSzPct", "a:buSzPts",
                "a:buFontTx", "a:buFont", "a:buNone", "a:buAutoNum", "a:buChar"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    anchor = None
    for tag in ("a:defRPr", "a:extLst"):
        el = pPr.find(qn(tag))
        if el is not None:
            anchor = el
            break
    if anchor is not None:
        anchor.addprevious(buFont)
        anchor.addprevious(buChar)
    else:
        pPr.append(buFont)
        pPr.append(buChar)


_CELL_BULLET_RE = re.compile(r"^(\s*)([•·▶◦●○■◆◇]|[-])\s+")


def _bulletize_text_frame(tf):
    """표 셀 안에서 마커(▶/•/-)로 시작하는 문단을 '진짜' 글머리 기호로 바꿈(흉내 글자 제거).
       ▶/• = 상위 항목(들여쓰기 0단), '-' = 하위 항목(1단 들여쓰기). 줄바꿈 시 정렬 깔끔."""
    for p in tf.paragraphs:
        m = _CELL_BULLET_RE.match(p.text)
        if not m:
            continue
        ch = m.group(2)
        sub = (ch == "-")
        n = len(m.group(0))                      # 마커(+공백) 글자 수만큼 런에서 제거
        for run in p.runs:
            if n <= 0:
                break
            tx = run.text
            if len(tx) <= n:
                n -= len(tx)
                run.text = ""
            else:
                run.text = tx[n:]
                n = 0
        _set_para_bullet(p, ch, marL=(0.36 if sub else 0.18), font=_FONT_LIGHT)


def _add_textbox(slide, L, T, W, H, text, *, size=10.5, bold=False,
                 color=None, align=PP_ALIGN.LEFT, red_set=None, ul_set=None, bullet=False):
    """글상자 생성 — 내부여백 0, 텍스트 내용에 딱 맞는 높이(군더더기 여백 없음). 반환=(shape, 높이in).
       red_set/ul_set 지정 시 원본 빨간 글씨/밑줄에 해당하는 부분만 빨강·밑줄로 표시.
       bullet=True 면 각 줄을 '진짜' 글머리 기호(buChar)로 — 줄 맨 앞 마커(•/▶/-)는 떼고 적용."""
    text = str(text).rstrip("\n ")        # 끝의 빈 줄/공백 제거(빈 여백 방지)
    est_h = _est_text_height(text, W, size)
    tb = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(est_h))
    tf = tb.text_frame
    tf.word_wrap = True
    # 자동맞춤 끔 — PPT가 안 줄여주는 빈 여백 문제 방지(계산한 높이 그대로 사용)
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    nm = _FONT_BOLD if bold else _FONT_LIGHT
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.font.size = Pt(size)                       # 문단 기본(빈 줄에도 적용)
        p.font.name = nm
        seg_line = line
        if bullet and line.strip():
            m = _CELL_BULLET_RE.match(line)
            if m:                                    # ▶/•/- 마커 보존(- = 하위 1단)
                ch = m.group(2)
                seg_line = line[len(m.group(0)):]
                _set_para_bullet(p, ch, marL=(0.36 if ch == "-" else 0.18), font=_FONT_LIGHT)
            else:                                    # 마커 없는 평문 → 기본 • 불릿
                _set_para_bullet(p, "•", marL=0.18, font=_FONT_LIGHT)
        for seg, is_red, is_ul in _emph_segments(seg_line, red_set, ul_set):
            run = p.add_run()
            run.text = seg
            run.font.name = nm
            run.font.size = Pt(size)
            run.font.color.rgb = _RED if is_red else (color or _BLACK)
            if is_ul:
                run.font.underline = True
    return tb, est_h


def _drop_empty_cols(rows):
    """완전히 빈 열(phantom column) 제거."""
    if not rows:
        return rows
    ncol = max(len(r) for r in rows)
    keep = [c for c in range(ncol)
            if any(str((r[c] if c < len(r) else "") or "").strip() for r in rows)]
    return [[(r[c] if c < len(r) else "") for c in keep] for r in rows]


def _looks_like_prose(rows):
    """표가 아니라 문장이 칸칸이 쪼개진 가짜 표인지 판정."""
    norm = _to_label_value(rows)
    if not norm:
        return True
    empty_first = sum(1 for r in rows if not str((r[0] if r else "") or "").strip())
    long_labels = sum(1 for lab, _ in norm if len(lab) > 18)
    return empty_first > len(rows) * 0.4 or long_labels > len(norm) * 0.5


def _add_norm_table(slide, rows, L, T, W):
    """정규화 표 삽입 → (graphicFrame, 사용높이in). 칸 너비를 내용 형태에 맞춤."""
    label_value = _is_label_value(rows)
    if label_value:
        data = _to_label_value(rows)
        ncol = 2
    else:
        data = _drop_empty_cols(rows)
        ncol = max((len(r) for r in data), default=0)
    if not data or ncol == 0:
        return None, 0.0

    nrow = len(data)
    row_h = 0.34
    height = max(row_h, nrow * row_h)
    gf = slide.shapes.add_table(nrow, ncol, Inches(L), Inches(T),
                                Inches(W), Inches(height))
    t = gf.table

    if label_value:
        lab_w = min(1.6, W * 0.24)
        t.columns[0].width = Inches(lab_w)
        t.columns[1].width = Inches(W - lab_w)
        for ri, pair in enumerate(data):
            _replace_text_keep_runs(t.cell(ri, 0).text_frame, pair[0])
            _replace_text_keep_runs(t.cell(ri, 1).text_frame,
                                    pair[1] if len(pair) > 1 else "")
            t.cell(ri, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        style_table(gf, has_header=False, label_cols=(0,),
                    label_fill=PALETTE["label_gray"])
    else:
        widths = _content_col_widths(data, ncol, W, 10.5)
        for c in range(ncol):
            t.columns[c].width = Inches(widths[c])
        for ri, row in enumerate(data):
            for ci in range(ncol):
                val = str((row[ci] if ci < len(row) else "") or "").strip()
                _replace_text_keep_runs(t.cell(ri, ci).text_frame, val)
        style_table(gf, has_header=True, label_cols=(),
                    header_fill=PALETTE["navy_dark"])
    return gf, height


def build_universal_content_slide(prs, *, title: str, subtitle: str, intro: str,
                                  tables: list, business_name: str = "",
                                  source: str = None):
    """대전 실측 골격을 재현하는 범용 본문 슬라이드."""
    slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)

    # 1) 섹션라벨/소제목은 템플릿 자리에, 본문 자리(plh[2])는 비움
    _fill_header(slide, title, subtitle, "")

    # 2) 인트로 = 별도 글상자 (Light 10.5, 여백0). ★제목칸에 안 넣음
    top = _BODY_TOP_START
    if intro and intro.strip():
        _add_textbox(slide, _INTRO_L, _INTRO_T, _INTRO_W, 0.45,
                     intro.strip(), size=10.5, bold=False)
        top = max(top, _INTRO_T + 0.55)

    # 3) 표 — 가짜표(문장쪼개짐) 제외, 정규화해 세로 스택
    real_tables = [tb for tb in (tables or []) if tb and not _looks_like_prose(tb)]
    for tb in real_tables:
        if top > _BODY_BOTTOM - 0.5:
            break   # TODO: 페이지 분할 (1/2)(2/2) — 후속 작업
        avail = _BODY_BOTTOM - top
        gf, used = _add_norm_table(slide, tb, _TBL_L, top, _TBL_W)
        if gf is None:
            continue
        # 표가 남은 공간보다 크면 그대로 두되 다음 표는 폴백 한계에 막힘
        top += min(used, avail) + 0.25

    # 4) 푸터(사업명) + 출처 글상자(여백0)
    _replace_footer_business_name(slide, business_name)
    if source and source.strip():
        _add_textbox(slide, _SRC_L, _SRC_T, _SRC_W, 0.22,
                     f"* 출처 : {source.strip()}", size=9,
                     color=PALETTE["gray_text"])
    return slide


# ──────────────────────────────────────────────────────
# 구조화(LLM) 슬라이드 빌더 — llm_structure.structure_page() 결과를 렌더
# ──────────────────────────────────────────────────────
_TOTAL_KW = ("합계", "소계", "총계", "합 계", "총 계", "총 합계", "총합계")


def _classify_total_rows(data, ncol):
    """행 인덱스를 (소계행, 합계/총합계행) 두 집합으로 분류.
       ★'라벨 셀'(짧은 칸 ≤12자)에 키워드가 있을 때만 인정 — 긴 값 속 '합계'(예: 연면적
         '지상…/지하…/합계 172,…')는 제외. 소계는 '소계'만, 합계/총계/총합계는 grand-total."""
    # ★재무상태표 항목(자산총계/부채총계/자본총계)은 합계행 아님 — 원본도 색칠 안 함 → 분류 제외
    _FIN_EXCLUDE = ("자산총계", "부채총계", "자본총계", "자 산 총 계", "부 채 총 계", "자 본 총 계")
    subtotal, grandtotal = set(), set()
    for i, row in enumerate(data):
        cells = [str((row[c] if c < len(row) else "") or "").strip() for c in range(ncol)]
        is_sub = any(len(cx) <= 12 and cx not in _FIN_EXCLUDE and any(k in cx for k in ("소계", "소 계"))
                     for cx in cells)
        #  '총…'(총매출/총비용/총계/총합계 등 짧은 합산 라벨)도 합계행으로 인정
        is_grand = any(len(cx) <= 12 and cx not in _FIN_EXCLUDE
                       and (cx.startswith("총")
                            or any(k in cx for k in ("합계", "총계", "총합계", "합 계", "총 계")))
                       for cx in cells)
        if is_sub:
            subtotal.add(i)
        elif is_grand:
            grandtotal.add(i)
    return subtotal, grandtotal


def _header_groups(header):
    """플랫 헤더에서 다단(그룹) 헤더 감지(시세표: 면적(평)/실거래가·분양가).
       반환 (group_row, sub_row, vmerge_cols, hmerges) 또는 None.
        - group_row: 상위 헤더행(그룹 시작칸=그룹명, 그룹 내 나머지=빈칸, 비그룹칸=원래 이름)
        - sub_row  : 하위 헤더행(그룹칸=세부명, 비그룹칸=빈칸)
        - vmerge_cols: 위·아래 세로병합할 비그룹 칸 인덱스
        - hmerges  : [(start,end,label)] 상위행에서 가로병합할 그룹 범위"""
    h = [str(x or "").strip() for x in (header or [])]
    n = len(h)
    if n < 4:
        return None

    # ★'<그룹> 항목/구분 | <그룹> 금액 | …'(자금용도 Cash-In/Cash-Out) → 2단 헤더로 분리.
    #   원본: 윗행 Cash-In/Cash-Out(각 2칸 묶음), 아랫행 항목(또는 구분) | 금액 | 항목 | 금액.
    #   '항목/구분' 앞에 그룹명이 있을 때만(맨 '항목'/'구분'/'금액' 단순 헤더는 건드리지 않음).
    row0 = [""] * n
    row1 = [""] * n
    hmerges = []
    used = [False] * n
    i = 0
    while i < n:
        sub_kw = next((kw for kw in ("항목", "구분") if kw in h[i]), None)
        lab = h[i].replace(sub_kw, "").strip() if sub_kw else ""
        if sub_kw and lab and i + 1 < n and "금액" in h[i + 1]:
            row0[i] = lab
            row1[i] = sub_kw
            row1[i + 1] = h[i + 1].replace(lab, "").strip() or h[i + 1]   # 금액칸에서 그룹명 제거
            hmerges.append((i, i + 1, lab))
            used[i] = used[i + 1] = True
            i += 2
        else:
            i += 1
    if hmerges:
        for k in range(n):
            if not used[k]:
                row0[k] = h[k]
        return row0, row1, [k for k in range(n) if not used[k]], hmerges

    def _runs(idxs):
        runs = []
        for i in idxs:
            if runs and i == runs[-1][1] + 1:
                runs[-1][1] = i
            else:
                runs.append([i, i])
        return [(a, b) for a, b in runs if b > a]   # 2칸 이상만

    grouped = [False] * n
    row0 = list(h)
    row1 = [""] * n
    hmerges = []
    for (s, e) in _runs([i for i, x in enumerate(h) if "면적(평)" in x or "면적(㎡)" in x]):
        hmerges.append((s, e, "면적(평)"))
        for i in range(s, e + 1):
            grouped[i] = True
            row1[i] = h[i].replace("(평)", "").replace("(㎡)", "")
    for (s, e) in _runs([i for i, x in enumerate(h)
                         if any(k in x for k in ("시점", "거래", "분양가", "평단가"))]):
        hmerges.append((s, e, "실거래가/분양가"))
        for i in range(s, e + 1):
            grouped[i] = True
            row1[i] = h[i]
    if not hmerges:
        return None
    for (s, e, lab) in hmerges:
        row0[s] = lab
        for i in range(s + 1, e + 1):
            row0[i] = ""
    vmerge = [i for i in range(n) if not grouped[i]]
    return row0, row1, vmerge, hmerges


def _merge_vertical_runs(table, data, ncol, header_rows=0):
    """열에서 값 있는 칸 아래로 이어지는 빈 칸들을 세로 병합(반복값=첫 칸만, 아래 빈칸 → 병합).

    ★그룹(병합)열 자동판정: 비-합계 행에서 '값 뒤에 빈칸'이 나오는 열만 세로 그룹열로 본다.
      → 연번(1,2,3…)처럼 매 행 값이 채워진 열은 그룹열이 아니므로 옆 열 병합을 막지 않는다
        (대전 토지개요: c0=연번이지만 구역·소유자·감정평가는 정상 병합돼야 함).
    ★그룹열은 소계행을 관통해 병합(공동주택·아파트·토지비가 자기 소계까지 한 칸), 합계 라벨이
      그 열을 가로로 점유하는 행에서만 끊는다. 비그룹열은 모든 합계행에서 끊는다.
    ★왼쪽의 더 큰 그룹열(앵커) 값이 새로 바뀌면 = 새 그룹 시작 → 그 행의 빈칸은 윗값의 연속이 아님
      (Equity: '롯데건설'(변경소유자)이 '초기사업비' 행으로 잘못 번지던 문제 방지)."""
    nrow = len(data)
    subtotal_rows, grandtotal_rows = _classify_total_rows(data, ncol)
    total_rows = subtotal_rows | grandtotal_rows

    def _cell(ri, ci):
        return str((data[ri][ci] if ri < len(data) and ci < len(data[ri]) else "") or "").strip()

    # ★숫자 측정값 열(전용면적·공급면적·세대수·면적 등)은 세로병합하지 않는다 — 각 행(층/타입)이
    #   고유 수치이므로 빈칸은 '윗값의 연속'이 아니라 '값 없음'이다(운수시설 전용면적 788.63이
    #   빈 지하2층으로 잘못 번지던 문제). 단 Equity 투입현황표(헤더에 '에쿼티')는 토지매매대금 46억이
    #   2개 하위행에 걸쳐야 하므로 예외로 둔다.
    hdr_blob = " ".join(_cell(r, c) for r in range(header_rows) for c in range(ncol))
    _allow_num = "에쿼티" in hdr_blob

    def _is_numeric_col(ci):
        vals = [_cell(ri, ci) for ri in range(header_rows, nrow)
                if ri not in total_rows and _cell(ri, ci)]
        if not vals:
            return False
        num = sum(1 for v in vals
                  if v.replace(",", "").replace(".", "").replace("%", "").isdigit())
        return num >= max(1, int(len(vals) * 0.6))

    def _is_remark_col(ci):     # 비고/비율 등 행마다 독립적인 설명열 → 세로병합 금지
        return any(k in _cell(r, ci) for r in range(header_rows) for k in ("비고", "비율"))
    # 숫자 측정열·비고열은 세로병합 대상에서 제외(Equity 표는 예외)
    no_merge = (set() if _allow_num
                else {ci for ci in range(ncol) if _is_numeric_col(ci) or _is_remark_col(ci)})

    # 그룹(세로 병합) 열 = 비-합계 행에서 '값 뒤 빈칸' 패턴이 있는 열(연번/숫자 측정열 제외)
    def _is_group_col(ci):
        seen = False
        for ri in range(header_rows, nrow):
            if ri in total_rows:
                continue
            if _cell(ri, ci):
                seen = True
            elif seen:
                return True
        return False
    group_cols = [ci for ci in range(ncol) if _is_group_col(ci) and ci not in no_merge]
    gset = set(group_cols)

    # 합계행이 가로 라벨로 점유하는 열 집합 = [최좌측 값 ~ 그 다음 값 직전]
    def _label_cols(ri):
        firsts = [c for c in range(ncol) if _cell(ri, c)]
        if not firsts:
            return set()
        lo = firsts[0]
        hi = next((c for c in firsts if c > lo), ncol)
        return set(range(lo, hi))
    label_cols = {ri: _label_cols(ri) for ri in total_rows}

    primary = group_cols[0] if group_cols else None    # 표의 최상위(맨 왼쪽) 그룹열
    for ci in range(ncol):
        if ci in no_merge:          # 숫자·비고열은 세로병합 자체를 하지 않음(각 행 고유값)
            continue
        is_group = ci in gset
        # 앵커 = 최상위 그룹열. 그 값이 바뀌면 새 그룹 → 하위 열 병합 끊음.
        #   (감정평가는 '구역' 기준으로 묶여야지, 중간 '소유자' 열이 바뀐다고 끊기면 안 됨.
        #    Equity '롯데건설'은 '구분' 바뀔 때만 끊김.)
        left_anchor = primary if (primary is not None and ci != primary) else None
        ri = header_rows
        while ri < nrow:
            # 시작 칸이 합계 라벨에 점유되거나(라벨 행), 비그룹열인데 합계행이면 시작 안 함
            if ri in total_rows and (ci in label_cols.get(ri, ()) or not is_group):
                ri += 1
                continue
            if not _cell(ri, ci):
                ri += 1
                continue
            end = ri
            for rj in range(ri + 1, nrow):
                # 합계행이 이 열을 라벨로 점유 → 경계 / 비그룹열은 모든 합계행이 경계
                if rj in total_rows and (ci in label_cols.get(rj, ()) or not is_group):
                    break
                # 왼쪽 그룹열(앵커) 값이 새로 바뀌면 새 그룹 시작 → 끊음
                if left_anchor is not None and _cell(rj, left_anchor):
                    break
                if _cell(rj, ci):       # 같은 열에 새 값 → 끊음
                    break
                end = rj
            if end > ri:
                try:
                    table.cell(ri, ci).merge(table.cell(end, ci))
                except Exception:
                    pass
            ri = end + 1


def _merge_total_rows(table, data, ncol):
    """합계/소계/총매출 등 합산행에서 라벨을 빈 칸들에 가로 병합
       (예 '터미널 합계(A)'·'총 매출 (A)'가 옆 빈 칸들을 차지)."""
    subtotal, grandtotal = _classify_total_rows(data, ncol)
    total_set = subtotal | grandtotal
    for ri, row in enumerate(data):
        if ri not in total_set:
            continue
        txts = [str((row[c] if c < len(row) else "") or "").strip() for c in range(ncol)]
        first = next((i for i, tx in enumerate(txts) if tx), None)
        if first is None:
            continue
        end = first
        for j in range(first + 1, ncol):
            if txts[j] == "":
                end = j
            else:
                break
        if end > first:
            try:
                table.cell(ri, first).merge(table.cell(ri, end))
            except Exception:
                pass


def _font_for_rowh(row_h):
    if row_h >= 0.30:
        return 10.5
    if row_h >= 0.24:
        return 9
    if row_h >= 0.19:
        return 8
    if row_h >= 0.155:
        return 7
    return 6


def _uniform_font_rowh(nrow, avail):
    """표 전체에 통일 적용할 (글자크기, 행높이). 10.5pt 우선, 안 들어가면 통일로 낮춤(최저 8pt)."""
    for fp in (10.5, 10, 9, 8):
        rh = fp * 1.5 / 72.0 + 0.05   # 한 줄 행높이 근사(in)
        if nrow * rh <= avail:
            return fp, rh
    return 8, 8 * 1.5 / 72.0 + 0.05   # 최저 8pt(넘치면 페이지분할 대상)


def _set_para_default_font(p, font_pt, name):
    """문단 기본 글꼴(defRPr) 크기·글꼴(라틴+한글 ea) 지정 — ★빈 칸에도 적용되어
       빈 칸이 큰 기본 글씨로 행높이를 키우는 문제 제거."""
    p.font.size = Pt(font_pt)
    p.font.name = name                      # a:latin
    sz = str(int(font_pt * 100))
    try:
        defRPr = p._p.get_or_add_pPr().get_or_add_defRPr()
        defRPr.set('sz', sz)
        ea = defRPr.find(qn('a:ea'))
        if ea is None:
            ea = defRPr.makeelement(qn('a:ea'), {})
            defRPr.append(ea)
        ea.set('typeface', name)
        # 빈 문단의 줄 높이를 좌우하는 endParaRPr 크기도 축소
        endpr = p._p.find(qn('a:endParaRPr'))
        if endpr is None:
            endpr = p._p.makeelement(qn('a:endParaRPr'), {})
            p._p.append(endpr)
        endpr.set('sz', sz)
    except Exception:
        pass


def _compact_grid(table, font_pt, row_h, kind="grid", has_header=False, header_rows=None):
    """표 글자크기·행높이 통일 + Bold/Light.
       - grid: 줄바꿈 OFF(한 줄)·통일 행높이.
       - label_value: 내용 칸(ci=1)은 줄바꿈 ON + 행높이 가변(row_h가 리스트면 행별 높이).
       ★헤더행·구분열(label_value col0)은 Bold, 나머지 Light — 빈 칸 포함 모든 칸에 적용."""
    heights = row_h if isinstance(row_h, (list, tuple)) else None
    hdr_n = header_rows if header_rows is not None else (1 if has_header else 0)
    for ri, row in enumerate(table.rows):
        if heights is not None:
            row.height = Inches(heights[ri] if ri < len(heights) else heights[-1])
        elif hdr_n and ri < hdr_n:
            row.height = Cm(0.6)        # ★표 헤더행(다단이면 여러 줄)은 0.6cm
        else:
            row.height = Inches(row_h)
        for ci, cell in enumerate(row.cells):
            is_header_cell = (ri < hdr_n)   # ★다단 헤더: 첫 hdr_n 행 모두 헤더(Bold)
            is_bold = is_header_cell or (kind == "label_value" and ci == 0)
            name = _FONT_BOLD if is_bold else _FONT_LIGHT
            cell.margin_left = cell.margin_right = Inches(0.04)
            cell.margin_top = cell.margin_bottom = Inches(0.02 if kind == "label_value" else 0)
            tf = cell.text_frame
            # label_value 내용칸(긴 약정문)만 줄바꿈 허용. 세로정렬은 전부 가운데 통일(사용자 지시)
            if kind == "label_value" and ci == 1 and not is_header_cell:
                tf.word_wrap = True
            else:
                tf.word_wrap = False
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf.paragraphs:
                # ★grid 셀 정렬 일관 통일(가운데) — 템플릿 빈칸 상속으로 인한 뒤죽박죽 정렬 방지
                if kind == "grid":
                    p.alignment = PP_ALIGN.CENTER
                _set_para_default_font(p, font_pt, name)   # 빈 칸 포함
                for run in p.runs:                          # 채워진 칸: 크기·글꼴 통일(헤더 Bold 보장)
                    run.font.size = Pt(font_pt)
                    run.font.name = name


_RED = RGBColor(0xC0, 0x00, 0x00)


def _cell_is_red(txt, red_set):
    """셀 텍스트가 원본 빨간 글씨와 매칭되면 True.
       ★셀 텍스트가 '빨간 주석 안에 포함'된 경우(예: 검정 '루시드' ⊂ 빨간 주석)는 제외.
         셀 == 빨간조각, 또는 빨간조각이 셀의 큰 부분일 때만 빨강 처리."""
    t = (txt or "").strip()
    if len(t) < 2 or not red_set:
        return False
    for r in red_set:
        r = (r or "").strip()
        if len(r) < 2:
            continue
        if r == t:
            return True
        # 빨간 조각이 셀 안에 있고, 셀 길이의 절반 이상을 차지하면 강조로 인정
        if r in t and len(r) >= max(2, len(t) * 0.5):
            return True
    return False


def _set_cell_red(cell):
    """셀 글자만 빨강 (원본은 '빨간 글씨' — 테두리 금지). 빨간 테두리는 _set_cell_red_border 별도."""
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = _RED


def _set_cell_red_border(cell, sides=("L", "R", "T", "B")):
    """셀의 지정한 변만 빨간 테두리 1pt. sides 예: ('T','B') (위·아래만).
       ★'본건' 행처럼 '행 바깥쪽만' 빨갛게 하려면 셀 위치별로 변을 골라 호출
         (모든 셀 T/B, 첫 셀 +L, 끝 셀 +R → 내부 세로선은 빨갛지 않음)."""
    # ★OOXML tcPr 자식 순서 규칙: a:lnL, a:lnR, a:lnT, a:lnB … 가 '채우기(solidFill)'보다 앞이어야 함.
    #   고정 index(insert(i))로 넣으면 셀에 fill이 이미 있을 때 테두리가 fill '뒤'로 들어가 순서가
    #   깨지고 PowerPoint가 무시함(테두리가 안 그려지던 근본 원인). → 항상 올바른 위치에 삽입.
    order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    want = {("a:ln" + s) for s in sides}
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in order:
        if tag not in want:
            continue
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
        ln = tcPr.makeelement(qn(tag), {"w": "12700", "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        sf = ln.makeelement(qn("a:solidFill"), {})
        sf.append(sf.makeelement(qn("a:srgbClr"), {"val": "C00000"}))
        ln.append(sf)
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
        # 삽입 위치 = 순서상 자기 앞에 와야 하는 ln 태그 중 이미 존재하는 개수(그 뒤, fill 앞)
        idx = sum(1 for pre in order[:order.index(tag)] if tcPr.find(qn(pre)) is not None)
        tcPr.insert(idx, ln)


def _grid_font(ncol):
    """표 글자크기 — 사용자 지시: 모든 표 10.5pt 통일."""
    return 10.5


def _rowh(font_pt):
    return font_pt * 1.32 / 72.0 + 0.02   # 한 줄 내용행 높이(in) — 최대한 축소(사진 공간 확보)


def _set_one_cell_fill_alpha(cell, hex_color, alpha_pct):
    """셀 1개 채움을 hex_color + 투명도(alpha_pct% 불투명)로."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    sf = tcPr.makeelement(qn("a:solidFill"), {})
    clr = sf.makeelement(qn("a:srgbClr"), {"val": hex_color})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))}))
    sf.append(clr)
    ins = 0
    for i, ch in enumerate(tcPr):
        if ch.tag.endswith(("lnL", "lnR", "lnT", "lnB")):
            ins = i + 1
    tcPr.insert(ins, sf)


def _content_col_widths(data, ncol, W, font_pt):
    """★열 너비를 내용 길이에 비례 배분(균등분할 금지). 짧은 칸은 좁게, 긴 칸은 넓게.
       전체 폭 W에 맞춰 스케일(짧은 열의 상대적 좁음은 유지)."""
    char_w = font_pt * 0.78 / 72.0
    maxlen = [1] * ncol
    for row in data:
        for c in range(ncol):
            txt = str(row[c] if c < len(row) else "")
            ll = max((len(s) for s in txt.split("\n")), default=0)
            if ll > maxlen[c]:
                maxlen[c] = ll
    raw = [max(0.45, maxlen[c] * char_w + 0.14) for c in range(ncol)]
    tot = sum(raw)
    scale = W / tot if tot > 0 else 1.0
    return [w * scale for w in raw]


_CORP_KEYS = ("회사명", "업체명", "대표자", "사업자번호", "법인등록번호",
              "설립일", "설립일자", "결산월", "결산일", "법인구분", "업종명", "주소", "본점소재지")


def _is_corp_info(tdef):
    """기업개요(기본정보) label_value 표인지 — 4열(구분|내용|구분|내용)로 바꿀 대상."""
    if (tdef.get("kind") or "") != "label_value":
        return False
    labels = [str((r[0] if r else "") or "") for r in (tdef.get("rows") or [])]
    hits = sum(1 for L in labels if any(k in L for k in _CORP_KEYS))
    return hits >= 3


def _corp_to_4col(tdef):
    """기업개요 label_value(세로 나열) → 4열 grid(구분|내용|구분|내용). 연속 2행을 좌/우로."""
    rows = [[str((r[0] if len(r) > 0 else "") or ""), str((r[1] if len(r) > 1 else "") or "")]
            for r in (tdef.get("rows") or [])]
    rows = [r for r in rows if r[0] or r[1]]
    out = []
    for i in range(0, len(rows), 2):
        left = rows[i]
        right = rows[i + 1] if i + 1 < len(rows) else ["", ""]
        out.append([left[0], left[1], right[0], right[1]])
    return {"title": tdef.get("title", ""), "kind": "grid",
            "header": ["구분", "내용", "구분", "내용"], "rows": out}


def _parse_tdef(tdef):
    """LLM 표 정의 → (kind, header, body_rows, ncol)."""
    kind = (tdef.get("kind") or "").strip()
    if kind == "label_value":
        body = [[str((r[0] if len(r) > 0 else "") or ""),
                 str((r[1] if len(r) > 1 else "") or "")] for r in (tdef.get("rows") or [])]
        body = [r for r in body if r[0] or r[1]]
        # ★옆으로 된 표는 맨 위에 네이비 「구분 | 내용」 헤더행 추가(분할 시 매 장 반복)
        return "label_value", ["구분", "내용"], body, 2
    header = [str(h or "") for h in (tdef.get("header") or [])]
    body = [[str(c or "") for c in r] for r in (tdef.get("rows") or [])]
    ncol = max([len(header)] + [len(r) for r in body]) if (header or body) else 0
    return "grid", header, body, ncol


def _set_header_fill_alpha(table, hex_color, alpha_pct, header_rows=1):
    """헤더행(0..header_rows-1) 셀 채움을 hex_color + 투명도(alpha_pct% 불투명)로. 글자=검정.
       ★중첩표(표 안의 표)의 헤더는 다단(2줄)이어도 모두 같은 하늘 65% — 둘째 행(항목/금액)이
         남색으로 남던 문제 해결(사용자: 표안의표 헤더는 항상 하늘 65%)."""
    for ri in range(min(header_rows, len(table.rows))):
        for cell in table.rows[ri].cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
                for el in tcPr.findall(qn(tag)):
                    tcPr.remove(el)
            sf = tcPr.makeelement(qn("a:solidFill"), {})
            clr = sf.makeelement(qn("a:srgbClr"), {"val": hex_color})
            clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))}))
            sf.append(clr)
            # 테두리(lnL/R/T/B) 뒤, fill은 그 다음에 와야 하므로 적절히 삽입
            ins = 0
            for i, ch in enumerate(tcPr):
                if ch.tag.endswith(("lnL", "lnR", "lnT", "lnB")):
                    ins = i + 1
            tcPr.insert(ins, sf)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor(0, 0, 0)


def _split_grouped_gubun(header, rows):
    """grid 구분(col0)이 'MAIN (SUB)' 형태로 같은 MAIN이 2행 이상 연속이면, 구분을
       [MAIN | SUB] 2열로 분리(MAIN 세로병합·SUB 없는 행은 구분 가로병합 대상).
       예: '인출금액 (일시)'/'인출금액 (한도)' → 구분='인출금액'(병합), 부구분='일시'/'한도'.
       반환 (new_header, new_rows, split). split=False면 원본 유지."""
    if not rows or not header or len(header) < 2:
        return header, rows, False
    pat = re.compile(r"^(.+?)\s*\(([^()]+)\)\s*$")
    mains = []
    for r in rows:
        c0 = str((r[0] if r else "") or "").strip()
        m = pat.match(c0)
        mains.append(m.group(1).strip() if m else None)
    if not any(mains[i] and mains[i] == mains[i + 1] for i in range(len(mains) - 1)):
        return header, rows, False
    grouped = set()
    for i in range(len(mains)):
        if mains[i] and ((i > 0 and mains[i - 1] == mains[i])
                         or (i + 1 < len(mains) and mains[i + 1] == mains[i])):
            grouped.add(i)
    new_header = [header[0], ""] + list(header[1:])
    new_rows = []
    for i, r in enumerate(rows):
        c0 = str((r[0] if r else "") or "").strip()
        rest = list(r[1:])
        if i in grouped:
            m = pat.match(c0)
            new_rows.append([m.group(1).strip(), m.group(2).strip()] + rest)
        else:
            new_rows.append([c0, ""] + rest)
    return new_header, new_rows, True


def _render_table_chunk(slide, kind, header, rows, ncol, L, T, W, font_pt, row_h, red_counts=None,
                        hdr_fill_hex=None, hdr_alpha=None, anchor_rows=None, is_last_chunk=True,
                        fill_set=None):
    """헤더(있으면)+행들을 렌더 → 높이(in).
       row_h: 스칼라(그리드 통일 행높이) 또는 본문행별 높이 리스트(label_value 가변).
       red_set: 원본 빨간 글씨 집합 — 매칭 셀은 빨간 글씨+빨간 테두리.
       hdr_fill_hex/hdr_alpha: 헤더행 특수 색(중첩표용 3E95BE+투명도)."""
    # ★구분 다단('MAIN (SUB)' 연속 동일 MAIN) → [MAIN|SUB] 2열 분리(MAIN 세로병합). 예 인출금액 일시/한도
    gsub = False
    if kind != "label_value" and header:
        header, rows, gsub = _split_grouped_gubun(header, rows)
        if gsub:
            ncol = len(header)
        elif (len(header) >= 3 and "구분" in str(header[0] or "")
              and not str(header[1] or "").strip()):
            gsub = True   # 이미 2열 구분(구분|빈 부구분헤더, 예 토지확보) → 헤더/빈SUB행 가로병합 적용
    # ★다단(그룹) 헤더(시세표 등): 헤더를 2줄(상위 그룹 + 하위 세부)로 구성
    hdr_grp = None if gsub else (_header_groups(header) if (kind != "label_value" and header) else None)
    if hdr_grp:
        _g_row0, _g_row1, _g_vmerge, _g_hmerges = hdr_grp
        header_list = [_g_row0, _g_row1]
    else:
        header_list = [header] if header else []
    n_hdr = len(header_list)
    data = header_list + rows
    if not data or ncol == 0:
        return 0.0
    nrow = len(data)
    if isinstance(row_h, (list, tuple)):
        hdr_h = _rowh(font_pt)
        row_heights = [hdr_h] * n_hdr + list(row_h)
        height = sum(row_heights)
        row_h = row_heights        # 아래 _compact_grid에 행별 높이 전달
    else:
        height = nrow * row_h
    gf = slide.shapes.add_table(nrow, ncol, Inches(L), Inches(T), Inches(W), Inches(height))
    t = gf.table
    if kind == "label_value":
        lab_w = min(1.7, W * 0.24)
        t.columns[0].width = Inches(lab_w)
        t.columns[1].width = Inches(W - lab_w)
        hdr_rows = 1 if header else 0
        for ri, pair in enumerate(data):
            _replace_text_keep_runs(t.cell(ri, 0).text_frame, pair[0])
            vtf = t.cell(ri, 1).text_frame
            _replace_text_keep_runs(vtf, pair[1])
            if ri >= hdr_rows:   # 헤더행(구분|내용)은 가운데, 값행은 모든 줄 좌측정렬 통일
                for p in vtf.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                _bulletize_text_frame(vtf)   # 값 셀 안 ▶/•/- → 진짜 글머리 기호
        # 헤더행=네이비+흰색Bold, 구분열=회색Bold (분할 청크마다 헤더 반복)
        style_table(gf, has_header=bool(header), label_cols=(0,),
                    header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])
    else:
        widths = _content_col_widths(data, ncol, W, font_pt)
        for c in range(ncol):
            t.columns[c].width = Inches(widths[c])
        for ri, row in enumerate(data):
            for ci in range(ncol):
                ctf = t.cell(ri, ci).text_frame
                _replace_text_keep_runs(
                    ctf, str((row[ci] if ci < len(row) else "") or ""))
                if ri >= n_hdr:   # 데이터 셀 안 ▶/•/- → 진짜 글머리 기호(헤더 제외)
                    _bulletize_text_frame(ctf)
        style_table(gf, has_header=bool(header), label_cols=(), header_fill=PALETTE["navy_dark"])
        # ★다단 헤더: 둘째 헤더행(세부)도 네이비+흰색 Bold, 그룹 병합(상위 가로 / 비그룹 세로)
        if hdr_grp:
            for ci in range(ncol):
                cl = t.cell(1, ci)
                cl.fill.solid(); cl.fill.fore_color.rgb = PALETTE["navy_dark"]
                for p in cl.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.name = _FONT_BOLD
                        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for (s, e, _lab) in _g_hmerges:        # 상위행 그룹 가로병합
                try:
                    t.cell(0, s).merge(t.cell(0, e))
                except Exception:
                    pass
            for c in _g_vmerge:                    # 비그룹 칸 세로병합(상·하위 헤더)
                try:
                    t.cell(0, c).merge(t.cell(1, c))
                except Exception:
                    pass
        # ★기업개요 4열(구분|내용|구분|내용)은 세로병합 금지 — 각 행이 독립(업종명이 주소 행으로
        #   번지는 버그 방지). 대신 우측 쌍이 빈 행(주소 등)은 값(col1)을 우측 끝까지 가로 병합(전폭).
        is_corp4 = (ncol == 4 and len(header) == 4
                    and "구분" in str(header[0]) and "구분" in str(header[2]))
        if not is_corp4:
            _merge_vertical_runs(t, data, ncol, header_rows=n_hdr)  # 반복값 세로병합
        _merge_total_rows(t, data, ncol)                                     # 합계행 가로병합
        if gsub:
            # 구분 다단: 헤더의 구분은 부구분칸까지 가로병합 / SUB 빈 데이터행은 구분이 부구분칸까지 가로병합
            try:
                t.cell(0, 0).merge(t.cell(0, 1))
            except Exception:
                pass
            for ri in range(n_hdr, nrow):
                sub = str((data[ri][1] if ri < len(data) and len(data[ri]) > 1 else "") or "").strip()
                if not sub:
                    try:
                        t.cell(ri, 0).merge(t.cell(ri, 1))
                    except Exception:
                        pass
        if is_corp4:
            for ri in range(n_hdr, nrow):
                r = data[ri]
                c2 = str((r[2] if len(r) > 2 else "") or "").strip()
                c3 = str((r[3] if len(r) > 3 else "") or "").strip()
                if not c2 and not c3 and str((r[1] if len(r) > 1 else "") or "").strip():
                    try:
                        t.cell(ri, 1).merge(t.cell(ri, 3))   # 주소 등 전폭 값
                    except Exception:
                        pass
        # ★사업수지(구분|세부항목|세대수/내역|…): 세대수/내역이 빈 일반행(발코니 확장·상가)은
        #   세부항목+세대수/내역을 가로 병합(원본처럼 "발코니 확장"이 두 칸 차지). 합계행은 위에서 처리됨.
        is_suji = (len(header) >= 3 and "구분" in str(header[0])
                   and ("세부" in str(header[1]) or "항목" in str(header[1])))
        if is_suji:
            _sub, _grand = _classify_total_rows(data, ncol)
            _tot = _sub | _grand
            for ri in range(n_hdr, nrow):
                if ri in _tot:
                    continue
                r = data[ri]
                c0 = str((r[0] if len(r) > 0 else "") or "").strip()
                c1 = str((r[1] if len(r) > 1 else "") or "").strip()
                c2 = str((r[2] if len(r) > 2 else "") or "").strip()
                if c0 and not c1 and not c2:
                    # 본건 시행이익(A-B)·Equity 반영 시 = 라벨이 구분+세부항목+세대수내역에 걸침
                    try:
                        t.cell(ri, 0).merge(t.cell(ri, 2))
                    except Exception:
                        pass
                elif c1 and not c2:        # 발코니 확장·상가 = 세부항목+세대수내역 가로병합
                    try:
                        t.cell(ri, 1).merge(t.cell(ri, 2))
                    except Exception:
                        pass
    _compact_grid(t, font_pt, row_h, kind=kind, has_header=bool(header), header_rows=n_hdr)
    if header and hdr_fill_hex:   # 중첩표 헤더 = 밝은색 + 투명도(다단이면 모든 헤더행)
        _set_header_fill_alpha(t, hdr_fill_hex, hdr_alpha if hdr_alpha is not None else 35,
                               header_rows=n_hdr)
    hdr_rows = n_hdr
    # ★표 안의 표 앵커 행은 글씨를 위에 두고 그 아래 grid를 얹으므로 위 정렬(나머지는 가운데)
    if anchor_rows:
        for li in anchor_rows:
            r = hdr_rows + li
            if r < nrow:
                t.cell(r, 1).vertical_anchor = MSO_ANCHOR.TOP
    # 구분(라벨)열 인덱스 — 아래 합계 음영 후 다시 F2F2F2로 덮어씀(구분칸은 항상 밝은 회색)
    is_content = (kind != "label_value")
    gubun_cols = ([c for c in range(ncol) if c < len(header) and "구분" in str(header[c] or "")]
                  if (is_content and header) else [])
    if gsub and 0 in gubun_cols and 1 not in gubun_cols:
        gubun_cols.append(1)   # 구분 다단: 부구분칸(일시/한도 등)도 구분열 회색으로
    # ★합계행 음영(원본 규칙, 사용자 재확인):
    #   - 표 안의 표(중첩 grid, hdr_fill_hex 지정) → 합계 색칠 '무조건 안 함'
    #   - 그 외 표: 표의 '맨 밑' 합계/총합계(마지막 청크의 최하단 합계행) 1개만 하늘 민트(3E95BE 65%).
    #     중간에 끼인 소계·합계·총합계는 전부 한 톤 어두운 회색(D9D9D9).
    is_nested = bool(hdr_fill_hex)
    if not is_nested:
        subtotal_rows, grandtotal_rows = _classify_total_rows(data, ncol)
        total_rows = subtotal_rows | grandtotal_rows
        bottom_total = max(total_rows) if total_rows else -1
        for ri in range(hdr_rows, nrow):
            if ri not in total_rows:
                continue
            if ri == bottom_total and is_last_chunk:
                for ci in range(ncol):
                    _set_one_cell_fill_alpha(t.cell(ri, ci), "3E95BE", 35)   # 맨 밑 합계 = 민트 65%
            else:
                for ci in range(ncol):
                    cl = t.cell(ri, ci)
                    cl.fill.solid()
                    cl.fill.fore_color.rgb = PALETTE["gray"]                 # 중간 소계/합계 = D9D9D9
    # ★구분열은 항상 밝은 회색 F2F2F2(합계 음영보다 우선). 단 '합계행'(total로 분류된 행)은 합계 음영 유지.
    #   행이 합계가 아니면 라벨에 '총'이 들어가도(자산총계·부채총계 등 항목명) 전부 F2F2F2 — 일관되게.
    #   ★표 안의 표(중첩 grid)는 예외 — 헤더만 하늘 65%, 구분열 데이터칸은 색칠 안 함(사용자 지시).
    if gubun_cols and not is_nested:
        _g_sub, _g_grand = _classify_total_rows(data, ncol)
        _g_totrows = _g_sub | _g_grand
        for ri in range(hdr_rows, nrow):
            if ri in _g_totrows:        # 진짜 합계행만 음영 유지
                continue
            for gc in gubun_cols:
                cl = t.cell(ri, gc)
                cl.fill.solid()
                cl.fill.fore_color.rgb = PALETTE["label_gray"]
    # ★원본 '포인트 색칠' 칸(시세표 공급평단가 주황 등) → 하늘색 65%로 재현(약속한 포인트색)
    #   단, '구분/라벨 열'은 항상 밝은 회색이어야 하므로 제외(자금용도 등이 민트되던 문제).
    if fill_set:
        skip_cols = set()
        if kind == "label_value":
            skip_cols = {0}
        elif header:
            skip_cols = {c for c in range(ncol) if c < len(header) and "구분" in str(header[c] or "")}
        for ri in range(hdr_rows, nrow):
            for ci in range(ncol):
                if ci in skip_cols:
                    continue
                ct = t.cell(ri, ci).text.strip()
                if ct and ct in fill_set:
                    _set_one_cell_fill_alpha(t.cell(ri, ci), "3E95BE", 35)
    # ★원본 빨간 글씨 재현: '모든 동일텍스트 칸이 빨강이었을 때만' 칠함(과다 방지) + 빨간 테두리
    if red_counts:
        body_counts = Counter()
        for ri in range(hdr_rows, nrow):
            for ci in range(ncol):
                ct = t.cell(ri, ci).text.strip()
                if ct:
                    body_counts[ct] += 1
        for ri in range(hdr_rows, nrow):
            for ci in range(ncol):
                cell = t.cell(ri, ci)
                ct = cell.text.strip()
                if ct and red_counts.get(ct, 0) > 0 and red_counts[ct] >= body_counts[ct]:
                    _set_cell_red(cell)
    # ★원본에서 '본건' 행 외곽이 빨간 강조(시세표) → 행 '바깥쪽만' 빨간 테두리 1pt
    #   (모든 셀 위·아래, 첫 셀 +왼쪽, 끝 셀 +오른쪽 — 내부 세로선은 빨갛지 않음)
    for ri in range(hdr_rows, nrow):
        c0 = str((data[ri][0] if ri < len(data) and data[ri] else "") or "").strip()
        if c0 in ("본건", "본 건"):
            for ci in range(ncol):
                _sides = ["T", "B"]
                if ci == 0:
                    _sides.append("L")
                if ci == ncol - 1:
                    _sides.append("R")
                _set_cell_red_border(t.cell(ri, ci), _sides)
            # ★공유하는 위·아래 선이 옆 행 회색선에 덮이지 않도록, 윗 행 아래·아랫 행 위 테두리도 빨강
            if ri > 0:
                for ci in range(ncol):
                    _set_cell_red_border(t.cell(ri - 1, ci), ["B"])
            if ri + 1 < nrow:
                for ci in range(ncol):
                    _set_cell_red_border(t.cell(ri + 1, ci), ["T"])
    # ★원본에서 '소유자' 열의 '루시드'(신탁 미이전 필지)는 빨간 글씨로 강조(토지개요)
    own_col = next((c for c in range(ncol)
                    if c < len(header) and "소유자" in str(header[c] or "")
                    and "위탁" not in str(header[c] or "")), None)
    if own_col is not None:
        for ri in range(hdr_rows, nrow):
            cl = t.cell(ri, own_col)
            if cl.text.strip() == "루시드":
                _set_cell_red(cl)
    # ★원본에서 '에쿼티 금액' 열 전체가 빨간 겉테두리(Equity 투입현황표) → 그 열 바깥만 빨강
    #   (데이터 행만, 헤더 제외. 열 좌·우 + 첫 데이터 칸 위 + 끝 칸 아래. 왼쪽 공유선은 옆 칸 우테두리도 빨강)
    eq_col = next((c for c in range(ncol)
                   if c < len(header) and "에쿼티" in str(header[c] or "")), None)
    if eq_col is not None:
        for ri in range(hdr_rows, nrow):
            _sides = ["L", "R"]
            if ri == hdr_rows:
                _sides.append("T")
            if ri == nrow - 1:
                _sides.append("B")
            _set_cell_red_border(t.cell(ri, eq_col), _sides)
            if eq_col > 0:                       # 왼쪽 공유선이 회색에 덮이지 않게 옆 칸 우테두리도 빨강
                _set_cell_red_border(t.cell(ri, eq_col - 1), ["R"])
        # ★박스 위쪽 모서리 완성: 헤더('에쿼티 금액(원)') 아래선도 빨강(데이터 첫 칸 위와 공유)
        if hdr_rows > 0:
            _set_cell_red_border(t.cell(hdr_rows - 1, eq_col), ["B"])
    # ★원본 정렬 재현(G2) — 모든 병합·색칠 뒤 맨 마지막에 적용(앞서 _compact_grid가 가운데로 리셋하므로).
    #   원본처럼: '숫자/비율 전용' 컬럼 = 오른쪽 / '텍스트 섞인 긴(≥10자)' 컬럼 = 왼쪽 / 그 외(짧은) = 가운데.
    #   (금액·세대수·면적·비율·매매대금 → 오른쪽 / 진행일정·관계법령·단지명·구분 라벨 → 왼쪽)
    if kind != "label_value" and header:
        _NUMERIC = re.compile(r"^[\d,.\-~%()\s원억평㎡천만원/]+$")
        for ci in range(ncol):
            if ci in gubun_cols:   # ★구분/부구분 라벨열은 항상 가운데(원본대로) — 길어도 왼쪽 안 함
                continue
            vals = [str((data[ri][ci] if ci < len(data[ri]) else "") or "").strip()
                    for ri in range(hdr_rows, nrow)]
            vals = [v for v in vals if v]
            if not vals:
                continue
            all_numeric = all(_NUMERIC.match(v) for v in vals)
            longish = any(len(v) >= 10 for v in vals)
            _al = (PP_ALIGN.RIGHT if all_numeric
                   else (PP_ALIGN.LEFT if longish else None))
            if _al is not None:
                for ri in range(hdr_rows, nrow):
                    for p in t.cell(ri, ci).text_frame.paragraphs:
                        p.alignment = _al
    return height


def _place_images_row(slide, imgs, L, top, W, bottom, *, labels=None):
    """남은 본문 공간에 이미지들을 한 줄로 배치(가로 균등, 높이 맞춰 축소). 사용높이(in) 반환.
       labels: 각 이미지 위 작은 라벨(조감도/광역입지 등) 리스트(선택)."""
    imgs = [im for im in (imgs or []) if im.get("data")]
    if not imgs:
        return 0.0
    avail_h = bottom - top
    if avail_h < 0.6:
        return 0.0
    n = len(imgs)
    gap = 0.2
    cell_w = (W - gap * (n - 1)) / n
    lab_h = 0.22 if labels else 0.0
    x = L
    used = 0.0
    for i, im in enumerate(imgs):
        iw, ih = im.get("width", 1), im.get("height", 1)
        scale = min(cell_w / iw, (avail_h - lab_h) / ih)
        w_in, h_in = iw * scale, ih * scale
        if labels and i < len(labels) and labels[i]:
            _add_textbox(slide, x, top, cell_w, lab_h, labels[i],
                         size=9, bold=True, align=PP_ALIGN.CENTER, color=PALETTE["navy_dark"])
        px = x + (cell_w - w_in) / 2
        try:
            slide.shapes.add_picture(io.BytesIO(im["data"]), Inches(px), Inches(top + lab_h),
                                     Inches(w_in), Inches(h_in))
        except Exception:
            pass
        used = max(used, lab_h + h_in)
        x += cell_w + gap
    return used


def _place_images_grid(slide, imgs, L, top, W, bottom):
    """사진 여러 장(사업지 전경 등): 첫 장 크게(상단) + 나머지 3열 그리드 — 원본 레이아웃 재현."""
    imgs = [im for im in (imgs or []) if im.get("data")]
    if not imgs:
        return 0.0
    if len(imgs) <= 2:
        return _place_images_row(slide, imgs, L, top, W, bottom)
    avail = bottom - top
    gap = 0.14
    rest = imgs[1:]
    cols = 3
    grows = (len(rest) + cols - 1) // cols
    big_h = min(avail * 0.46, W * 0.48)
    _place_images_row(slide, imgs[:1], L, top, W, top + big_h)
    gy = top + big_h + gap
    if grows > 0:
        row_h = (bottom - gy - gap * (grows - 1)) / grows
        if row_h > 0.5:
            for r in range(grows):
                chunk = rest[r * cols:(r + 1) * cols]
                _place_images_row(slide, chunk, L, gy + r * (row_h + gap), W,
                                  gy + r * (row_h + gap) + row_h)
    return avail


def _img_labels_for(subtitle, n):
    """소제목으로 이미지 라벨 추정(조감도/광역입지/현장사진/위치도 등)."""
    s = subtitle or ""
    if "승인" in s or "공문" in s or "인허가" in s:
        base = ["사업계획승인서", "첨부"]
    elif "현장" in s:
        base = ["현장사진", "위치도"]
    elif "건축" in s or "조감" in s:
        base = ["조감도", "광역입지"]
    elif "사업 개요" in s or "사업개요" in s:
        base = ["조감도", "위치도"]
    elif "입지" in s or "시세" in s or "분양" in s or "토지" in s:
        base = ["위치도", "조감도"]
    else:
        base = ["참고 이미지", "참고 이미지"]
    return [base[i] if i < len(base) else base[-1] for i in range(n)]


def _place_images_col(slide, imgs, L, top, W, bottom, *, labels=None):
    """오른쪽에 '구분 | 내용' 2열 표를 만들고 '내용' 칸에 사진을 얹는다(따로따로 사진 X).
       구분=라벨(조감도/광역입지), 내용=이미지 — 원본 '표 안 사진' 형식(사용자 지시)."""
    imgs = [im for im in (imgs or []) if im.get("data")]
    if not imgs:
        return 0.0
    avail_h = bottom - top
    if avail_h < 0.7:
        return 0.0
    n = len(imgs)
    hdr_h = 0.30
    row_h = (avail_h - hdr_h) / n
    lab_w = 0.95
    content_w = W - lab_w
    gf = slide.shapes.add_table(1 + n, 2, Inches(L), Inches(top), Inches(W), Inches(avail_h))
    tb = gf.table
    tb.columns[0].width = Inches(lab_w)
    tb.columns[1].width = Inches(content_w)
    tb.cell(0, 0).text = "구분"
    tb.cell(0, 1).text = "내용"
    for i in range(n):
        tb.cell(1 + i, 0).text = (labels[i] if labels and i < len(labels) else "참고 이미지")
    # 헤더=네이비/흰색, 구분열(col0)=밝은 회색
    style_table(gf, has_header=True, label_cols=(0,),
                header_fill=PALETTE["navy_dark"], label_fill=PALETTE["label_gray"])
    tb.rows[0].height = Inches(hdr_h)
    for i in range(n):
        tb.rows[1 + i].height = Inches(row_h)
    for ci in (0, 1):
        for p in tb.cell(0, ci).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
    for i in range(n):
        c = tb.cell(1 + i, 0)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(10)
        # '내용' 칸에 사진 얹기(셀 안쪽 여백 0.06)
        im = imgs[i]
        cw, ch = content_w - 0.12, row_h - 0.12
        iw, ih = im.get("width", 1), im.get("height", 1)
        scale = min(cw / iw, ch / ih) if iw and ih else 1.0
        w_in, h_in = iw * scale, ih * scale
        px = L + lab_w + (content_w - w_in) / 2
        py = top + hdr_h + i * row_h + (row_h - h_in) / 2
        try:
            slide.shapes.add_picture(io.BytesIO(im["data"]), Inches(px), Inches(py),
                                     Inches(w_in), Inches(h_in))
        except Exception:
            pass
    return avail_h


_GAP = 0.22       # 표 간 간격
_LABEL_H = 0.34   # 표 미니라벨 높이


def build_structured_slide(prs, struct: dict, *, business_name: str = "",
                           section_label: str = None, subtitle: str = None,
                           images=None, red_texts=None, underline_texts=None, fill_texts=None):
    """LLM 구조화 결과 → 1개 이상의 본문 슬라이드. 내용이 길면 헤더 반복하며 "(i/n)"로 분할."""
    section_label = (section_label if section_label is not None
                     else struct.get("section_label") or "").strip()
    subtitle = (subtitle if subtitle is not None
                else struct.get("subtitle") or "").strip()
    intro = (struct.get("intro") or "").strip()
    bullets = [b for b in (struct.get("bullets") or []) if str(b).strip()]
    tables = struct.get("tables") or []
    # ★기업개요(시행사·시공사 기본정보) label_value → 4열(구분|내용|구분|내용)로(원본처럼)
    tables = [_corp_to_4col(t) if _is_corp_info(t) else t for t in tables]
    source = (struct.get("source") or "").strip()

    # ★표 안의 표: 앵커 라벨 → grid 파싱 (해당 label_value 행 내용칸에 grid를 얹음)
    nested_grids = struct.get("_nested_grids") or []
    nested_parsed = []   # (anchor_label, (kind,header,body,ncol), note, unit)
    for item in nested_grids:
        anchor, gdef = item[0], item[1]
        note = item[2] if len(item) > 2 else ""
        gk, gh, gb, gn = _parse_tdef(gdef)
        # grid title의 '(단위 …)' → 표 위 캡션으로(원본: Cash 표 위 '(단위: 백만원)')
        _gt = str((gdef.get("title") if isinstance(gdef, dict) else "") or "")
        _ui = _gt.find("(단위")
        unit = _gt[_ui:(_gt.find(")", _ui) + 1 or len(_gt))] if _ui != -1 else ""
        if gn > 0 and gb:
            nested_parsed.append((anchor, (gk, gh, gb, gn), note, unit))

    def _grid_render_h(gh, gb, gn):
        """중첩 grid 렌더 높이 — Cash-in/Cash-out 2행 헤더는 헤더 2줄로 계산."""
        hdr_rows = 1 if gh else 0
        if gh and any("Cash" in str(h) for h in gh):
            hdr_rows = 2
        return (hdr_rows + len(gb)) * _rowh(_grid_font(gn)) + 0.10

    def _nested_for(label):
        for a, parsed, note, unit in nested_parsed:
            if a and a in str(label):
                return parsed
        return None

    def _nested_note_for(label):
        for a, parsed, note, unit in nested_parsed:
            if a and a in str(label):
                return note
        return ""

    def _nested_unit_for(label):
        for a, parsed, note, unit in nested_parsed:
            if a and a in str(label):
                return unit
        return ""

    red_set = set(red_texts or [])             # 산문(글상자) 부분 빨강용
    red_counts = Counter(red_texts or [])      # 표 셀 빨강용(모든 동일칸 빨강일 때만)
    ul_set = set(underline_texts or [])        # 산문(글상자) 부분 밑줄용
    fill_set = set(fill_texts or [])           # 원본 포인트 색칠 칸 → 하늘 65%
    # ★각주 줄이 원본 빨강 문구(주2)·주3) 등)와 공백 무시하고 겹치면 그 줄 전체를 빨강 처리.
    #   red_set은 부분문자열 매칭이라 '루시드 2 필지'(원본 빨강) vs '루시드 2필지'(각주 원문)
    #   띄어쓰기 차이로 매칭이 안 되던 것 보완(주3) 담보신탁계약·주2) Equity 초기사업비).
    _red_despaced = [re.sub(r"\s+", "", str(r)) for r in red_set
                     if str(r).lstrip().startswith("주")]

    def _maybe_red_note(line):
        d = re.sub(r"\s+", "", str(line))
        if not d:
            return
        for rp in _red_despaced:
            if len(rp) >= 6 and (rp in d or d in rp):
                red_set.add(str(line))
                return
    # ★주1)2)3) 등 각주는 본문이 아니라 맨 아래로(원본처럼). 그 외 설명만 본문 불릿.
    def _is_note(b):
        s = str(b).lstrip()
        return s.startswith(("주", "*", "※")) or s[:2] in ("1)", "2)", "3)", "4)", "5)")
    notes = [b for b in bullets if _is_note(b)]
    body_bullets = [b for b in bullets if b not in notes]
    btext = "\n".join(str(b) for b in body_bullets) if body_bullets else ""

    # ★출처 필드에 섞여 들어온 주N)/※ 각주 분리 — 출처 글상자는 '출처'만, 주N)는 표 밑 각주로
    def _looks_note(ln):
        s = ln.lstrip("(").lstrip()
        return s.startswith(("주", "*", "※")) and (
            s[:1] in ("*", "※") or (len(s) > 1 and s[1].isdigit()) or s.startswith("주 "))
    if source:
        _slines = [x.strip() for x in source.replace(" / ", "\n").split("\n") if x.strip()]
        _src_keep = [x for x in _slines if not _looks_note(x)]
        _src_notes = [x for x in _slines if _looks_note(x)]
        source = " ".join(_src_keep).strip()
        notes = notes + _src_notes
    for _ln in notes:
        _maybe_red_note(_ln)
    notes_text = "\n".join(notes) if notes else ""

    # ── 본문 사진 배치 규칙(사람 제안서 실측) ──
    #   • 원본이 '표 안 사진'(조감도·광역위치도 등 라벨 이미지) → 표박스(우측, 1:1) = side_box
    #   • 원본이 '그냥 사진'(담보토지·인근시장 지도) → 표 위 전폭 바 이미지   = top_bare
    #   • 표 없는 부록(현장사진·승인서) → 표 없이 크게 중앙
    big_imgs = _big_images(images)
    has_tbl = any(_parse_tdef(t)[3] > 0 and _parse_tdef(t)[2] for t in tables)
    labeled_img = bool(big_imgs) and any(k in subtitle for k in
                                         ("건축", "입지", "조감", "위치", "사업 개요", "사업개요"))
    side_box = bool(big_imgs) and has_tbl and labeled_img
    top_bare = bool(big_imgs) and has_tbl and not labeled_img
    _IMG_W, _IMG_GAP_LR = 4.0, 0.25
    _sidebox_tw = _TBL_W - _IMG_W - _IMG_GAP_LR
    tw = _sidebox_tw if side_box else _TBL_W
    img_col_L = _TBL_L + tw + _IMG_GAP_LR
    _IMG_TOP_H = 2.8 if top_bare else 0.0
    # ★side_box(조감도/광역입지) 이미지는 '첫 plan'(건축개요 사업명·대지위치 표) 옆에 둔다.
    #   → 용도별공급계획(다음 plan)은 전폭 유지(좁아져 이상해지던 문제 해결).
    img_plan_idx = 0

    # ── 표를 청크로 분할 + 라벨에 "(k/m)" 부착 (★우측상단 표식 아님, 표 라벨에 이어 적음) ──
    #   각 표마다 빈 슬라이드 1장 용량 기준으로 m등분 → "(k/m)"이 안정적으로 결정됨.
    blocks = []   # (label, kind, header, rows, ncol, fp, rh)  rh=스칼라(grid)|리스트(label_value)
    for tdef in tables:
        kind, header, body, ncol = _parse_tdef(tdef)
        if ncol == 0 or not body:
            continue
        fp = _grid_font(ncol)
        title = (tdef.get("title") or "").strip()
        has_h = bool(header)
        label_h = _LABEL_H if title else 0.0
        tnotes = [str(x).strip() for x in (tdef.get("_notes") or []) if str(x).strip()]  # 이 표 전용 각주
        for _ln in tnotes:
            _maybe_red_note(_ln)

        if kind == "label_value":
            # ★내용이 긴 약정문 → 행별 높이 가변. 행높이 합이 한 장 용량 넘으면 분할(헤더 반복).
            #   ★앵커 행(주요 대출조건·자금용도)은 grid 높이만큼 키워 그 안에 grid를 얹는다(표 안의 표).
            lab_w = min(1.7, tw * 0.24)
            val_w = tw - lab_w

            def _calc_rhs(fp_):
                out = []
                for r in body:
                    ng = _nested_for(r[0])
                    if ng:
                        _, gh, gb, gn = ng
                        grid_h = _grid_render_h(gh, gb, gn)
                        # ★앵커 행 = (단위캡션)+(grid)+(주N) 각주)+(요약 글씨) — 표 먼저, 글 나중
                        txt_h = (_est_text_height(r[1], val_w - 0.08, fp_) + 0.04
                                 if str(r[1]).strip() else 0.0)
                        nt = _nested_note_for(r[0])
                        note_h = (_est_text_height(nt, val_w - 0.10, 9) + 0.04) if nt.strip() else 0.0
                        unit_h = 0.22 if _nested_unit_for(r[0]) else 0.0
                        out.append(grid_h + txt_h + note_h + unit_h)
                    else:
                        out.append(max(_est_text_height(r[0], lab_w - 0.08, fp_),
                                       _est_text_height(r[1], val_w - 0.08, fp_)) + 0.04)
                return out

            hdr_h = _rowh(fp)
            avail = _BODY_BOTTOM - _INTRO_T - label_h - hdr_h
            rhs_all = _calc_rhs(fp)
            # ★side_box: 기본 10.5 유지. '10.5에서 넘칠 때만' 폰트 축소(입지분석). 사업개요는 10.5 그대로.
            _no_anchor = not any(_nested_for(r[0]) for r in body)
            if side_box and _no_anchor and sum(rhs_all) > avail:
                fp = max(8.0, fp - 2.0)
                hdr_h = _rowh(fp)
                avail = _BODY_BOTTOM - _INTRO_T - label_h - hdr_h
                rhs_all = _calc_rhs(fp)
            chunks, cr, crh, cacc = [], [], [], 0.0
            for r, h in zip(body, rhs_all):
                if cr and cacc + h > avail:
                    chunks.append((cr, crh)); cr, crh, cacc = [], [], 0.0
                cr.append(r); crh.append(h); cacc += h
            if cr:
                chunks.append((cr, crh))
            # ★side_box 단일 청크: 우측 이미지표 높이에 맞춰 행 높이를 늘림(좌우 표 길이 맞춤 — 사업개요).
            if side_box and _no_anchor and len(chunks) == 1:
                _csum = sum(chunks[0][1])
                if 0 < _csum < avail:
                    _sc = avail / _csum
                    chunks[0] = (chunks[0][0], [h * _sc for h in chunks[0][1]])
            m = len(chunks)
            for k, (rows, rhs) in enumerate(chunks):
                lbl = (f"{title}({k + 1}/{m})" if (title and m > 1) else title)
                anchors = [(i, rows[i][0]) for i in range(len(rows)) if _nested_for(rows[i][0])]
                _bn = tnotes if k == m - 1 else []   # 표 각주는 마지막 청크 밑에만
                blocks.append((lbl, kind, header, rows, ncol, fp, rhs, anchors, _bn))
        else:
            rh = _rowh(fp)
            cap_full = max(1, int((_BODY_BOTTOM - _INTRO_T - label_h) / rh) - (1 if has_h else 0))
            # ★상단 전폭 지도(top_bare)가 있는 표는 '첫 청크'를 지도 높이만큼 줄여, 지도+첫 묶음만
            #   1페이지에(예: 토지개요 → 지도 + 터미널 합계까지). 나머지는 다음 페이지로.
            cap_first = cap_full
            if top_bare and _IMG_TOP_H > 0:
                cap_first = max(1, cap_full - int(_IMG_TOP_H / rh) - 1)
            # ★카테고리(소계/합계로 끝나는 묶음)를 페이지 사이에서 쪼개지 않게 — 소계 경계에서 분할.
            _sub, _grand = _classify_total_rows(([header] if has_h else []) + body, ncol)
            _hoff = 1 if has_h else 0
            boundary_after = {ti - _hoff for ti in (_sub | _grand) if 0 <= ti - _hoff < len(body)}
            starts = [0]
            cur_n, last_b = 0, -1
            i = 0
            while i < len(body):
                cur_n += 1
                if i in boundary_after:
                    last_b = i                      # 여기서 끊으면 안전(소계/합계 뒤)
                cap = cap_first if len(starts) == 1 else cap_full   # 첫 청크만 지도분 축소
                if cur_n >= cap and i + 1 < len(body):
                    cut = last_b if last_b >= starts[-1] else i   # 경계 없으면 어쩔 수 없이 여기서
                    starts.append(cut + 1)
                    cur_n, last_b, i = 0, -1, cut + 1
                    continue
                i += 1
            bounds = [(starts[s], starts[s + 1] if s + 1 < len(starts) else len(body))
                      for s in range(len(starts))]
            m = len(bounds)
            for k, (start, _end) in enumerate(bounds):
                chunk = [list(r) for r in body[start:_end]]   # 복사(수정 위해)
                if not chunk:
                    continue
                # ★다음 장으로 넘어간 세로병합값: 첫 행의 빈 칸을 직전 비어있지 않은 값으로 재표기
                #   (A구역처럼 위에서 병합된 라벨이 빈칸으로 시작하지 않도록)
                if start > 0:
                    for c in range(ncol):
                        cell0 = str(chunk[0][c] if c < len(chunk[0]) else "").strip()
                        if cell0:
                            continue
                        for r in range(start - 1, -1, -1):
                            prev = str(body[r][c]) if c < len(body[r]) else ""
                            if prev.strip():
                                while len(chunk[0]) <= c:
                                    chunk[0].append("")
                                chunk[0][c] = prev
                                break
                lbl = (f"{title}({k + 1}/{m})" if (title and m > 1) else title)
                _bn = tnotes if k == m - 1 else []   # 표 각주는 마지막 청크 밑에만
                blocks.append((lbl, kind, header, chunk, ncol, fp, rh, None, _bn))

    # ── 블록을 슬라이드에 채움 (각 청크는 빈 슬라이드 1장에 반드시 들어감) ──
    #   ★제목의 내용(인트로)이 없으면 본문을 인트로 자리(_INTRO_T)부터 시작 = 전체적으로 위로.
    plans = []
    cur = {"intro": intro or None, "items": []}
    top = _INTRO_T
    if intro:
        top = _INTRO_T + _est_text_height(intro, _INTRO_W, 9) + 0.12
    top += _IMG_TOP_H   # 슬1 상단 전폭 사진 공간(top_bare)
    # ★본문 산문(bullets)·각주(주N)는 '표 아래'에 배치(원본: 글이 표 끝에) → 하단 공간 예약
    reserve_bottom = 0.0
    if btext:
        reserve_bottom += _est_text_height(btext, tw, 9) + 0.18
    if notes_text:
        reserve_bottom += _est_text_height(notes_text, _TBL_W, 9) + 0.10
    pack_bottom = _BODY_BOTTOM - reserve_bottom

    def flush():
        nonlocal cur, top
        plans.append(cur)
        cur = {"intro": None, "items": []}
        top = _INTRO_T          # 연속 슬라이드(인트로 없음)는 인트로 자리부터 = 위로

    for blk in blocks:
        lbl, kind, header, rows, ncol, fp, rh, _anchors, _tn = blk
        label_h = _LABEL_H if lbl else 0.0
        _tn_h = (_est_text_height("\n".join(_tn), tw, 9) + 0.10) if _tn else 0.0
        if isinstance(rh, (list, tuple)):
            need = label_h + (_rowh(fp) if header else 0) + sum(rh) + _tn_h + _GAP
        else:
            need = label_h + ((1 if header else 0) + len(rows)) * rh + _tn_h + _GAP
        if cur["items"] and top + need > pack_bottom:
            flush()
        cur["items"].append(blk)
        top += need
    flush()

    # ── 렌더 ──
    n = len(plans)
    for idx, plan in enumerate(plans):
        slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)
        _fill_header(slide, section_label, subtitle, "")   # 소제목은 목차와 동일
        # ★plan별 표 폭: side_box 이미지가 들어가는 plan만 좁게, 나머지는 전폭
        if side_box:
            tw = _sidebox_tw if idx == img_plan_idx else _TBL_W
            img_col_L = _TBL_L + tw + _IMG_GAP_LR
        t = _INTRO_T
        if plan["intro"]:
            # ★제목의 내용 글상자 = Bold (사용자 지시), 원본 빨간 글씨 부분만 빨강
            _, h = _add_textbox(slide, _INTRO_L, _INTRO_T, _INTRO_W, 0.40, plan["intro"],
                                size=9, bold=True, red_set=red_set, ul_set=ul_set)
            t = _INTRO_T + h + 0.12
        # (본문 산문 bullets는 더 이상 표 위에 찍지 않음 — 표 아래로 이동, 아래 렌더 참조)
        # ★'그냥 사진'(지도 등)은 표 위에 전폭으로 (사람 제안서: 담보토지·인근시장)
        #   ★단 시세/분양 비교 페이지는 지도 왼쪽 + 분석문 오른쪽(원본 16번처럼).
        if idx == 0 and top_bare:
            _side_text = (btext if (btext and idx == 0
                          and any(k in str(subtitle) for k in ("시세", "분양", "입지"))) else "")
            if _side_text:
                _map_w = _TBL_W * 0.56
                _place_images_row(slide, big_imgs[:1], _TBL_L, t, _map_w, t + _IMG_TOP_H)
                _add_textbox(slide, _TBL_L + _map_w + 0.2, t, _TBL_W - _map_w - 0.2,
                             _IMG_TOP_H, _side_text, size=9, bullet=True,
                             red_set=red_set, ul_set=ul_set)
                btext = ""   # 지도 옆에 이미 적었으니 표 아래엔 안 찍음
            else:
                _place_images_row(slide, big_imgs[:1], _TBL_L, t, _TBL_W, t + _IMG_TOP_H)
            t += _IMG_TOP_H + 0.12
        tbl_start = t   # 표가 시작되는 y(사진 옆배치 기준)
        for (lbl, kind, header, rows, ncol, fp, rh, anchors, tbl_notes) in plan["items"]:
            if lbl:
                # ★'(단위 …)'는 표 우측 위에 Light 9pt 별도 글상자로(원본처럼), 라벨에선 분리
                unit = ""
                mi = lbl.find("(단위")
                if mi != -1:
                    me = lbl.find(")", mi)
                    if me != -1:
                        unit = lbl[mi:me + 1]
                        lbl = (lbl[:mi] + lbl[me + 1:]).strip()
                    else:
                        unit = lbl[mi:]; lbl = lbl[:mi].strip()
                _add_textbox(slide, _TBL_L, t, tw - 2.2, 0.28, lbl,
                             size=12, bold=True, color=PALETTE["navy_dark"])
                if unit:
                    _add_textbox(slide, _TBL_L + tw - 2.4, t + 0.06, 2.4, 0.20, unit,
                                 size=9, bold=False, align=PP_ALIGN.RIGHT, color=PALETTE["gray_text"])
                t += _LABEL_H
            tbl_top = t
            _anchor_li = {li for (li, _a) in anchors} if anchors else None
            # ★분할표(k/m)에서 '맨 밑 합계=민트'는 마지막 청크에만 적용 → is_last_chunk 판단
            _is_last_chunk = True
            if lbl and lbl.rstrip().endswith(")") and "(" in lbl:
                _tail = lbl[lbl.rfind("(") + 1:lbl.rfind(")")]
                if "/" in _tail:
                    _a, _b = (_tail.split("/") + [""])[:2]
                    if _a.strip().isdigit() and _b.strip().isdigit():
                        _is_last_chunk = int(_a) >= int(_b)
            # ★Equity 투입현황표(위탁자/에쿼티 헤더)는 원본에 빨간 글씨 없음 →
            #   같은 슬라이드의 담보토지 빨강(루시드)이 번지지 않도록 red 억제
            _hdr_blob = " ".join(str(h) for h in (header or []))
            _rc = None if any(k in _hdr_blob for k in ("위탁자", "에쿼티")) else red_counts
            # ★표 안의 표(앵커 행): 원본은 '표 먼저 → 글' 순서 → 본 표 셀에선 값 텍스트를 비워
            #   위에 안 찍히게 하고, grid 아래에 따로 그린다(아래 overlay).
            _anchor_vals = {}
            if anchors and isinstance(rh, (list, tuple)):
                for (li, alabel) in anchors:
                    if _nested_for(alabel) and li < len(rows) and len(rows[li]) > 1:
                        _anchor_vals[li] = str(rows[li][1] or "")
                        rows[li][1] = ""
            used = _render_table_chunk(slide, kind, header, rows, ncol, _TBL_L, t, tw, fp, rh,
                                       red_counts=_rc, anchor_rows=_anchor_li,
                                       is_last_chunk=_is_last_chunk, fill_set=fill_set)
            # ★앵커 행 내용칸: (단위 캡션) → grid → 주N) 각주 → 요약 글씨 순으로 얹음(표 먼저, 글 나중)
            if anchors and isinstance(rh, (list, tuple)):
                lab_w = min(1.7, tw * 0.24)
                val_x = _TBL_L + lab_w
                val_w = tw - lab_w
                hh = _rowh(fp)
                for (li, alabel) in anchors:
                    ng = _nested_for(alabel)
                    if not ng:
                        continue
                    gk, gh, gb, gn = ng
                    grid_h = _grid_render_h(gh, gb, gn)
                    cur_y = tbl_top + hh + sum(rh[:li]) + 0.05
                    # (단위: …) — grid 위 우측 캡션(원본: Cash 표 위 '(단위: 백만원)')
                    _unit = _nested_unit_for(alabel)
                    if _unit:
                        _add_textbox(slide, val_x + 0.05, cur_y, val_w - 0.10, 0.18, _unit,
                                     size=9, align=PP_ALIGN.RIGHT, color=PALETTE["gray_text"])
                        cur_y += 0.22
                    try:
                        _render_table_chunk(slide, "grid", gh, gb, gn,
                                            val_x + 0.05, cur_y, val_w - 0.10,
                                            _grid_font(gn), _rowh(_grid_font(gn)), red_counts=red_counts,
                                            hdr_fill_hex="3E95BE", hdr_alpha=35)
                        cur_y += grid_h + 0.02
                        # grid 바로 밑 주N) 각주(원본처럼)
                        gnote = _nested_note_for(alabel)
                        if gnote.strip():
                            _, _gnh = _add_textbox(slide, val_x + 0.05, cur_y, val_w - 0.10, 0.2,
                                                   gnote, size=9, color=PALETTE["gray_text"],
                                                   red_set=red_set, ul_set=ul_set)
                            cur_y += (_gnh or 0.18) + 0.02
                        # ★표 다음에 요약 글씨(자금관리 ▶ 계좌 설명 등) — 표 먼저, 글 나중
                        _atext = _anchor_vals.get(li, "")
                        if _atext.strip():
                            _add_textbox(slide, val_x + 0.05, cur_y, val_w - 0.10, 0.2, _atext,
                                         size=fp, red_set=red_set, ul_set=ul_set, bullet=True)
                    except Exception as _e:
                        print(f"[nested grid] 실패: {_e}")
            t += used
            # ★이 표 '전용 각주'(주N)는 이 표 바로 밑에(다른 표/아래 페이지로 안 넘어감) ── 표별 각주 분리
            if tbl_notes:
                _tnt = "\n".join(tbl_notes)
                _tnh = _est_text_height(_tnt, _TBL_W, 9)
                _, _tnhh = _add_textbox(slide, _TBL_L, t + 0.04, _TBL_W, _tnh, _tnt,
                                        size=9, color=PALETTE["gray_text"], red_set=red_set, ul_set=ul_set)
                t += _tnhh + 0.06
            t += _GAP
        # ★본문 사진 배치(top_bare는 위에서 처리).
        #   side_box(조감도/광역입지)=첫 plan(건축개요 표) 오른쪽 / 부록=마지막 plan 표없이 크게
        if side_box and idx == img_plan_idx and big_imgs:
            labels = _img_labels_for(subtitle, min(2, len(big_imgs)))
            # 원본이 '표 안 사진'(조감도·광역위치도) → 라벨 헤더 표박스(각 1:1), 표 오른쪽
            _place_images_col(slide, big_imgs[:2], img_col_L, tbl_start,
                              _IMG_W, _BODY_BOTTOM, labels=labels)
        elif idx == n - 1 and big_imgs and not top_bare and not side_box and not has_tbl:
            # 표 없는 사진 페이지 → 4장 이상이면 그리드(첫장 크게+3열, 사업지 전경 6장 등), 아니면 한 줄
            if len(big_imgs) >= 4:
                _place_images_grid(slide, big_imgs, _TBL_L, tbl_start, _TBL_W, _BODY_BOTTOM)
            else:
                _place_images_row(slide, big_imgs[:2], _TBL_L, tbl_start,
                                  _TBL_W, _BODY_BOTTOM)
        # ★순서(원본): 표 → 주N) 각주(표 바로 밑, 9pt 회색) → 본문 산문(10pt 검정).
        #   주N)는 표에 딱 붙고, 분석 프로즈는 그 아래(원본: 인근 주요…는 주N) 다음).
        if idx == n - 1 and notes_text:
            nh = _est_text_height(notes_text, _TBL_W, 9)
            _, _nhh = _add_textbox(slide, _TBL_L, t + 0.04, _TBL_W, nh, notes_text,
                                   size=9, color=PALETTE["gray_text"], red_set=red_set, ul_set=ul_set)
            t += _nhh + 0.10
        if idx == n - 1 and btext:
            bh = _est_text_height(btext, tw, 10)
            _, _bhh = _add_textbox(slide, _TBL_L, t + 0.06, tw, bh, btext,
                                   size=10, bold=False, color=RGBColor(0x00, 0x00, 0x00),
                                   red_set=red_set, ul_set=ul_set, bullet=True)
            t += _bhh + 0.10
        _add_combined_footer(slide, business_name)
        # ★출처는 항상 좌측 하단 고정(별도 글상자) — 원본에 있을 때만
        if idx == n - 1 and source:
            _add_textbox(slide, _SRC_L, _SRC_T, _SRC_W, 0.22,
                         f"* 출처 : {source}", size=9, color=PALETTE["gray_text"])
    return None


# ──────────────────────────────────────────────────────
# 투자구조도 — 도형(상자/선/글상자)으로 직접 그림 (이미지 캡처 아님)
# ──────────────────────────────────────────────────────
def _set_line_arrow(conn, *, color="595959", width_pt=1.25, end=True, begin=False):
    """커넥터 선 색·두께·화살표 머리 설정(XML)."""
    ln = conn.line._get_or_add_ln()
    ln.set("w", str(int(width_pt * 12700)))
    # 색
    for tag in ("a:solidFill", "a:noFill"):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    sf = ln.makeelement(qn("a:solidFill"), {})
    sf.append(sf.makeelement(qn("a:srgbClr"), {"val": color}))
    ln.insert(0, sf)
    if begin:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if end:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def _dbox(slide, l, t, w, header, bodies, *, hdr_cm=0.9, body_cm=1.7, header_fill=None,
          body_fill=None, body_white=False):
    """다이어그램 상자 = 2~N행 1열 표(헤더 색 + 본문). graphicFrame 반환.
       ★고정 사양: 헤더행 0.9cm, 본문행 1.7cm(차주 기준), 글꼴 전부 Bold 10.5pt.
       body_fill 지정 시 본문 셀 색 변경(대주단 트랜치=진남색+흰글자 등)."""
    if isinstance(bodies, str):
        bodies = [bodies] if bodies else []
    bodies = [str(b) for b in bodies if str(b).strip()]
    nrow = 1 + len(bodies)
    total_h = Cm(hdr_cm + body_cm * len(bodies))
    gf = slide.shapes.add_table(nrow, 1, Inches(l), Inches(t), Inches(w), total_h)
    tb = gf.table
    # ★v22처럼 표 자체(tblPr)에 드롭섀도 추가(blurRad 50800/dist 38100/dir 2700000/검정 alpha 40%)
    _tblPr = tb._tbl.find(qn("a:tblPr"))
    if _tblPr is not None:
        for _e in _tblPr.findall(qn("a:effectLst")):
            _tblPr.remove(_e)
        _eff = _tblPr.makeelement(qn("a:effectLst"), {})
        _shd = _eff.makeelement(qn("a:outerShdw"),
                                {"blurRad": "50800", "dist": "38100", "dir": "2700000",
                                 "algn": "tl", "rotWithShape": "0"})
        _c = _shd.makeelement(qn("a:prstClr"), {"val": "black"})
        _c.append(_c.makeelement(qn("a:alpha"), {"val": "40000"}))
        _shd.append(_c)
        _eff.append(_shd)
        _tblPr.insert(0, _eff)
    tb.cell(0, 0).text = header
    for i, b in enumerate(bodies):
        tb.cell(1 + i, 0).text = b
    # ★헤더=지정색(기본 네이비), 본문=회색 D9D9D9(기본) 또는 body_fill
    style_table(gf, has_header=True, label_cols=(),
                header_fill=header_fill or PALETTE["navy_dark"],
                value_fill=body_fill or PALETTE["gray"])
    # ★행 높이는 style_table(헤더 0.6cm로 덮어씀) 뒤에 다시 강제 — 헤더 0.9cm / 본문 1.7cm
    tb.rows[0].height = Cm(hdr_cm)
    for i in range(1, nrow):
        tb.rows[i].height = Cm(body_cm)
    for ri in range(nrow):
        cell = tb.cell(ri, 0)
        cell.margin_top = cell.margin_bottom = Inches(0.01)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(13 if ri == 0 else 12)   # ★표 박스 글씨 헤더13/본문12pt(사용자 지시)
                r.font.name = _FONT_BOLD                    # ★전부 Bold 고정
                if ri >= 1 and body_white:
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # ★v22처럼 4변 흰색 테두리 2.25pt(28575) — 박스 사이 흰 간격
        tcPr = cell._tc.get_or_add_tcPr()
        for i, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
            for el in tcPr.findall(qn(tag)):
                tcPr.remove(el)
            ln = tcPr.makeelement(qn(tag), {"w": "28575", "cap": "flat", "cmpd": "sng", "algn": "ctr"})
            sf = ln.makeelement(qn("a:solidFill"), {})
            sf.append(sf.makeelement(qn("a:srgbClr"), {"val": "FFFFFF"}))
            ln.append(sf)
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
            tcPr.insert(i, ln)
    return gf


def _dlabel(slide, l, t, w, text):
    """선 위 설명 글상자(가운데, Bold 10.5pt). ★글상자 10.5pt + 전부 Bold + 검정(사용자 지시)."""
    _add_textbox(slide, l, t, w, 0.20, text, size=10.5, bold=True, align=PP_ALIGN.CENTER,
                 color=RGBColor(0x00, 0x00, 0x00))


def build_invest_diagram_slide(prs, data: dict, *, section_label, subtitle,
                               intro="", business_name=""):
    """투자구조도(2.1) — 참여기관을 '표 박스(네이비 헤더+본문)'와 연결선으로 그림(v22 방식)."""
    ent = (data or {}).get("entities", {}) or {}
    tranches = (data or {}).get("tranches", []) or []

    def _name(key, default):
        v = str(ent.get(key, "") or "").strip() or default
        for sep in ("—", "–", " - ", ":"):     # '신탁사 — 한국투자…' → '한국투자…'
            if sep in v:
                v = v.split(sep)[-1].strip()
        return v

    slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)
    _fill_header(slide, section_label, subtitle, "")
    if intro and intro.strip():
        _add_textbox(slide, _INTRO_L, _INTRO_T, _INTRO_W, 0.40, intro.strip(), size=9, bold=True)

    # ★페이지 상·하단 가로 구분선(예시처럼)
    for ly in (1.50, 7.05):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          Inches(0.43), Inches(ly), Inches(10.40), Inches(ly))
        _set_line_arrow(line, color="BFBFBF", width_pt=1.0, end=False)

    # ★v22 3색 헤더: 차주=진남색, 신탁사·시공사·SPC·사채권자=중간파랑, 대주단=적갈색
    navy = PALETTE["navy_dark"]            # 08377C
    blue = RGBColor(0x00, 0x63, 0xA1)      # 0063A1
    maroon = RGBColor(0x8C, 0x4A, 0x59)    # 8C4A59

    tr_lines = []
    for tr in tranches:
        nm = tr.get("name", "") if isinstance(tr, dict) else str(tr)
        amt = tr.get("amount", "") if isinstance(tr, dict) else ""
        tr_lines.append((nm + (f" {amt}" if amt else "")).strip())
    if not tr_lines:
        tr_lines = ["Tr.A", "Tr.B", "Tr.C"]

    # ── 박스 배치는 딜별 원본 레이아웃대로(in). 기본값 'T.B.D.'(대전 회사명 하드코딩 금지) ──
    is_dj = "터미널" in str(business_name)
    if is_dj:
        # 대전: 신탁사·차주 같은 행, 시공사 차주 아래, 대주 우측(승인 레이아웃)
        trustee = _dbox(slide, 0.70, 2.75, 2.05, "신탁사", _name("trustee", "T.B.D."), header_fill=blue)
        borrower = _dbox(slide, 4.20, 2.75, 2.05, "차주", _name("borrower", "T.B.D."), header_fill=navy)
        constructor = _dbox(slide, 4.20, 4.70, 2.05, "시공사", _name("constructor", "T.B.D."), header_fill=blue)
        lenders = _dbox(slide, 7.90, 2.45, 2.20, "본건 PF 대주단", tr_lines,
                        hdr_cm=0.9, body_cm=0.95, header_fill=maroon)
    else:
        # 천안형(일반): 신탁사 상단중앙 / 시공사 좌 · 차주 중앙 · 대주 우 (한 행)
        trustee = _dbox(slide, 4.20, 1.85, 2.05, "신탁사", _name("trustee", "T.B.D."), header_fill=blue)
        constructor = _dbox(slide, 0.70, 3.95, 2.05, "시공사", _name("constructor", "T.B.D."), header_fill=blue)
        borrower = _dbox(slide, 4.20, 3.95, 2.05, "차주/시행사", _name("borrower", "T.B.D."), header_fill=navy)
        lenders = _dbox(slide, 7.90, 3.60, 2.20, "본건 PF 대주단", tr_lines,
                        hdr_cm=0.9, body_cm=0.95, header_fill=maroon)
    # ★원본 구조도는 박스 4개(신탁사·차주·시공사·대주단)만 — SPC/사채권자 없음

    def cx(sh, fx=0.5):
        return (sh.left + int(sh.width * fx)) / 914400

    def cy(sh, fy=0.5):
        return (sh.top + int(sh.height * fy)) / 914400

    def conn(x1, y1, x2, y2, **kw):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        _set_line_arrow(c, **kw)
        return c

    # ── 연결선 라벨 = 딜별(원본대로). 대전(사업명에 '터미널')은 승인된 라벨 고정, 그 외 딜은
    #   LLM relationships에서 추출(예: 천안=공사도급계약/대출약정/담보신탁계약). 대전 라벨이 타 딜에
    #   복붙되던 문제 해결. 시공사→대주 '물상보증'은 대전 전용(천안 등 원본엔 없음 → 생략).
    if is_dj:
        lab_tb, lab_bc, lab_bl = "사업부지 담보신탁 계약 체결", "공동사업약정 체결", "대출실행 / 원리금상환"
        lab_tl, lab_cl = "순위별 우선수익권 제공", "터미널 부지 물상보증"
    else:
        _rels = data.get("relationships") or []

        def _flbl(*kws, default=""):
            for r in _rels:
                s = str(r)
                lab = (s.split(":", 1)[1].strip().split("(")[0].strip()) if ":" in s else ""
                if lab and any(k in lab for k in kws):
                    return lab
            return default
        lab_tb = _flbl("담보신탁", "신탁계약", default="담보신탁계약")
        lab_bc = _flbl("공사도급", "도급", "공동사업", default="공사도급계약")
        lab_bl = _flbl("대출약정", "대출", default="대출약정")
        lab_tl = _flbl("우선수익", default="담보신탁 우선수익권")
        lab_cl = None      # 시공사→대주(물상보증)는 대전 외 딜엔 없음

    if is_dj:
        # 신탁사 → 차주(수평)
        conn(cx(trustee, 1.0), cy(trustee), cx(borrower, 0.0), cy(borrower))
        _dlabel(slide, 2.80, 3.12, 1.55, lab_tb)
        # 차주 → 시공사(수직)
        conn(cx(borrower), cy(borrower, 1.0), cx(constructor), cy(constructor, 0.0), begin=True)
        _dlabel(slide, 6.30, 4.10, 1.5, lab_bc)
        # 차주 ↔ 대주단(수평)
        conn(cx(borrower, 1.0), cy(borrower), cx(lenders, 0.0), cy(lenders, 0.55), begin=True)
        _dlabel(slide, 6.35, 3.02, 1.55, lab_bl)
        # 시공사 → 대주단(물상보증, 대각)
        if lab_cl:
            conn(cx(constructor, 1.0), cy(constructor), cx(lenders, 0.0), cy(lenders, 0.9))
            _dlabel(slide, 6.35, 4.95, 1.55, lab_cl)
        # 신탁사 → 대주단(우선수익권, 상단 아치)
        ax = 2.30
        conn(cx(trustee), cy(trustee, 0.0), cx(trustee), ax)
        conn(cx(trustee), ax, cx(lenders), ax)
        conn(cx(lenders), ax, cx(lenders), cy(lenders, 0.0), end=True)
        _dlabel(slide, (cx(trustee) + cx(lenders)) / 2 - 1.0, ax - 0.22, 2.0, lab_tl)
    else:
        # 천안형: 신탁사↓차주(수직) · 시공사↔차주(수평) · 차주↔대주(수평) · 신탁사→대주(상단 아치)
        conn(cx(trustee), cy(trustee, 1.0), cx(borrower), cy(borrower, 0.0), begin=True, end=True)
        _dlabel(slide, cx(trustee) + 0.12, (cy(trustee, 1.0) + cy(borrower, 0.0)) / 2 - 0.16, 1.5, lab_tb)
        conn(cx(constructor, 1.0), cy(constructor), cx(borrower, 0.0), cy(borrower), begin=True, end=True)
        _dlabel(slide, (cx(constructor, 1.0) + cx(borrower, 0.0)) / 2 - 0.75, cy(borrower) - 0.52, 1.5, lab_bc)
        conn(cx(borrower, 1.0), cy(borrower), cx(lenders, 0.0), cy(lenders, 0.5), begin=True, end=True)
        _dlabel(slide, (cx(borrower, 1.0) + cx(lenders, 0.0)) / 2 - 0.75, cy(borrower) - 0.52, 1.5, lab_bl)
        ax = cy(trustee, 0.5)
        conn(cx(trustee, 1.0), ax, cx(lenders, 0.5), ax)
        conn(cx(lenders, 0.5), ax, cx(lenders, 0.5), cy(lenders, 0.0), end=True)
        _dlabel(slide, (cx(trustee, 1.0) + cx(lenders, 0.5)) / 2 - 1.0, ax - 0.30, 2.2, lab_tl)

    _add_combined_footer(slide, business_name)
    return None


# ──────────────────────────────────────────────────────
# 투자구조도 — 원본 PDF의 구조도 영역을 이미지로 캡처해 삽입(딜마다 달라 자동그리기 대신 캡처)
# ──────────────────────────────────────────────────────
def render_structure_image(pdf_path, dpi_scale=3):
    """원본 PDF에서 '구조도'가 있는 페이지의 다이어그램 영역을 PNG로 렌더 → {data,w,h} 또는 None."""
    try:
        import fitz
    except Exception:
        return None
    doc = fitz.open(pdf_path)
    try:
        target = None
        for i in range(len(doc)):
            if "구조도" in (doc[i].get_text("text") or ""):
                target = i
                break
        if target is None:
            return None
        page = doc[target]
        # '구조도' 라벨 y 위치
        ky = None
        for b in page.get_text("dict").get("blocks", []):
            if b.get("type") != 0:
                continue
            t = "".join(s["text"] for ln in b.get("lines", []) for s in ln.get("spans", []))
            if "구조도" in t:
                ky = b["bbox"][1]
                break
        # 라벨 아래 벡터 드로잉들의 합집합 bbox
        xs, ys = [], []
        for d in page.get_drawings():
            r = d["rect"]
            if r.y0 >= (ky - 4 if ky else 0) and r.width < page.rect.width and r.height < page.rect.height:
                xs += [r.x0, r.x1]
                ys += [r.y0, r.y1]
        if not xs:
            return None
        pad = 6
        clip = fitz.Rect(max(0, min(xs) - pad), max(0, min(ys) - pad),
                         min(page.rect.width, max(xs) + pad),
                         min(page.rect.height, max(ys) + pad))
        if clip.width < 40 or clip.height < 40:
            return None
        pix = page.get_pixmap(clip=clip, matrix=fitz.Matrix(dpi_scale, dpi_scale))
        return {"data": pix.tobytes("png"), "w": pix.width, "h": pix.height}
    finally:
        doc.close()


def build_invest_image_slide(prs, *, section_label, subtitle, intro, img,
                             business_name=""):
    """투자구조도 슬라이드 — 섹션라벨/소제목/인트로 + 캡처한 구조도 이미지(본문 중앙 배치)."""
    slide = clone_slide_layout(prs, "content", skip_graphic_frames=True)
    _fill_header(slide, section_label, subtitle, "")
    top = _BODY_TOP_START
    if intro and intro.strip():
        _, h = _add_textbox(slide, _INTRO_L, _INTRO_T, _INTRO_W, 0.40, intro.strip(), size=9)
        top = max(top, _INTRO_T + h + 0.12)
    if img and img.get("data"):
        avail_h = _BODY_BOTTOM - top - 0.1
        avail_w = _TBL_W
        scale = min(avail_w / img["w"], avail_h / img["h"])
        iw, ih = img["w"] * scale, img["h"] * scale
        px = _TBL_L + (avail_w - iw) / 2
        slide.shapes.add_picture(io.BytesIO(img["data"]), Inches(px), Inches(top),
                                 Inches(iw), Inches(ih))
    _add_combined_footer(slide, business_name)
    return slide


# ──────────────────────────────────────────────────────
# 푸터 — 우측에 "사업명ㅣ쪽번호" 합쳐서 (대전 방식). 쪽번호는 슬라이드번호 필드.
# ──────────────────────────────────────────────────────
def _add_combined_footer(slide, business_name):
    # 기존 템플릿 푸터(사업명 자리표시자 + '슬라이드 번호' 도형) 제거 → 중복 방지
    for sh in list(slide.shapes):
        nm = sh.name or ""
        is_pagenum = "슬라이드 번호" in nm
        is_biz_footer = (sh.has_text_frame and (sh.left or 0) / 360000 < 1.0
                         and (sh.top or 0) / 360000 > 17.5)
        if is_pagenum or is_biz_footer:
            try:
                sh._element.getparent().remove(sh._element)
            except Exception:
                pass

    tb = slide.shapes.add_textbox(Inches(7.18), Inches(7.19), Inches(3.32), Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    run = p.add_run()
    run.text = f"{business_name or ''}｜"   # 전각 세로막대 ｜ (대전 원본 방식)
    run.font.name = _FONT_LIGHT
    run.font.size = Pt(8)
    run.font.color.rgb = PALETTE["gray_text"]

    # 슬라이드 번호 필드 (자동 쪽번호)
    fld = p._p.makeelement(qn('a:fld'), {
        'id': '{B7E3A1C2-0000-4000-9000-000000000001}', 'type': 'slidenum'})
    rPr = fld.makeelement(qn('a:rPr'), {'lang': 'ko-KR', 'sz': '800'})
    sf = rPr.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:srgbClr'), {'val': '808080'})
    sf.append(clr)
    rPr.append(sf)
    latin = rPr.makeelement(qn('a:latin'), {'typeface': _FONT_LIGHT})
    ea = rPr.makeelement(qn('a:ea'), {'typeface': _FONT_LIGHT})
    rPr.append(latin)
    rPr.append(ea)
    t_el = fld.makeelement(qn('a:t'), {})
    t_el.text = "1"
    fld.append(rPr)
    fld.append(t_el)
    p._p.append(fld)
    return tb
