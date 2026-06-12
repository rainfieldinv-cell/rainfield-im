"""5단계 '내용 검수' — 생성 PPT가 원본(PDF/워드)과 글자·숫자 단위로 1:1 일치하는지 정밀 점검.

★읽기 전용: 본문 생성 로직을 건드리지 않고, 문제를 '찾아서 보여주기만' 한다(자동 수정 X).
점검 항목
  1) 정밀 대조 — 원본의 '값'(숫자 전체형·데이터성 이름)이 PPT에 정확히 옮겨졌는지
       · 숫자 누락 / 숫자 오류 의심(1,640→1,460) / 형식·단위·쉼표 차이 / 글자 잘림(더함도시개발→더함도시개) / 내용 누락
       · 페이지별 '원본 대비 일치율(%)'을 계산(낮은 페이지 우선 정렬)
  2) 빈 표 셀 — 데이터 누락 의심 칸
노이즈 제거(의도적 변환은 오류로 안 잡음): 페이지번호(2/26)·줄바꿈·공백·괄호 템플릿([1,640])·표 배치는 무시.
★맞춤법(hanspell)은 비활성 — 목적은 '원본 일치'이지 문법 검사가 아님.
반환 item 스키마: {"page","type","original","ppt","rate"}.
"""
import io
import re
from collections import Counter, defaultdict

from pptx import Presentation

# 값에 붙는 단위(이게 붙은 숫자만 '중요 값'으로 본다)
_UNIT = r"(억원|억|만원|만|천원|원|％|%|평|㎡|세대|개월|개동|동|호|명|위|배|건|개|차|호선|년|개소)"
# 숫자 전체형: 선행 @ + 천단위쉼표/소수 + 후행 단위
_FULLNUM_RE = re.compile(r"@?\s?\d[\d,]*(?:\.\d+)?\s*" + _UNIT + r"?")
# 데이터성 한글 이름(회사·기관·단지 등) — 특정 접미사로 끝나는 토큰
_NAME_RE = re.compile(
    r"[가-힣A-Za-z0-9·]{2,}(?:신탁|건설|증권|개발|산업|자산운용|캐피탈|은행|토건|중공업|"
    r"에쿼티|파트너스|투자|금융|보험|공사|법인|주식회사|아파트|디벨로퍼|엔지니어링)")
# 페이지번호 류(비교 제외)
_PAGEFOOT_RE = re.compile(r"^\s*(\d+\s*/\s*\d+|페이지\s*\d+|\d+\s*페이지|- ?\d+ ?-)\s*$")
# '값'으로 의미 있는 단위(금액·비율·면적·세대/명). 차/호/년/월 등 날짜·서수는 1:1 값이 아니라 제외.
_VALUE_UNITS = {"억원", "억", "만원", "만", "천원", "원", "%", "％", "평", "㎡", "세대", "명"}
# 이름 접미사로 끝나지만 고유명사가 아닌 '일반 분류어'(누락 오탐 방지)
_STOP_NAMES = {"부동산금융", "프로젝트금융", "구조화금융", "도시개발", "부동산개발", "주택개발",
               "자산운용", "종합건설", "일반건설", "전문건설", "해외건설", "주택건설",
               "토목건설", "부동산투자", "간접투자", "직접투자", "공동주택", "민간투자"}


def _digits(s: str) -> str:
    return re.sub(r"\D", "", s)


def _form(s: str) -> str:
    """비교용 정규형 — 공백·대괄호·소괄호 제거(쉼표·단위·@·% 는 유지해서 표기차이를 잡는다)."""
    return re.sub(r"[\s\[\]()]", "", s)


def _ctx(text: str, start: int, end: int, pad: int = 16) -> str:
    """원본에서 토큰 주변(앞뒤 pad자)을 잘라 '어디에 있던 값인지' 맥락 제공."""
    a, b = max(0, start - pad), min(len(text), end + pad)
    s = re.sub(r"\s{2,}", " ", text[a:b].replace("\n", " ")).strip()
    return ("…" if a > 0 else "") + s + ("…" if b < len(text) else "")


def _is_summary_page(text: str) -> bool:
    """Executive Summary / 요약 페이지인지 — 이런 페이지는 원본을 '의도적으로 요약·생략'하므로
       1:1 대조(누락 판정)에서 제외한다(단, 잘못 들어간 값=잘림/숫자오류는 계속 검출)."""
    if re.search(r"executive\s*summary", text, re.I):
        return True
    for ln in text.splitlines()[:12]:                     # 상단 제목/머리말 영역만
        s = ln.strip()
        if 0 < len(s) <= 16 and "요약" in s:
            return True
    return False


# ──────────────────────────────────────────────────────────
# PPT 텍스트 추출
# ──────────────────────────────────────────────────────────
def _ppt_full_text(prs) -> str:
    parts = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text:
                parts.append(sh.text_frame.text)
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        if c.text:
                            parts.append(c.text)
    return "\n".join(parts)


def _ppt_korean_lines(prs):
    """(슬라이드번호, 텍스트) — 맞춤법 검사용 한글 포함 줄."""
    out = []
    for i, s in enumerate(prs.slides, start=1):
        for sh in s.shapes:
            if sh.has_text_frame:
                for ln in sh.text_frame.text.split("\n"):
                    t = ln.strip()
                    if len(t) >= 6 and re.search(r"[가-힣]", t):
                        out.append((i, t))
    return out


# ──────────────────────────────────────────────────────────
# 1) 정밀 대조 — 원본의 '값'(숫자·이름)이 PPT에 정확히 옮겨졌는지 + 페이지별 일치율
# ──────────────────────────────────────────────────────────
def _checkable_tokens(text: str):
    """원본 텍스트에서 1:1로 옮겨져야 하는 '값' 토큰: (종류, 토큰, 주변맥락)."""
    out = []
    for m in _FULLNUM_RE.finditer(text):
        tok = re.sub(r"\s+", "", m.group(0))
        d = _digits(tok)
        unit = m.group(1)
        # 값으로 의미있는 것만: 천단위쉼표 큰수 / @단가 / 금액·비율·면적·세대 단위.
        # 날짜(22.07)·서수(26차)·기간(24년)은 단위가 빠지거나 값단위가 아니라 자동 제외.
        if len(d) >= 2 and (("," in tok) or "@" in tok or (unit in _VALUE_UNITS)):
            out.append(("num", tok, _ctx(text, m.start(), m.end())))
    for m in _NAME_RE.finditer(text):
        nm = m.group(0).strip()
        if len(nm) >= 4 and nm not in _STOP_NAMES:
            out.append(("name", nm, _ctx(text, m.start(), m.end())))
    return out


def _near_number(d, by_len, maxdiff=2):
    """같은 자릿수인데 1~2자리만 다른 PPT 숫자(오타/전치 의심) 반환 — 없으면 None.
       후보가 둘 이상이면 애매하므로 None(오탐 방지)."""
    cands = [c for c in by_len.get(len(d), ()) if c != d
             and 1 <= sum(1 for a, b in zip(d, c) if a != b) <= maxdiff]
    return cands[0] if len(cands) == 1 else None


def _fmt_commas(d):
    """숫자 문자열에 천단위 쉼표(소수 없는 정수만)."""
    return f"{int(d):,}" if d.isdigit() else d


def compare_content(pages_text, prs):
    """원본 페이지별로 값 토큰을 PPT와 정밀 대조 → (items, page_rate).
       page_rate: {원본페이지: 일치율(%)}. items 각 행에 해당 페이지 일치율 포함."""
    ptext = _ppt_full_text(prs)
    pform = _form(ptext)
    numset, by_len = set(), defaultdict(set)
    for m in _FULLNUM_RE.finditer(ptext):
        d = _digits(m.group(0))
        if len(d) >= 2:
            numset.add(d)
            by_len[len(d)].add(d)

    _ABSENT = ("숫자 누락", "내용 누락")
    items, page_rate, global_seen = [], {}, set()
    for pi, raw in enumerate(pages_text or [], start=1):
        # 페이지번호 줄 제거(노이즈)
        text = "\n".join(ln for ln in (raw or "").splitlines()
                         if not _PAGEFOOT_RE.match(ln.strip()))
        is_sum = _is_summary_page(text)                   # 요약/ES 페이지는 누락 판정 제외
        page_rows, seen = [], set()
        total = matched = 0
        for kind, tok, ctx in _checkable_tokens(text):
            f = _form(tok)
            key = (kind, f)
            if not f or key in seen:
                continue
            seen.add(key)
            total += 1
            if f in pform:                                # 표기까지 정확히 일치
                matched += 1
                continue
            if kind == "num":
                d = _digits(tok)
                if d in numset:                           # 수치는 맞고 표기만 다름
                    ty, ppt = "형식/단위/쉼표 차이", "수치는 있으나 표기 다름"
                else:
                    near = _near_number(d, by_len)
                    ty, ppt = (("숫자 오류 의심", f"PPT '{_fmt_commas(near)}'") if near
                               else ("숫자 누락", "없음"))
            else:                                         # name
                pref = next((f[:L] for L in (len(f) - 1, len(f) - 2)
                             if L >= 4 and f[:L] in pform), None)
                ty, ppt = (("글자 잘림 의심", f"PPT '{pref}…'") if pref
                           else ("내용 누락", "없음"))
            # 요약/ES 페이지의 '누락'은 의도적 생략이므로 대조 대상에서 제외(잘림/오류는 유지)
            if is_sum and ty in _ABSENT:
                total -= 1
                continue
            page_rows.append((ty, tok, ppt, ctx))
        rate = round(matched / total * 100) if total else 100
        if total and not is_sum:                          # 요약 페이지는 일치율 산정 제외
            page_rate[pi] = rate
        for ty, orig, ppt, ctx in page_rows:
            gkey = (pi, ty, orig)
            if gkey in global_seen:
                continue
            global_seen.add(gkey)
            items.append({"page": f"원본 {pi}p", "type": ty, "original": orig,
                          "context": ctx, "ppt": ppt,
                          "rate": ("요약" if is_sum else rate)})
    return items, page_rate


# ──────────────────────────────────────────────────────────
# 2) 빈 표 셀 — 병합(연속)칸 제외, 헤더행 제외한 진짜 빈 칸
# ──────────────────────────────────────────────────────────
def _is_merge_continuation(cell):
    tc = cell._tc
    return (tc.get("hMerge") in ("1", "true") or tc.get("vMerge") in ("1", "true"))


def find_empty_cells(prs):
    """병합 연속칸·헤더·설계상 빈 열은 제외하고, '대부분 채워진 열에 뚫린 빈칸'(데이터 누락 의심)만 표시.
       — 발코니확장/근린생활시설처럼 원래 비는 칸(열 자체가 듬성)은 노이즈라 걸러냄."""
    rows = []
    for i, s in enumerate(prs.slides, start=1):
        for sh in s.shapes:
            if not sh.has_table:
                continue
            t = sh.table
            ncol, nrow = len(t.columns), len(t.rows)
            ndata = max(1, nrow - 1)
            # 열별 채움률(데이터행 기준)
            filled = [0] * ncol
            for ci in range(ncol):
                for ri in range(1, nrow):
                    if t.rows[ri].cells[ci].text.strip():
                        filled[ci] += 1
            for ri in range(1, nrow):
                _c0 = t.rows[ri].cells[0].text.strip()
                if any(k in _c0 for k in ("합계", "소계", "합 계", "소 계", "총계", "총합계", "총 계")):
                    continue                              # 합계/소계행은 빈칸이 정상 → 제외
                # 행 자체가 듬성(데이터칸 절반 이상 빔)하면 비용항목 등 설계상 빈행 → 제외
                _rowfill = sum(1 for cj in range(1, ncol) if t.rows[ri].cells[cj].text.strip())
                if _rowfill < (ncol - 1) * 0.5:
                    continue
                for ci, c in enumerate(t.rows[ri].cells):
                    if c.text.strip() or _is_merge_continuation(c):
                        continue
                    if filled[ci] < ndata * 0.6:          # 열 자체가 듬성하면(설계상) 제외
                        continue
                    lbl = t.rows[ri].cells[0].text.strip() or (
                        t.rows[0].cells[ci].text.strip() if ncol > ci else "") or "?"
                    rows.append({"page": f"슬라이드 {i}", "type": "빈 표 셀",
                                 "original": f"'{lbl[:18]}' 행({ri+1}행 {ci+1}열)",
                                 "context": "표 내 빈칸", "ppt": "빈 셀", "rate": "-"})
    return rows


# ──────────────────────────────────────────────────────────
# 3) 오타 / 맞춤법
# ──────────────────────────────────────────────────────────
def check_spelling(prs):
    """맞춤법 검사 — hanspell(네이버) 설치/동작 시에만 수행. 없으면 '건너뛰기'.
       ★py-hanspell은 클라우드 설치가 자주 실패하므로 requirements에 넣지 않는다.
         라이브러리 미설치/호출 실패 시 빈 결과 + '비활성' 표시 → 앱은 절대 죽지 않음.
         (내용 1:1 대조·빈 표 셀 검출은 이 함수와 무관하게 정상 동작)"""
    try:
        from hanspell import spell_checker            # 설치돼 있을 때만(기본 미설치)
    except Exception:
        return [], "맞춤법 검사 비활성(라이브러리 미설치)"

    rows, ok = [], False
    for i, t in _ppt_korean_lines(prs):
        try:
            res = spell_checker.check(t[:480])         # hanspell 길이 제한 대비
        except Exception:
            continue
        ok = True
        errs = getattr(res, "errors", 0)
        checked = getattr(res, "checked", "")
        if errs and checked and checked.strip() != t.strip():
            rows.append({"page": f"슬라이드 {i}", "type": "맞춤법",
                         "content": t[:40], "suggestion": checked[:60]})
    if not ok:
        return [], "맞춤법 검사 비활성(호출 실패 — 네트워크)"
    return rows, "hanspell(네이버 맞춤법)"


# ──────────────────────────────────────────────────────────
# 통합 검수
# ──────────────────────────────────────────────────────────
def review_presentation(ppt_bytes: bytes, pages_text):
    """PPT bytes + 원본 페이지텍스트 → 정밀 검수 결과 dict.
       items: [{page,type,original,ppt,rate}] — 일치율 낮은 페이지(=문제 많은 곳) 우선 정렬.
       page_rate: {원본페이지: 일치율} (낮은 순)."""
    prs = Presentation(io.BytesIO(ppt_bytes))
    compared, page_rate = compare_content(pages_text, prs)
    empty = find_empty_cells(prs)
    spelling, spell_engine = check_spelling(prs)        # 비활성(목적이 문법검사가 아님)
    items = compared + empty + spelling

    def _sort_key(it):
        r = it.get("rate")
        return (0, r) if isinstance(r, (int, float)) else (1, 999)  # 일치율 낮은 행 먼저, 빈셀 뒤로

    items.sort(key=_sort_key)
    counts = dict(Counter(it["type"] for it in items))
    return {
        "ok": len(items) == 0,
        "items": items,
        "counts": counts,
        "page_rate": dict(sorted(page_rate.items(), key=lambda kv: kv[1])),
        "spell_engine": spell_engine,
    }
