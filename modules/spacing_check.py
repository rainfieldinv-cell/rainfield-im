"""6단계 '간격 점검' — PPT 텍스트의 띄어쓰기·공백·오버플로우·빈줄 문제를 점검(읽기전용, 자동수정 X).

★맞춤법은 5단계(content_review)에서 처리하므로 여기선 간격/띄어쓰기/넘침만 본다.
점검 항목
  1) 연속 공백(스페이스 2개 이상)·셀/문단 끝 불필요 공백
  2) 흔한 띄어쓰기 오류(의존명사 '수' 등 — 가능한 범위의 규칙 기반)
  3) 오버플로우 — 표 셀/글상자 폭 대비 텍스트가 길어 넘칠 가능성(폰트·폭 추정)
  4) 과도한 빈 문단/빈 줄(연속 2개 이상)
반환 형식은 5단계와 동일: {"ok", "items":[{page,type,content,suggestion}], "counts"}.
"""
import io
import re

from pptx import Presentation

# 흔한 '붙여쓴 띄어쓰기 오류'(있으면 거의 확실히 틀림) — 의존명사 '수' 중심으로 보수적으로.
_SPACING_PATTERNS = [
    (re.compile(r"(할|볼|갈|올|될|들|쓸|줄|살|일|날)수\s*(있|없|밖)"), "‘~ㄹ 수 있/없’ 띄어쓰기"),
    (re.compile(r"것입니다같"), "‘것입니다’ 뒤 띄어쓰기"),
    (re.compile(r"등을통해|등의경우|의경우엔?[가-힣]"), "조사 뒤 띄어쓰기 확인"),
]


def _font_size(tf, default=10.5):
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size:
                return r.font.size.pt
    return default


def _text_width_in(text, size):
    """텍스트 한 줄의 대략적 표시 폭(in) — 한글/전각 ~0.82×, 그 외 ~0.45×(피플폰트 기준 보수적)."""
    w = 0.0
    for ch in text:
        w += size * (0.82 if ord(ch) > 0x2000 else 0.45) / 72.0
    return w


def _strip_section_label(t):
    """'01  사모사채 개요'처럼 섹션번호 뒤 2칸은 디자인(노이즈) → 점검에서 앞부분 제거."""
    return re.sub(r"^\s*\d{1,2}\s{1,3}", "", t)


# ──────────────────────────────────────────────────────────
# 1) 연속/불필요 공백
# ──────────────────────────────────────────────────────────
def _check_spaces(prs):
    rows = []
    for i, s in enumerate(prs.slides, start=1):
        # 글상자
        for sh in s.shapes:
            if sh.has_text_frame:
                for ln in sh.text_frame.text.split("\n"):
                    t = ln.rstrip("\n")
                    body = _strip_section_label(t)
                    if re.search(r"\S {2,}\S", body):
                        rows.append({"page": f"슬라이드 {i}", "type": "연속 공백",
                                     "content": t.strip()[:40], "suggestion": "공백 1칸으로 정리"})
                    elif t != t.rstrip() and t.strip():
                        rows.append({"page": f"슬라이드 {i}", "type": "끝 공백",
                                     "content": t.strip()[:40], "suggestion": "줄 끝 공백 제거"})
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        ct = c.text
                        if not ct.strip():
                            continue
                        if re.search(r"\S {2,}\S", ct):
                            rows.append({"page": f"슬라이드 {i}", "type": "연속 공백",
                                         "content": f"표 '{ct.strip()[:24]}'", "suggestion": "공백 1칸으로 정리"})
                        elif ct != ct.strip():
                            rows.append({"page": f"슬라이드 {i}", "type": "끝 공백",
                                         "content": f"표 '{ct.strip()[:24]}'", "suggestion": "셀 앞뒤 공백 제거"})
    return rows


# ──────────────────────────────────────────────────────────
# 2) 띄어쓰기(보수적 규칙)
# ──────────────────────────────────────────────────────────
def _check_spacing_rules(prs):
    rows = []
    for i, s in enumerate(prs.slides, start=1):
        for sh in s.shapes:
            texts = []
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
            if sh.has_table:
                texts += [c.text for r in sh.table.rows for c in r.cells]
            for t in texts:
                for pat, msg in _SPACING_PATTERNS:
                    m = pat.search(t)
                    if m:
                        rows.append({"page": f"슬라이드 {i}", "type": "띄어쓰기",
                                     "content": m.group(0)[:24], "suggestion": msg})
    return rows


# ──────────────────────────────────────────────────────────
# 3) 오버플로우 — 폭 대비 텍스트 길이로 추정
# ──────────────────────────────────────────────────────────
def _check_overflow(prs):
    rows = []
    for i, s in enumerate(prs.slides, start=1):
        for sh in s.shapes:
            if sh.has_table:
                t = sh.table
                ncol = len(t.columns)
                colw = [t.columns[k].width / 914400 for k in range(ncol)]
                for r in t.rows:
                    for ci, c in enumerate(r.cells):
                        txt = c.text.strip()
                        tc = c._tc
                        if not txt or tc.get("hMerge") or tc.get("vMerge"):
                            continue
                        if c.text_frame.word_wrap:        # 줄바꿈 허용 셀은 가로 넘침 아님
                            continue
                        gs = int(tc.get("gridSpan") or 1)
                        cw = sum(colw[ci:ci + gs])
                        size = _font_size(c.text_frame)
                        longest = max((_text_width_in(ln, size) for ln in txt.split("\n")), default=0)
                        if longest > cw * 1.2:            # 명백히 넘칠 때만(추정 오차 감안 보수적)
                            rows.append({"page": f"슬라이드 {i}", "type": "오버플로우",
                                         "content": f"표 '{txt[:18]}' (폭 {cw:.1f}\")",
                                         "suggestion": "폰트 축소/열 확장 검토"})
            # 글상자는 앱이 텍스트 높이에 맞춰 크기를 잡으므로 가로 넘침 오탐이 많아 제외.
    return rows


# ──────────────────────────────────────────────────────────
# 4) 과도한 빈 줄/빈 문단
# ──────────────────────────────────────────────────────────
def _check_blank_lines(prs):
    rows = []
    for i, s in enumerate(prs.slides, start=1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            paras = sh.text_frame.paragraphs
            consec = 0
            for p in paras:
                if p.text.strip() == "":
                    consec += 1
                    if consec == 2:                       # 연속 빈 문단 2개 이상
                        rows.append({"page": f"슬라이드 {i}", "type": "빈 줄 과다",
                                     "content": "연속 빈 문단", "suggestion": "불필요한 빈 줄 제거"})
                else:
                    consec = 0
    return rows


# ──────────────────────────────────────────────────────────
# 통합
# ──────────────────────────────────────────────────────────
def check_spacing(ppt_bytes: bytes):
    prs = Presentation(io.BytesIO(ppt_bytes))
    sp = _check_spaces(prs)
    rule = _check_spacing_rules(prs)
    over = _check_overflow(prs)
    blank = _check_blank_lines(prs)
    items = sp + rule + over + blank
    return {
        "ok": len(items) == 0,
        "items": items,
        "counts": {
            "연속/끝 공백": len(sp),
            "띄어쓰기": len(rule),
            "오버플로우": len(over),
            "빈 줄 과다": len(blank),
        },
    }
